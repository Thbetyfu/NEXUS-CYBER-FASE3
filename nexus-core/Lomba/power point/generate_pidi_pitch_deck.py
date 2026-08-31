#!/usr/bin/env python3
"""
Generate Nexus Cyber PIDI Capstone pitch deck.

Structure:
  PART A — 14 slide inti (wajib PIDI Digdaya x Hackathon 2026)
  PART B — Additional slides untuk Q&A (data, metrik, appendix)

PS resmi: PS1 Penguatan Ketahanan dan Inovasi Keuangan · Sub: Manajemen Risiko
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Nexus-Cyber-PIDI-Capstone.pptx"

# Brand
BG = RGBColor(0x07, 0x0B, 0x10)
CARD = RGBColor(0x0E, 0x16, 0x1C)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
AMBER = RGBColor(0xF5, 0xB3, 0x01)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
ROSE = RGBColor(0xFB, 0x71, 0x85)
LINE = RGBColor(0x1E, 0x29, 0x3B)
CORE_BADGE = RGBColor(0x0D, 0x3D, 0x38)
ADD_BADGE = RGBColor(0x3D, 0x2A, 0x0D)

W, H = Inches(13.333), Inches(7.5)

PS_MAIN = "PS1 · Penguatan Ketahanan dan Inovasi Keuangan"
PS_SUB = "Sub-PS · Manajemen Risiko"


def _run(run, size=18, bold=False, color=WHITE):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _bg(slide):
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill(fill, BG)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
    _fill(bar, TEAL)


def _text(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    _run(r, size=size, bold=bold, color=color)


def _bullets(slide, l, t, w, h, items, size=15, color=SOFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = f"•  {item}"
        _run(r, size=size, color=color)


def _card(slide, l, t, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    _fill(sh, CARD)
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)


def _footer(slide, n, total, part: str):
    _text(slide, Inches(0.45), Inches(7.12), Inches(9), Inches(0.28),
          f"Nexus Cyber · PIDI Capstone · {part}", size=10, color=MUTED)
    _text(slide, Inches(11.5), Inches(7.12), Inches(1.5), Inches(0.28),
          f"{n}/{total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def _badge(slide, text, color, top=Inches(0.22)):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), top, Inches(2.8), Inches(0.38))
    _fill(sh, color)
    sh.line.fill.background()
    tf = sh.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    _run(r, size=10, bold=True, color=WHITE)


def _title(slide, eyebrow, title, subtitle=None):
    _text(slide, Inches(0.5), Inches(0.28), Inches(9.5), Inches(0.32), eyebrow, size=11, bold=True, color=TEAL)
    _text(slide, Inches(0.5), Inches(0.58), Inches(12), Inches(0.65), title, size=28, bold=True)
    if subtitle:
        _text(slide, Inches(0.5), Inches(1.18), Inches(12), Inches(0.42), subtitle, size=14, color=MUTED)


def _table(slide, l, t, w, headers, rows, col_w=None, row_h=Inches(0.42), hdr_size=11, cell_size=10):
    ncols = len(headers)
    nrows = len(rows) + 1
    if col_w is None:
        col_w = [w / ncols] * ncols
    y = t
    # header
    x = l
    for i, h in enumerate(headers):
        _card(slide, x, y, col_w[i], row_h)
        _text(slide, x + Inches(0.06), y + Inches(0.08), col_w[i] - Inches(0.1), row_h,
              h, size=hdr_size, bold=True, color=TEAL)
        x += col_w[i]
    y += row_h
    for row in rows:
        x = l
        for i, cell in enumerate(row):
            _card(slide, x, y, col_w[i], row_h)
            col = TEAL if i == 0 else SOFT
            _text(slide, x + Inches(0.06), y + Inches(0.08), col_w[i] - Inches(0.1), row_h,
                  str(cell), size=cell_size, color=col)
            x += col_w[i]
        y += row_h


def new_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    return s


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    builders: list = []

    def add(fn):
        builders.append(fn)

    # ── COVER ──
    def s0(slide, n, total):
        _badge(slide, "PIDI 2026", CORE_BADGE)
        _text(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.35),
              "FINAL SUBMISSION · CAPSTONE PROJECT", size=13, bold=True, color=TEAL)
        _text(slide, Inches(0.6), Inches(2.1), Inches(12), Inches(0.9), "Nexus Cyber", size=44, bold=True)
        _text(slide, Inches(0.6), Inches(3.1), Inches(11), Inches(0.5), PS_MAIN, size=16, color=AMBER)
        _text(slide, Inches(0.6), Inches(3.65), Inches(11), Inches(0.5), PS_SUB, size=16, color=TEAL)
        _text(slide, Inches(0.6), Inches(4.4), Inches(11), Inches(1.2),
              "Edge Antibody Cowork (GaaS) — wasit risiko kanal web/API\n"
              "Evidence over claim · Built vs Planned · Human L0/L1",
              size=17, color=SOFT)
        _text(slide, Inches(0.6), Inches(6.0), Inches(11), Inches(0.4),
              "Part A: 14 slide inti presentasi  ·  Part B: Additional untuk Q&A",
              size=13, color=MUTED)
        _footer(slide, n, total, "Cover")

    add(s0)

    # ═══ PART A — 14 CORE ═══
    def core(n_core, title, eyebrow, fn):
        def wrapper(slide, n, total):
            _badge(slide, f"CORE {n_core}/14", CORE_BADGE)
            _title(slide, eyebrow, title)
            fn(slide)
            _footer(slide, n, total, f"Part A · Slide {n_core}")
        add(wrapper)

    # 1 Solution at a Glance
    def body1(slide):
        _card(slide, Inches(0.5), Inches(1.75), Inches(12.2), Inches(4.8))
        _text(slide, Inches(0.75), Inches(2.0), Inches(11.5), Inches(0.4),
              "Value proposition (1 kalimat)", size=14, bold=True, color=TEAL)
        _text(slide, Inches(0.75), Inches(2.5), Inches(11.5), Inches(0.9),
              "Wasit terukur untuk kanal digital: ukur celah HTTP → kendalikan di tepi → "
              "uji replay → serahkan artefak risiko — bukan PDF scanner saja.",
              size=18, bold=True, color=WHITE)
        _bullets(slide, Inches(0.75), Inches(3.6), Inches(11.5), Inches(2.8), [
            f"Problem Statement: {PS_MAIN}",
            f"Sub-Problem: {PS_SUB} (identifikasi · ukur · pantau · kurangi risiko operasional/keamanan)",
            "Target: pemilik risiko kanal (fintech/ITSK), integrator web keuangan, IT corporat",
            "Offtaker/adopter: direksi IT, compliance ops, agensi yang deploy kanal untuk klien keuangan",
            "Entry funnel UMKM (Channel Starter) — bukan fokus PS, hanya pintu volume",
        ], size=14)

    core(1, "Solution at a Glance", "Slide 1 · Wajib PIDI", body1)

    # 2 Problem
    def body2(slide):
        _bullets(slide, Inches(0.55), Inches(1.75), Inches(6.2), Inches(5), [
            "Kanal digital (portal, API, onboarding) dirilis cepat → permukaan HTTP bertambah.",
            "Setiap rilis menambah risiko operasional: deface, abuse form, IDOR, gap origin vs tepi.",
            "Institusi keuangan wajib bukti pengendalian — bukan hanya temuan audit statis.",
            "Insiden kanal publik merusak kepercayaan konsumen & stabilitas sistem (PS1).",
        ], size=15)
        _card(slide, Inches(6.9), Inches(1.75), Inches(5.8), Inches(5))
        _text(slide, Inches(7.15), Inches(2.0), Inches(5.3), Inches(0.35),
              "Konsekuensi bila tidak diselesaikan", size=14, bold=True, color=ROSE)
        _bullets(slide, Inches(7.1), Inches(2.5), Inches(5.4), Inches(4), [
            "Firefight manual pasca-insiden (biaya ops + reputasi).",
            "Keputusan direksi tanpa jejak residual jujur.",
            "Regulator/pemilik risiko sulit verifikasi “sudah aman”.",
            "Gap WAF vs origin tidak terdeteksi → hijau palsu.",
        ], size=14)

    core(2, "Problem & Why It Matters", "Slide 2 · Wajib PIDI", body2)

    # 3 Validation & Root Cause
    def body3(slide):
        _text(slide, Inches(0.5), Inches(1.72), Inches(12), Inches(0.35),
              "Root cause: siklus manajemen risiko kanal terputus (identifikasi–pengendalian–pemantauan)",
              size=14, bold=True, color=AMBER)
        _table(slide, Inches(0.5), Inches(2.15), Inches(12.2),
               ["Sumber bukti", "Temuan validasi (contoh/pilot)"],
               [
                   ["Observasi alur kerja IT", "Scanner → PDF; WAF → log; tidak ada loop uji replay"],
                   ["Wawancara stakeholder*", "Pemilik kanal minta bukti pengendalian berulang, bukan one-shot audit"],
                   ["Data operasional publik", "Deface/abus kanal UMKM & institusi masih sering dilaporkan (media/BSSN advisories)"],
                   ["Regulasi & praktik industri", "Manajemen risiko ITSK/fintech mensyaratkan identifikasi–pengendalian–pemantauan"],
               ],
               col_w=[Inches(3.5), Inches(8.7)], row_h=Inches(0.72), cell_size=12)
        _text(slide, Inches(0.5), Inches(6.35), Inches(12), Inches(0.35),
              "*Lengkapi transkrip wawancara di Additional slide A16 — jangan klaim tanpa bukti nyata tim Anda.",
              size=11, color=MUTED)

    core(3, "Problem Validation & Root Cause", "Slide 3 · Wajib PIDI", body3)

    # 4 Solution & Core Use Case
    def body4(slide):
        _text(slide, Inches(0.5), Inches(1.72), Inches(12), Inches(0.35),
              "Core use case: satu Job Cowork menutup loop wasit pada satu host kanal",
              size=14, bold=True, color=TEAL)
        steps = ["Scope host", "Ukur delta", "Draft antibodi", "Approve L0/L1", "Replay uji", "Artefak MD/JSON"]
        x = 0.35
        for i, s in enumerate(steps):
            _card(slide, Inches(x), Inches(2.2), Inches(2.05), Inches(1.5))
            _text(slide, Inches(x + 0.1), Inches(2.45), Inches(1.85), Inches(1.0), f"{i+1}. {s}", size=13, bold=True, color=TEAL)
            x += 2.15
        _bullets(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(2.8), [
            "Before: celah origin terbuka / tepi tidak terbukti · After: status CLOSED_OK atau CLOSED_GAP (jujur)",
            "Alur A (always-on): trafik → WAF :8080 → origin/403 + Reflex + antibodi cache",
            "Alur B (Job): orkestrasi NEX-RED jobs/ + bridge :3004 + sync PostgreSQL",
            "Alur C (artefak): export MD/JSON untuk pemilik risiko kanal — bukan login SOC pelanggan",
        ], size=14)

    core(4, "Solution & Core Use Case", "Slide 4 · Wajib PIDI", body4)

    # 5 Differentiation
    def body5(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Dimensi", "Scanner/PDF", "WAF/CDN", "SOC manual", "Nexus GaaS"],
               [
                   ["Bukti origin vs tepi", "Tidak", "Sebagian", "Proses berat", "Defense delta"],
                   ["Uji replay jujur", "Tidak", "Jarang", "Ad hoc", "replay_missed→GAP"],
                   ["Artefak risiko", "PDF statis", "Log", "Tiket", "Job MD/JSON"],
                   ["Human oversight", "N/A", "N/A", "Analyst", "L0/L1 gerbang"],
                   ["Harga entry UMKM", "Mahal sekali", "Menengah", "Tidak feasible", "Starter terpisah*"],
               ],
               col_w=[Inches(2.4), Inches(2.3), Inches(2.3), Inches(2.5), Inches(2.8)],
               row_h=Inches(0.78), cell_size=11)
        _text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.3),
              "*Channel Starter lab v0.1 — bukan Loop GaaS di Rp 20rb. Detail di Additional A7–A10.",
              size=11, color=MUTED)

    core(5, "Value Proposition & Differentiation", "Slide 5 · Wajib PIDI", body5)

    # 6 Prototype State
    def body6(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Komponen", "Status", "Bukti"],
               [
                   ["Gateway WAF + Reflex (Alur A)", "BUILT", "deploy-local · go run gateway · :8080"],
                   ["Job Cowork NEX-RED (Alur B)", "BUILT", "NEX-RED/jobs/ · bridge :3004 · nexred.py job"],
                   ["Artefak MD/JSON (Alur C)", "BUILT", "jobs/data/artifacts/ · export API"],
                   ["Operator GaaS Console", "BUILT", "nexus-admin-dashboard :3001"],
                   ["Channel Portal v0.1", "BUILT (lab)", "nexus-channel-portal :3003 · WA manual"],
                   ["Channel Starter deploy massal", "PLANNED", "channel-starter/ lab v0.1"],
                   ["Billing Midtrans otomatis", "PLANNED", "ditunda — fail-closed"],
                   ["Edge on-prem B2G packaging", "PLANNED", "pitching · belum DC produksi"],
                   ["eBPF / SOC 24/7 / pentest exploit", "NOT CLAIMED", "LIMITATIONS.md"],
               ],
               col_w=[Inches(4.2), Inches(1.8), Inches(6.3)],
               row_h=Inches(0.52), cell_size=10)
        _text(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.3),
              "Level prototype: Functional Prototype (Level 2) — core use case dapat di-demo end-to-end di lab.",
              size=12, bold=True, color=TEAL)

    core(6, "Prototype & Current Product State", "Slide 6 · Wajib PIDI", body6)

    # 7 3-Layer Technology
    def body7(slide):
        layers = [
            ("LAYER 1 · User Experience",
             "Pemilik risiko / operator: portal order · site di belakang WAF · Operator Console (internal) · artefak Job"),
            ("LAYER 2 · System Logic",
             "Input HTTP jinak → defense delta (WAF vs origin) → draft antibodi → L0/L1 approve → vaccine-probe + replay → status Job"),
            ("LAYER 3 · Technical Architecture",
             "Caddy/tunnel → Gateway :8080 → origin · Control :8081 · PG/Redis · NEX-RED :3004 · nex-ai-protect/reflex (opsional)"),
        ]
        y = 1.75
        for title, body in layers:
            _card(slide, Inches(0.5), Inches(y), Inches(12.2), Inches(1.55))
            _text(slide, Inches(0.75), Inches(y + 0.15), Inches(11.5), Inches(0.35), title, size=14, bold=True, color=TEAL)
            _text(slide, Inches(0.75), Inches(y + 0.55), Inches(11.5), Inches(0.85), body, size=13, color=SOFT)
            y += 1.7
        _text(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.25),
              "Diagram detail & data flow → Additional A2, A5", size=11, color=MUTED)

    core(7, "How the Technology Works", "Slide 7 · Wajib PIDI · 3-Layer", body7)

    # 8 Technical Testing
    def body8(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Kategori uji", "Metode", "Hasil lab (contoh)", "Catatan"],
               [
                   ["Defense delta", "Request identik WAF vs origin", "Label waf_blocked / origin_open / replay_*", "NEX-RED Sprint 1"],
                   ["Antibody loop", "Replay 403 di tepi", "replay_held vs replay_missed", "Sprint 2 · CLOSED_GAP jujur"],
                   ["Job lifecycle", "Orkestrasi Job Cowork", "OPEN→…→CLOSED_OK/GAP", "test_job_cowork.py"],
                   ["Gateway proxy", "Unit + integration tests", "go test ./internal/proxy", "CI pre-push"],
                   ["Security surface", "Telemetry di :8080", "GET /api/telemetry → 404", "CAPABILITIES.md"],
                   ["Hotspot harness", "Port privat tertutup", ":8081/:3001 closed from hotspot", "NEX_RED_HOTSPOT_HARNESS"],
               ],
               col_w=[Inches(2.5), Inches(3.0), Inches(3.5), Inches(3.3)],
               row_h=Inches(0.62), cell_size=10)
        _text(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.35),
              "Log lengkap & error analysis → Additional A6. Bukan load test production DC.",
              size=11, color=MUTED)

    core(8, "Technical Testing & Performance", "Slide 8 · Wajib PIDI", body8)

    # 9 Impact Evidence
    def body9(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Metric", "Baseline (tanpa wasit)", "Hasil pilot/lab", "Target implementasi"],
               [
                   ["Loop risiko kanal tertutup", "Manual / tidak terukur", "1 Job = 1 siklus status jujur", "Loop berkala per rilis"],
                   ["Deteksi gap origin vs tepi", "Tidak terlihat", "Defense delta terlabel", "0 origin_open tanpa residual"],
                   ["Bukti replay antibodi", "N/A", "replay_held / replay_missed tercatat", "replay_missed → GAP, bukan OK"],
                   ["Waktu artefak ke pemilik risiko", "Minggu (audit)", "48–72 jam Job + export", "Same-day setelah approve"],
                   ["False sense of security", "Tinggi (PDF hijau)", "Diturunkan via CLOSED_GAP", "Transparansi residual"],
               ],
               col_w=[Inches(3.2), Inches(2.8), Inches(3.2), Inches(3.1)],
               row_h=Inches(0.72), cell_size=11)
        _text(slide, Inches(0.5), Inches(6.45), Inches(12), Inches(0.4),
              "Metodologi: lab portfolio.nexus-lab.test di belakang WAF; metrik dari Job Cowork + NEX-RED test suite.",
              size=11, color=MUTED)

    core(9, "Impact & Evidence of Effectiveness", "Slide 9 · Wajib PIDI", body9)

    # 10 Market Validation
    def body10(slide):
        _bullets(slide, Inches(0.55), Inches(1.75), Inches(6.0), Inches(5), [
            "ICP B2B: fintech, ITSK, integrator web keuangan (COWORK_B2B.md).",
            "Adoption ecosystem: pemilik risiko kanal + operator Nexus (managed).",
            "Validasi v1: pembayaran/manual WA + order portal (62895603358692).",
            "Pilot distribusi: PC operator 24/7 + tunnel — bukan SLA DC.",
            "LOI formal: [ISI jika sudah ada] — jujur jika belum.",
            "User testing: [ISI feedback nyata tim Anda] — lihat template A16.",
        ], size=14)
        _card(slide, Inches(6.7), Inches(1.75), Inches(6.0), Inches(5))
        _text(slide, Inches(6.95), Inches(2.0), Inches(5.5), Inches(0.35),
              "Willingness to pay (ilustrasi pilot)", size=14, bold=True, color=TEAL)
        _bullets(slide, Inches(6.9), Inches(2.5), Inches(5.5), Inches(4), [
            "Job Cowork: Rp 200.000 / host (one-shot)",
            "Loop GaaS: Rp 300.000 / bulan / host",
            "B2G Edge: Rp 18 jt/tahun + Loop 3,5 jt/bln",
            "TAM/SOM tanpa wawancara = tidak cukup (aturan PIDI)",
        ], size=13)

    core(10, "Market / User / Offtaker Validation", "Slide 10 · Wajib PIDI", body10)

    # 11 Adoption & Sustainability
    def body11(slide):
        _card(slide, Inches(0.5), Inches(1.75), Inches(5.8), Inches(4.6))
        _text(slide, Inches(0.75), Inches(2.0), Inches(5.3), Inches(0.35), "Go-to-market", size=14, bold=True, color=TEAL)
        _bullets(slide, Inches(0.7), Inches(2.5), Inches(5.4), Inches(3.5), [
            "Portal segmen → WA order → deploy pilot",
            "Upsell: Starter → Tepi → Job → Loop → On-prem",
            "Demo lewat PROTECTED_HOST (bukan origin langsung)",
            "Artefak = deliverable ke pemilik risiko",
        ], size=13)
        _card(slide, Inches(6.5), Inches(1.75), Inches(6.2), Inches(4.6))
        _text(slide, Inches(6.75), Inches(2.0), Inches(5.8), Inches(0.35), "Unit economics (ringkas)", size=14, bold=True, color=AMBER)
        _table(slide, Inches(6.55), Inches(2.45), Inches(6.1),
               ["Paket", "Jual", "Margin"],
               [
                   ["Job Cowork", "Rp 200rb", "~63%"],
                   ["Loop GaaS", "Rp 300rb/bln", "~68%"],
                   ["B2G Edge+Loop/thn", "~Rp 60jt", "~75%"],
               ],
               col_w=[Inches(2.2), Inches(2.0), Inches(1.9)],
               row_h=Inches(0.48), cell_size=11)
        _text(slide, Inches(6.75), Inches(4.2), Inches(5.8), Inches(1.8),
              "Detail COGS per segmen → Additional A7–A10\nAsumsi pilot PC+tunnel (infra ≈ 0)",
              size=12, color=MUTED)

    core(11, "Adoption & Sustainability Path", "Slide 11 · Wajib PIDI", body11)

    # 12 Team
    def body12(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Peran", "Skill", "Kontribusi nyata", "Q&A ownership"],
               [
                   ["[Nama] Lead Presenter", "Product / GTM", "Problem, impact, closing", "Slide 1–2, 9–10"],
                   ["[Nama] Technical Lead", "Go, NEX-RED, gateway", "Job Cowork, demo live", "Slide 6–8, A2–A6"],
                   ["[Nama] Business Rep", "Market, BMC", "Validasi, pricing, adopsi", "Slide 10–11, A7–A11"],
                   ["[Nama] Domain Support", "Security / compliance", "L0/L1, limitasi jujur", "Slide 14, A13, A17"],
               ],
               col_w=[Inches(2.8), Inches(2.5), Inches(3.5), Inches(3.5)],
               row_h=Inches(0.78), cell_size=11)
        _text(slide, Inches(0.5), Inches(6.35), Inches(12), Inches(0.45),
              "ISI nama & kontribusi nyata tim Anda. Aturan PIDI: dilarang hanya foto + universitas.",
              size=12, bold=True, color=AMBER)

    core(12, "Team & Execution Readiness", "Slide 12 · Wajib PIDI", body12)

    # 13 Roadmap
    def body13(slide):
        _table(slide, Inches(0.45), Inches(1.72), Inches(12.3),
               ["Milestone", "Outcome", "Resource", "Timeline"],
               [
                   ["M0 (now)", "Demo Job + portal + lab WAF", "PC operator + repo", "Hackathon"],
                   ["M1", "3–5 klien Loop hosted pilot", "Operator + WA sales", "+3 bulan"],
                   ["M2", "Packaging Edge berlisensi v1", "Dev + legal template", "+6 bulan"],
                   ["M3", "1 pilot on-prem staging B2G", "Install remote + Loop", "+9–12 bulan"],
               ],
               col_w=[Inches(2.0), Inches(4.0), Inches(3.5), Inches(2.8)],
               row_h=Inches(0.72), cell_size=11)
        _text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.0),
              "Biaya ilustrasi R&D+ops: Rp 8–15 jt/bln (founder blended) — sesuaikan aktual.\n"
              "Payback: 1× B2G tahun-1 ATAU ~40–80 Loop stabil (PRICING_UNIT_ECONOMICS.md).",
              size=13, color=SOFT)

    core(13, "Roadmap to Implementation", "Slide 13 · Wajib PIDI", body13)

    # 14 Risks
    def body14(slide):
        _table(slide, Inches(0.35), Inches(1.68), Inches(12.5),
               ["Kategori", "Risiko", "Mitigasi"],
               [
                   ["Teknis", "replay_missed / origin_open", "CLOSED_GAP jujur · residual tertulis · Loop"],
                   ["Adopsi", "Ekspektasi SOC 24/7 murah", "Copy portal + LIMITATIONS · L0/L1"],
                   ["Regulasi", "Klaim sertifikasi palsu", "BSSN/POJK = framing · bukan approve"],
                   ["Operasional", "PC single point of failure", "Roadmap VPS · backup demo video"],
                   ["Data", "Validasi pasar tipis", "Wawancara + LOI target · template A16"],
                   ["B2G", "Packaging belum produksi", "Pitching jujur · milestone M2–M3"],
               ],
               col_w=[Inches(1.8), Inches(4.5), Inches(6.2)],
               row_h=Inches(0.68), cell_size=10)
        _text(slide, Inches(0.5), Inches(6.55), Inches(12), Inches(0.35),
              "Risiko per segmen detail → Additional A10 (finance) · mitigasi UMKM/B2G terpisah jika juri tanya.",
              size=11, color=MUTED)

    core(14, "Key Risks & Next Priorities", "Slide 14 · Wajib PIDI", body14)

    # ═══ PART B — ADDITIONAL (Q&A) ═══
    def add_slide(title, eyebrow, items, table=None, note=None):
        def wrapper(slide, n, total):
            _badge(slide, "ADDITIONAL · Q&A", ADD_BADGE)
            _title(slide, eyebrow, title)
            if table:
                headers, rows, col_w, rh = table
                _table(slide, Inches(0.45), Inches(1.72), Inches(12.3), headers, rows, col_w=col_w, row_h=rh, cell_size=10)
            elif items:
                _bullets(slide, Inches(0.55), Inches(1.75), Inches(12), Inches(5.2), items, size=14)
            if note:
                _text(slide, Inches(0.5), Inches(6.6), Inches(12), Inches(0.35), note, size=11, color=MUTED)
            _footer(slide, n, total, "Part B · Q&A Support")
        add(wrapper)

    # B divider
    def s_div(slide, n, total):
        _badge(slide, "PART B", ADD_BADGE, top=Inches(0.22))
        _text(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(0.5),
              "ADDITIONAL SLIDES", size=36, bold=True, color=AMBER)
        _text(slide, Inches(0.6), Inches(3.3), Inches(11), Inches(1.5),
              "Tidak dipresentasikan penuh dalam 10 menit.\n"
              "Digunakan saat Q&A — jawab dengan data & rujukan slide, bukan abstrak.",
              size=20, color=SOFT)
        _bullets(slide, Inches(0.6), Inches(4.8), Inches(11), Inches(2), [
            "Teknis: wasit, Job, testing, keamanan, AI 8 poin",
            "Bisnis: unit economics per segmen, BMC, ROI",
            "Validasi: template wawancara, checklist eligibility",
        ], size=16)
        _footer(slide, n, total, "Part B · Pembuka")
    add(s_div)

    add_slide("Defense Delta & Label Wasit Jujur",
              "A1 · Data Teknis",
              None,
              table=(
                  ["Label", "Arti", "Implikasi Job"],
                  [
                      ["waf_blocked", "Tepi menahan", "Pengendalian di gateway OK"],
                      ["origin_open", "Origin terbuka tanpa WAF", "Residual → CLOSED_GAP"],
                      ["both_held", "Keduanya menahan", "Positif"],
                      ["replay_held", "Replay tetap 403", "Antibody loop OK"],
                      ["replay_missed", "Replay lolos", "WAJIB GAP — bukan hijau palsu"],
                  ],
                  [Inches(2.5), Inches(5.0), Inches(4.8)],
                  Inches(0.62),
              ),
              note="Sumber: NEX-RED defense delta · PRODUCT_MODEL §4 · aturan CLOSED_OK")

    add_slide("Job Cowork — Status & Lifecycle",
              "A2 · Data Teknis",
              None,
              table=(
                  ["Status", "Arti", "Aksi manusia"],
                  [
                      ["OPEN", "Job dibuka", "Set scope"],
                      ["MEASURING", "Wasit jalan", "—"],
                      ["PENDING_APPROVAL", "Draft antibodi", "Approve L0/L1"],
                      ["VERIFYING", "Replay/vaccine", "—"],
                      ["CLOSED_OK", "Loop OK", "Export artefak"],
                      ["CLOSED_GAP", "Masih gap", "Residual jujur"],
                      ["PARTIAL", "Agen gagal sebagian", "Tutup jujur"],
                  ],
                  [Inches(2.8), Inches(4.5), Inches(5.0)],
                  Inches(0.55),
              ))

    add_slide("AI / GaaS — 8 Poin Teknis (PIDI wajib jika klaim AI)",
              "A3 · AI Transparency",
              [
                  "1. Input: payload HTTP jinak, header, path, defense delta features",
                  "2. Sumber: log gateway lab, request wasit NEX-RED, origin loopback/RFC1918 saja",
                  "3. Preprocessing: normalisasi label wasit, immune memory per host (JSON)",
                  "4. Output: klasifikasi risiko (regex reflex) + reasoning opsional (nex-ai-protect JSON)",
                  "5. Trigger aksi: draft antibodi → gerbang L0 (artefak) / L1 (pasang tepi)",
                  "6. Metrik: recall kelas wasit (NEX-RED test suite); bukan klaim accuracy produksi massal",
                  "7. Limitasi: bukan pentest exploit; reflex=regex default; model opsional",
                  "8. Human oversight: L0/L1 wajib — unban/DNS/disk TIDAK otonom (PRODUCT_MODEL §6)",
                  "Pembedaan: arsitektur Job + wasit milik Nexus; model GGUF milik pemilik — bukan API OpenAI generic tanpa logika",
              ])

    add_slide("Keamanan Lab & Posture",
              "A4 · Security",
              None,
              table=(
                  ["Kontrol", "Implementasi lab", "Bukti"],
                  [
                      ["SOC tidak publik", ":3001/:8081 localhost", "DISTRIBUTION_PILOT"],
                      [":8080 telemetry", "404 / auth", "test_live_http"],
                      ["Origin direct NEX-RED", "HTTP privat saja", "netaddr guard"],
                      ["Cookie session WAF", "POST tanpa cookie → 401", "proxy_core.go"],
                      ["Tunnel juri", "Hanya site WAF", "JURY_PUBLIC_ACCESS.md"],
                  ],
                  [Inches(3.0), Inches(4.5), Inches(4.8)],
                  Inches(0.58),
              ))

    add_slide("Bukti Pengujian Otomatis (Repo)",
              "A5 · Testing Log",
              [
                  "NEX-RED: test_job_cowork, test_live_http (defense delta, replay, antibody)",
                  "Gateway: go test ./internal/proxy, ./tests (pre-push hook)",
                  "Aturan: replay_missed → CLOSED_GAP (test_job_cowork.TestJobClosure)",
                  "Juice Shop lab: recall kelas AUTH/INJ — equal_to_shannon_strix = false (jujur)",
                  "Cara reproduksi: deploy-local START-FOR-JURY.bat · nexred.py job run",
                  "Log lengkap: jalankan pytest / go test di mesin presenter — screenshot output untuk juri",
              ])

    add_slide("Unit Economics — Ringkasan Pilot",
              "A6 · Finance",
              None,
              table=(
                  ["Paket", "Jual", "COGS", "Margin"],
                  [
                      ["Job Cowork", "Rp 200rb", "Rp 75rb", "63%"],
                      ["Loop GaaS", "Rp 300rb/bln", "Rp 95rb", "68%"],
                      ["UMKM Website", "Rp 20rb/bln", "Rp 8,5rb", "58%"],
                      ["B2G Edge/thn", "Rp 18jt", "Rp 2,4jt", "87%"],
                      ["B2G Loop/bln", "Rp 3,5jt", "Rp 1,05jt", "70%"],
                  ],
                  [Inches(3.0), Inches(2.5), Inches(2.5), Inches(1.5)],
                  Inches(0.52),
              ),
              note="PRICING_UNIT_ECONOMICS.md · asumsi PC+tunnel · bukan laporan audit")

    add_slide("Finance · UMKM & Sekolah",
              "A7 · Finance Detail",
              None,
              table=(
                  ["Segmen · paket", "Jual/bln", "Margin"],
                  [
                      ["UMKM · Website Aman", "Rp 20rb", "58%"],
                      ["UMKM · GaaS entry", "Rp 35rb", "47%"],
                      ["UMKM · Pagar (sudah web)", "Rp 15rb", "60%"],
                      ["Sekolah · struktur sama", "Rp 20–35rb", "47–58%"],
                  ],
                  [Inches(4.5), Inches(2.5), Inches(2.0)],
                  Inches(0.55),
              ),
              note="Volume UMKM = margin % OK, laba absolut kecil — batch support wajib")

    add_slide("Finance · Startup & Corporat",
              "A8 · Finance Detail",
              None,
              table=(
                  ["Paket", "Jual", "Margin"],
                  [
                      ["Landing+pagar", "Rp 45rb/bln", "64%"],
                      ["Landing+Tepi", "Rp 75rb/bln", "63%"],
                      ["Job Wasit", "Rp 200rb", "63%"],
                      ["Loop GaaS", "Rp 300rb/bln", "68%"],
                      ["Corporat on-prem", "Quote", "Custom"],
                  ],
                  [Inches(4.0), Inches(3.5), Inches(2.0)],
                  Inches(0.55),
              ))

    add_slide("Finance · Pemerintah / B2G",
              "A9 · Finance Detail",
              [
                  "Lisensi Edge On-Prem: Rp 18.000.000 / tahun (1 zona/DC ilustrasi)",
                  "Loop On-Prem wajib: Rp 3.500.000 / bulan",
                  "Tahun 1 tipikal 1 DC: jual ~Rp 60 jt · COGS ~Rp 15 jt · laba ~Rp 45 jt",
                  "Source & SOC control plane TIDAK diserahkan (COWORK_B2G.md)",
                  "Pitching ≠ pengadaan SIPLah/E-Katalog selesai",
                  "Mitigasi: packaging M2 · kontrak NDA · pilot staging dulu",
              ])

    add_slide("Business Model Canvas (detail)",
              "A10 · BMC",
              None,
              table=(
                  ["Blok", "Isi"],
                  [
                      ["Value Prop", "Wasit HTTP + entry site · loop jujur"],
                      ["Customer", "Fintech · integrator · UMKM upsell · B2G DC"],
                      ["Channels", "Portal · WA · demo lab · pitch"],
                      ["Revenue", "Starter · Job · Loop · Edge license"],
                      ["Cost", "Operator time · listrik PC · R&D"],
                      ["Key Resources", "Gateway · NEX-RED · Job PG · IP proses"],
                  ],
                  [Inches(3.5), Inches(8.8)],
                  Inches(0.58),
              ))

    add_slide("Limitations & Klaim yang TIDAK Kami Buat",
              "A11 · Transparency (PIDI prinsip #6)",
              [
                  "Bukan SOC otonom 24/7 — operator L0/L1",
                  "Bukan pentest exploit / Shannon parity",
                  "Bukan eBPF/XDP nyata — eBPF stub",
                  "Bukan GRC bank penuh — irisan kanal web/API",
                  "Bukan approve regulator POJK/BSSN",
                  "Channel Starter ≠ Loop GaaS di Rp 20rb",
                  "B2G on-prem = pitching; packaging produksi belum",
                  "Billing Midtrans massal ditunda (fail-closed)",
                  "Satu PROTECTED_HOST per instance — bukan multi-tenant massal",
              ],
              note="LIMITATIONS.md · AGENTS.md — juri menghargai kejujuran")

    add_slide("Demo SOP — 60–120 Detik + Backup 3 Lapis",
              "A12 · Showcasing",
              [
                  "Alur demo: Operator → Job Cowork → approve L0 → artefak → site di belakang WAF",
                  "Tier 1: Live demo (START-FOR-JURY.bat · portfolio.nexus-lab.test)",
                  "Tier 2: Video rekaman 60–90 detik di harddisk lokal (WAJIB siapkan)",
                  "Tier 3: Screenshot flow di slide ini jika jaringan gagal",
                  "DILARANG: feature dump semua menu · login panjang di panggung · klaim fitur roadmap",
                  "Alokasi offline: ~4 menit demo (PANDUAN-SHOWCASING § Tahap 2)",
              ])

    add_slide("Distribusi Pilot — PC + Tunnel",
              "A13 · Ops",
              [
                  "Host: PC operator 24/7 (bukan VPS dulu) — DISTRIBUTION_PILOT.md",
                  "Publik: Cloudflare Tunnel → Caddy :80 → WAF :8080 → origin",
                  "TIDAK di-tunnel: SOC :3001, control :8081, DB, NEX-RED mentah",
                  "Copy jujur ke juri: pilot infra operator — bukan SLA data center",
                  "Script: deploy-local/jury/START-FOR-JURY.bat",
              ])

    add_slide("Template Validasi Lapangan (ISI data tim)",
              "A14 · Validation Evidence",
              [
                  "Wawancara #1: [Role · Instansi] — Pain: bukti pengendalian berulang",
                  "Wawancara #2: [Role · Instansi] — Pain: gap WAF vs origin tidak terlihat",
                  "Wawancara #3: [UMKM/fintech] — Willingness: Job vs Loop",
                  "Observasi: alur deploy rilis → tidak ada replay otomatis",
                  "LOI: [nama instansi / status: draft / signed / belum]",
                  "⚠ ISI dengan data nyata tim Anda sebelum final submit — jangan presentasi template kosong",
              ],
              note="Checklist PIDI #8: bukti di luar asumsi internal tim")

    add_slide("Checklist Eligibility — Jawaban Tim",
              "A15 · Q&A Cheat Sheet",
              None,
              table=(
                  ["#", "Pertanyaan PIDI", "Jawaban singkat + rujuk slide"],
                  [
                      ["1", "Bukti masalah?", "Ya — slide 2–3 + A14"],
                      ["2", "Alignment PS?", "PS1 Manajemen Risiko — slide 1"],
                      ["3", "Tanpa buzzword?", "Job Cowork = wasit loop — slide 4"],
                      ["4", "Core use case built?", "Ya Level 2 — slide 6 + demo"],
                      ["5", "Alur input→output?", "3-layer slide 7 + A1–A2"],
                      ["6", "Data testing?", "Ya — slide 8 + A5"],
                      ["7", "Dampak metrik?", "Tabel slide 9"],
                      ["8", "Validasi pasar?", "Pilot WA — slide 10 + A14"],
                      ["9", "Differentiation?", "Slide 5"],
                      ["10", "Tim jelas?", "Slide 12 — isi nama"],
                      ["11", "Roadmap?", "Slide 13"],
                      ["12", "Built vs planned?", "Slide 6 + A11"],
                  ],
                  [Inches(0.6), Inches(5.5), Inches(6.1)],
                  Inches(0.48),
              ))

    add_slide("Referensi Repo & Dokumen Hidup",
              "A16 · Appendix Index",
              [
                  "docs/PRODUCT_MODEL.md — model GaaS + Alur A/B/C",
                  "docs/PRICING_UNIT_ECONOMICS.md — angka finance",
                  "docs/COWORK_B2B.md · COWORK_B2G.md — GTM",
                  "docs/LIMITATIONS.md · CAPABILITIES.md — kejujuran",
                  "docs/DISTRIBUTION_PILOT.md · JURY_PUBLIC_ACCESS.md — demo",
                  "NEX-RED/jobs/ · nexus-core-gateway/ · nexus-admin-dashboard/",
                  "Repo: github.com/Thbetyfu/NEXUS-CYBER-FASE3",
              ])

    # Closing presentasi
    def s_close(slide, n, total):
        _badge(slide, "CLOSING", CORE_BADGE)
        _text(slide, Inches(0.6), Inches(2.2), Inches(12), Inches(0.5),
              "Nexus Cyber siap pilot — dengan bukti, bukan klaim.", size=32, bold=True)
        _bullets(slide, Inches(0.6), Inches(3.2), Inches(11), Inches(2.5), [
            "PS1 · Manajemen Risiko — loop wasit kanal digital",
            "Functional prototype · defense delta · replay jujur",
            "Milestone berikutnya: klien Loop + packaging Edge",
            "Terima kasih — kami siap Q&A (rujuk slide Additional)",
        ], size=18)
        _footer(slide, n, total, "Closing")
    add(s_close)

    total = len(builders)
    for i, fn in enumerate(builders, 1):
        fn(new_slide(prs), i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    core_count = 14
    add_count = total - core_count - 2  # minus cover and divider and closing approx
    print(f"Wrote {OUT}")
    print(f"Total slides: {total} (Cover + Core {core_count} + Additional ~{add_count} + divider/closing)")


if __name__ == "__main__":
    build()
