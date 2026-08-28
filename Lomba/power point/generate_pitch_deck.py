#!/usr/bin/env python3
"""Generate Nexus Cyber investor pitch deck (python-pptx)."""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "Nexus-Cyber-Investor-Pitch.pptx"

# Brand — ink + teal (selaras portal, hindari ungu AI-generic)
BG = RGBColor(0x07, 0x0B, 0x10)
CARD = RGBColor(0x0E, 0x16, 0x1C)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
TEAL_DIM = RGBColor(0x14, 0xB8, 0xA6)
AMBER = RGBColor(0xF5, 0xB3, 0x01)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
ROSE = RGBColor(0xFB, 0x71, 0x85)
LINE = RGBColor(0x1E, 0x29, 0x3B)

W, H = Inches(13.333), Inches(7.5)  # 16:9 widescreen


def _set_run(run, size=18, bold=False, color=WHITE, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _bg(slide):
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    _fill(fill, BG)
    # subtle top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.08))
    _fill(bar, TEAL)


def _textbox(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size=size, bold=bold, color=color)
    return box


def _bullets(slide, left, top, width, height, items, size=16, color=SOFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        _set_run(run, size=size, color=color)
    return box


def _card(slide, left, top, width, height):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    _fill(sh, CARD)
    sh.line.color.rgb = LINE
    sh.line.width = Pt(1)
    return sh


def _footer(slide, n, total):
    _textbox(slide, Inches(0.5), Inches(7.15), Inches(8), Inches(0.3), "Nexus Cyber · Investor Pitch · Rahasia", size=11, color=MUTED)
    _textbox(slide, Inches(11.2), Inches(7.15), Inches(1.6), Inches(0.3), f"{n}/{total}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def _title_block(slide, eyebrow, title, subtitle=None):
    _textbox(slide, Inches(0.55), Inches(0.28), Inches(12), Inches(0.35), eyebrow.upper(), size=12, bold=True, color=TEAL)
    _textbox(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.7), title, size=32, bold=True, color=WHITE)
    if subtitle:
        _textbox(slide, Inches(0.55), Inches(1.2), Inches(12), Inches(0.45), subtitle, size=15, color=MUTED)


def new_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    _bg(slide)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_plan = []  # filled as we go; footer uses final count after build — we'll two-pass: create then stamp

    # Collect slide builders
    builders = []

    def add(fn):
        builders.append(fn)

    # --- 01 Cover ---
    def s_cover(slide, n, total):
        _textbox(slide, Inches(0.7), Inches(1.8), Inches(11), Inches(0.4), "INVESTOR PITCH DECK · 2026", size=14, bold=True, color=TEAL)
        _textbox(slide, Inches(0.7), Inches(2.3), Inches(12), Inches(1), "Nexus Cyber", size=48, bold=True, color=WHITE)
        _textbox(
            slide,
            Inches(0.7),
            Inches(3.3),
            Inches(11),
            Inches(1.2),
            "Edge Antibody Cowork — GaaS wasit untuk kanal web/API\n+ Channel Starter sebagai pintu masuk UMKM",
            size=20,
            color=SOFT,
        )
        _textbox(
            slide,
            Inches(0.7),
            Inches(5.0),
            Inches(11),
            Inches(0.8),
            "Ukur → kendalikan → uji · residual jujur · bukan SOC otonom 24/7",
            size=14,
            color=MUTED,
        )
        _footer(slide, n, total)

    add(s_cover)

    # --- 02 Agenda ---
    def s_agenda(slide, n, total):
        _title_block(slide, "Outline", "Agenda pitching")
        cols = [
            ("Masalah & pasar", ["Masalah & urgensi", "Root cause + konteks Indonesia", "Persona & dampak"]),
            ("Produk", ["Solusi GaaS + Channel Starter", "Cara kerja & demo", "Arsitektur"]),
            ("Bisnis", ["Kompetitor & USP", "BMC", "Harga per segmen + ROI"]),
            ("Eksekusi", ["Dampak aplikasi", "Roadmap & balik modal", "Risiko & mitigasi per segmen"]),
        ]
        x = 0.5
        for title, items in cols:
            _card(slide, Inches(x), Inches(1.9), Inches(3.0), Inches(4.5))
            _textbox(slide, Inches(x + 0.2), Inches(2.1), Inches(2.6), Inches(0.4), title, size=16, bold=True, color=TEAL)
            _bullets(slide, Inches(x + 0.15), Inches(2.7), Inches(2.7), Inches(3.4), items, size=13)
            x += 3.15
        _footer(slide, n, total)

    add(s_agenda)

    # --- 03 Masalah ---
    def s_masalah(slide, n, total):
        _title_block(
            slide,
            "01 · Masalah",
            "Inovasi digital cepat — siklus risiko HTTP terpecah",
            "Institusi & UMKM menambah portal/API; alat yang ada berhenti di tengah jalan.",
        )
        rows = [
            ("Scanner / audit", "Berhenti di PDF temuan — tidak membuktikan tepi menahan."),
            ("WAF / CDN klasik", "Blok traffic — tidak membuktikan origin masih lemah."),
            ("Chat / Copilot AI", "Saran saja — tidak ada uji replay jujur."),
        ]
        y = 1.9
        for t, d in rows:
            _card(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.15))
            _textbox(slide, Inches(0.8), Inches(y + 0.2), Inches(11.5), Inches(0.35), t, size=18, bold=True, color=AMBER)
            _textbox(slide, Inches(0.8), Inches(y + 0.55), Inches(11.5), Inches(0.45), d, size=15, color=SOFT)
            y += 1.35
        _footer(slide, n, total)

    add(s_masalah)

    # --- 04 Urgensi ---
    def s_urgensi(slide, n, total):
        _title_block(slide, "02 · Urgensi", "Setiap rilis = permukaan serangan baru")
        _bullets(
            slide,
            Inches(0.6),
            Inches(1.9),
            Inches(6),
            Inches(4.5),
            [
                "Fintech, ITSK, agensi web: onboarding & API bertambah tiap sprint.",
                "Deface / abuse form / IDOR pada kanal publik merusak kepercayaan dalam jam.",
                "Regulator & pemilik risiko minta jejak pengendalian — bukan janji verbal.",
                "UMKM punya website murah tanpa pagar → target mudah & massal.",
                "Jendela: dari temuan sampai patch sering berhari-hari tanpa verifikasi ulang.",
            ],
            size=16,
        )
        _card(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(4.3))
        _textbox(slide, Inches(7.25), Inches(2.15), Inches(5.1), Inches(0.4), "Pertanyaan investor", size=16, bold=True, color=TEAL)
        _textbox(
            slide,
            Inches(7.25),
            Inches(2.7),
            Inches(5.1),
            Inches(3.2),
            "Siapa yang membuktikan\nsecara berulang bahwa\ncelah sudah ditutup di tepi\n— bukan hanya “sudah di-scan”?",
            size=20,
            bold=True,
            color=WHITE,
        )
        _footer(slide, n, total)

    add(s_urgensi)

    # --- 05 Root cause ---
    def s_root(slide, n, total):
        _title_block(slide, "02b · Root cause", "Tidak ada loop tertutup ukur → kendalikan → uji")
        causes = [
            ("Pemisahan tool", "Deteksi (scanner) terpisah dari pengendalian (WAF) dan dari verifikasi (uji ulang)."),
            ("Insentif hijau palsu", "Laporan cenderung “bersih” tanpa residual / replay_missed yang jujur."),
            ("Kurang manusia di titik kritis", "Aksi irreversibel tanpa gerbang L0/L1 → atau sebaliknya, semua manual lambat."),
            ("Distribusi salah ukuran", "UMKM dipaksa beli “SOC penuh”; institusi hanya dapat PDF sekali jalan."),
        ]
        y = 1.85
        for t, d in causes:
            _card(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.05))
            _textbox(slide, Inches(0.8), Inches(y + 0.15), Inches(3.2), Inches(0.7), t, size=16, bold=True, color=TEAL)
            _textbox(slide, Inches(4.1), Inches(y + 0.25), Inches(8.3), Inches(0.65), d, size=15, color=SOFT)
            y += 1.15
        _footer(slide, n, total)

    add(s_root)

    # --- 06 Indonesia ---
    def s_id(slide, n, total):
        _title_block(
            slide,
            "02c · Konteks Indonesia",
            "Kenapa celah ini berdampak buruk di Indonesia",
        )
        _bullets(
            slide,
            Inches(0.6),
            Inches(1.85),
            Inches(12),
            Inches(5),
            [
                "Digitalisasi UMKM & sekolah cepat (marketplace, PPDB, portal desa) — keamanan sering “nanti saja”.",
                "Banyak situs di hosting murah / shared / template; deface & form abuse masih kejadian rutin.",
                "Literasi risiko siber pemilik usaha rendah; vendor sering jual “aman” tanpa bukti wasit.",
                "Kanal keuangan & fintech lokal wajib jaga kepercayaan; satu insiden HTTP → reputasi & biaya operasional.",
                "Pengadaan B2G butuh narasi lisensi + retainer jujur (bukan source dump) — celah vendor yang hanya jual PDF.",
                "Biaya enterprise WAF/SOC global terlalu mahal untuk long tail; harga warung tanpa wasit = salah unit ekonomi.",
            ],
            size=16,
        )
        _footer(slide, n, total)

    add(s_id)

    # --- 07 Persona ---
    def s_persona(slide, n, total):
        _title_block(slide, "03 · User persona", "Siapa yang kami layani")
        personas = [
            ("Pemilik UMKM / sekolah", "Butuh website + pagar murah; tidak pegang SOC."),
            ("Founder startup", "Landing/API cepat; butuh Job wasit sebelum pitch/investor."),
            ("Pemilik risiko kanal (fintech)", "Bukti pengendalian + residual; approve L0/L1."),
            ("IT Corporat / integrator", "Hosted Job/Loop atau on-prem jika kritis."),
            ("Pengelola DC instansi", "Edge berlisensi di DC + Loop wajib; source tidak diserahkan."),
            ("Operator Nexus (internal)", "Menjalankan wasit — kokpit GaaS, bukan pelanggan."),
        ]
        positions = [(0.5, 1.85), (4.5, 1.85), (8.5, 1.85), (0.5, 4.3), (4.5, 4.3), (8.5, 4.3)]
        for (x, y), (t, d) in zip(positions, personas):
            _card(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.15))
            _textbox(slide, Inches(x + 0.2), Inches(y + 0.25), Inches(3.4), Inches(0.7), t, size=15, bold=True, color=TEAL)
            _textbox(slide, Inches(x + 0.2), Inches(y + 1.0), Inches(3.4), Inches(0.9), d, size=13, color=SOFT)
        _footer(slide, n, total)

    add(s_persona)

    # --- 08 Dampak masalah ---
    def s_dampak_masalah(slide, n, total):
        _title_block(slide, "04 · Dampak masalah", "Tanpa loop wasit — biaya nyata")
        items = [
            ("Reputasi", "Deface / data bocor di kanal publik → kepercayaan hilang lebih cepat dari recovery teknis."),
            ("Operasional", "Firefight manual, rollback rilis, downtime saat “baru ketahuan” setelah live."),
            ("Keuangan", "Biaya insiden + audit ulang + peluang hilang (UMKM & fintech sama-sama kena)."),
            ("Tata kelola", "Direksi/IT tidak punya artefak residual — keputusan berbasis asumsi."),
        ]
        y = 1.85
        for t, d in items:
            _card(slide, Inches(0.55), Inches(y), Inches(12.2), Inches(1.1))
            _textbox(slide, Inches(0.85), Inches(y + 0.25), Inches(2.5), Inches(0.6), t, size=18, bold=True, color=ROSE)
            _textbox(slide, Inches(3.5), Inches(y + 0.3), Inches(8.9), Inches(0.6), d, size=15, color=SOFT)
            y += 1.2
        _footer(slide, n, total)

    add(s_dampak_masalah)

    # --- 09 Solusi ---
    def s_solusi(slide, n, total):
        _title_block(slide, "05 · Solusi", "Dua lapisan: Channel Starter + Edge Antibody Cowork (GaaS)")
        _card(slide, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.5))
        _textbox(slide, Inches(0.75), Inches(2.1), Inches(5.5), Inches(0.4), "Entry · Channel Starter", size=18, bold=True, color=AMBER)
        _bullets(
            slide,
            Inches(0.7),
            Inches(2.7),
            Inches(5.5),
            Inches(3.3),
            [
                "Form → template → website (UMKM/sekolah).",
                "Harga entry ~Rp 15–35rb/bulan (pilot).",
                "Bukan Loop GaaS penuh di harga warung.",
                "Pintu upsell ke tepi + Job.",
            ],
            size=15,
        )
        _card(slide, Inches(6.8), Inches(1.85), Inches(6.0), Inches(4.5))
        _textbox(slide, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.4), "Inti · GaaS Cowork", size=18, bold=True, color=TEAL)
        _bullets(
            slide,
            Inches(7.0),
            Inches(2.7),
            Inches(5.5),
            Inches(3.3),
            [
                "Alur A: tepi always-on (WAF + antibodi).",
                "Alur B: Job Cowork ukur→kendalikan→uji.",
                "Alur C: artefak MD/JSON + residual jujur.",
                "L0/L1: manusia pemilik risiko tetap di kursi.",
            ],
            size=15,
        )
        _footer(slide, n, total)

    add(s_solusi)

    # --- 10 Cara kerja ---
    def s_cara(slide, n, total):
        _title_block(slide, "06 · Cara kerja", "Siklus wasit yang bisa ditutup jujur")
        steps = [
            ("1. Scope", "Host + izin HTTP jinak"),
            ("2. Ukur", "Defense delta WAF vs origin"),
            ("3. Draft", "Antibodi / pengendalian"),
            ("4. Approve", "L0 artefak / L1 pasang"),
            ("5. Uji", "Vaccine + replay"),
            ("6. Tutup", "CLOSED_OK atau CLOSED_GAP"),
        ]
        x = 0.4
        for t, d in steps:
            _card(slide, Inches(x), Inches(2.2), Inches(2.0), Inches(3.2))
            _textbox(slide, Inches(x + 0.12), Inches(2.45), Inches(1.75), Inches(0.8), t, size=14, bold=True, color=TEAL)
            _textbox(slide, Inches(x + 0.12), Inches(3.4), Inches(1.75), Inches(1.6), d, size=13, color=SOFT)
            x += 2.15
        _textbox(
            slide,
            Inches(0.55),
            Inches(5.7),
            Inches(12),
            Inches(0.6),
            "Aturan produk: replay_missed → CLOSED_GAP (bukan hijau palsu).",
            size=15,
            bold=True,
            color=AMBER,
        )
        _footer(slide, n, total)

    add(s_cara)

    # --- 11 Demo ---
    def s_demo(slide, n, total):
        _title_block(slide, "06b · Demo", "Apa yang ditunjukkan ke investor / juri")
        _bullets(
            slide,
            Inches(0.6),
            Inches(1.9),
            Inches(12),
            Inches(4.8),
            [
                "Site lab (portofolio) di belakang WAF — PROTECTED_HOST, bukan tembak origin langsung.",
                "Portal Channel (`nexus-channel-portal`) — pintu jual per segmen (UMKM → Pemerintah).",
                "Operator GaaS Console — jalankan Job, antrian approve L0/L1, unduh artefak (internal).",
                "Artefak MD/JSON: delta wasit, residual, status antibodi — bukti untuk pemilik risiko.",
                "Pilot distribusi: PC operator 24/7 + Cloudflare Tunnel (SOC tidak dipublikasikan).",
                "Jujur di slide demo: Channel Starter = lab v0.1; mesin Job Cowork sudah ada di lab.",
            ],
            size=16,
        )
        _footer(slide, n, total)

    add(s_demo)

    # --- 12 Arsitektur ---
    def s_arch(slide, n, total):
        _title_block(slide, "07 · Arsitektur", "Diagram alur Nexus Cyber (lab / pilot)")
        # boxes as architecture
        boxes = [
            (0.4, 2.3, "Pengunjung\n/ Internet"),
            (2.7, 2.3, "Caddy /\nTunnel"),
            (5.0, 2.3, "Gateway WAF\n:8080\nAlur A"),
            (7.5, 2.3, "Origin\nPortofolio /\nsite klien"),
            (5.0, 4.6, "Control plane\n:8081\n(internal)"),
            (7.5, 4.6, "Operator GaaS\n:3001"),
            (10.2, 3.4, "NEX-RED\nJob Cowork\n:3004"),
        ]
        for x, y, label in boxes:
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.0), Inches(1.5))
            _fill(sh, CARD)
            sh.line.color.rgb = TEAL_DIM
            tf = sh.text_frame
            tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = label
            _set_run(run, size=12, bold=True, color=WHITE)
        _textbox(
            slide,
            Inches(0.5),
            Inches(6.3),
            Inches(12),
            Inches(0.5),
            "Data plane publik ≠ control plane. SOC/DB/NEX-RED tidak di-tunnel ke juri.",
            size=14,
            color=MUTED,
        )
        _footer(slide, n, total)

    add(s_arch)

    # --- 13 Kompetitor ---
    def s_comp(slide, n, total):
        _title_block(slide, "08 · Kompetitor", "Analisis singkat (jujur)")
        headers = ["Dimensi", "Scanner/PDF", "WAF/CDN saja", "SOC enterprise", "Nexus GaaS"]
        rows = [
            ["Bukti tepi vs origin", "Tidak", "Sebagian", "Ya (mahal)", "Defense delta"],
            ["Uji replay jujur", "Tidak", "Jarang", "Proses berat", "replay_missed → GAP"],
            ["Harga UMKM", "Sekali mahal", "Menengah", "Terlalu mahal", "Entry 15–35rb*"],
            ["Job + artefak risiko", "PDF statis", "Log saja", "Tiket SIEM", "Job Cowork"],
            ["On-prem B2G jujur", "Kadang", "Ya", "Ya", "Edge + Loop wajib"],
        ]
        # header
        x0 = 0.4
        widths = [2.4, 2.3, 2.3, 2.5, 2.6]
        x = x0
        for i, h in enumerate(headers):
            _card(slide, Inches(x), Inches(1.85), Inches(widths[i] - 0.08), Inches(0.55))
            _textbox(slide, Inches(x + 0.08), Inches(1.95), Inches(widths[i] - 0.2), Inches(0.4), h, size=12, bold=True, color=TEAL)
            x += widths[i]
        y = 2.5
        for row in rows:
            x = x0
            for i, cell in enumerate(row):
                _card(slide, Inches(x), Inches(y), Inches(widths[i] - 0.08), Inches(0.7))
                col = TEAL if i == 4 else SOFT
                _textbox(slide, Inches(x + 0.08), Inches(y + 0.18), Inches(widths[i] - 0.2), Inches(0.45), cell, size=11, bold=(i == 0 or i == 4), color=col)
                x += widths[i]
            y += 0.78
        _textbox(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.35), "*Channel Starter lab v0.1 — bukan klaim billing produksi selesai.", size=12, color=MUTED)
        _footer(slide, n, total)

    add(s_comp)

    # --- 14 USP ---
    def s_usp(slide, n, total):
        _title_block(slide, "09 · USP", "Unique Selling Proposition")
        usps = [
            ("Wasit tertutup", "Satu Job mengikat ukur → patch → replay → status jujur."),
            ("Anti hijau palsu", "replay_missed tidak boleh CLOSED_OK."),
            ("Dua lapisan harga", "Entry UMKM terpisah dari Loop GaaS / on-prem."),
            ("Otonomi terbatas", "L0/L1 — manusia pemilik risiko tetap mengesahkan."),
            ("Moat teknis", "Jalur di depan origin + jejak risiko — bukan UI SOC."),
            ("B2G jujur", "Lisensi Edge + Loop; source & SOC tidak diserahkan."),
        ]
        positions = [(0.5, 1.85), (4.5, 1.85), (8.5, 1.85), (0.5, 4.35), (4.5, 4.35), (8.5, 4.35)]
        for (x, y), (t, d) in zip(positions, usps):
            _card(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.2))
            _textbox(slide, Inches(x + 0.2), Inches(y + 0.3), Inches(3.4), Inches(0.5), t, size=16, bold=True, color=TEAL)
            _textbox(slide, Inches(x + 0.2), Inches(y + 0.95), Inches(3.4), Inches(1.0), d, size=14, color=SOFT)
        _footer(slide, n, total)

    add(s_usp)

    # --- 15 BMC ---
    def s_bmc(slide, n, total):
        _title_block(slide, "10 · BMC", "Business Model Canvas (ringkas)")
        cells = [
            (0.35, 1.75, 2.4, 2.4, "Key Partners", "Tunnel/CDN· Hosting opsional· Integrator web· WA commerce"),
            (2.85, 1.75, 2.4, 2.4, "Key Activities", "Job Cowork· Edge ops· Artefak· Portal jual· Upsell"),
            (5.35, 1.75, 2.4, 1.15, "Value Prop", "Wasit HTTP + entry site murah"),
            (5.35, 3.0, 2.4, 1.15, "Key Resources", "Gateway· NEX-RED· IP proses· PC pilot"),
            (7.85, 1.75, 2.4, 2.4, "Customer Rel.", "WA manual· Artefak· Retainer Loop· L0/L1"),
            (10.35, 1.75, 2.5, 2.4, "Channels", "Portal segmen· WA· Demo lab· Pitch B2B/B2G"),
            (0.35, 4.3, 4.9, 1.9, "Cost Structure", "Waktu operator (Rp100rb/jam asumsi)· Listrik PC· R&D· Belum VPS di pilot"),
            (5.35, 4.3, 3.7, 1.9, "Customer Segments", "UMKM· Sekolah· Startup· Corporat· Pemerintah"),
            (9.15, 4.3, 3.7, 1.9, "Revenue Streams", "Starter bln· Job 200rb· Loop 300rb· Edge 18jt/thn· Loop on-prem 3,5jt/bln"),
        ]
        for x, y, w, h, title, body in cells:
            _card(slide, Inches(x), Inches(y), Inches(w), Inches(h))
            _textbox(slide, Inches(x + 0.1), Inches(y + 0.08), Inches(w - 0.2), Inches(0.35), title, size=11, bold=True, color=TEAL)
            _textbox(slide, Inches(x + 0.1), Inches(y + 0.45), Inches(w - 0.2), Inches(h - 0.55), body, size=11, color=SOFT)
        _footer(slide, n, total)

    add(s_bmc)

    # --- 16-20 Finance per segment ---
    finance = [
        (
            "11a · Finance · UMKM",
            "Unit ekonomi UMKM (pilot PC+tunnel)",
            [
                "Belum punya web · Website Aman: Rp 20rb/bln · margin ~58%",
                "Belum · GaaS entry: Rp 35rb/bln · margin ~47%",
                "Sudah punya web · Pagar: Rp 15rb/bln · margin ~60%",
                "Sudah · Pagar+status: Rp 28rb/bln · margin ~48%",
                "COGS utama = waktu operator batch — bukan VPS.",
                "Jangan jual Loop institusi di harga warung.",
            ],
        ),
        (
            "11b · Finance · Sekolah",
            "Struktur harga selaras UMKM (pilot)",
            [
                "Website / GaaS entry / Pagar / Pagar+alert — angka sama band UMKM.",
                "Nilai: anti-deface & ketersediaan info publik sekolah.",
                "Support harus template/batch agar margin bertahan.",
                "Upsell ke Job Wasit jika ada event PPDB / portal kritis.",
            ],
        ),
        (
            "11c · Finance · Startup",
            "Landing + tepi + wasit",
            [
                "Landing+pagar Rp 45rb/bln · margin ~64%",
                "Landing+Tepi / Tepi Rp 75rb/bln · margin ~63–71%",
                "Job Wasit Rp 200rb sekali · margin ~63%",
                "Loop GaaS Rp 300rb/bln · margin ~68%",
                "Cocok sebelum fundraising / launch fitur sensitif.",
            ],
        ),
        (
            "11d · Finance · Corporat",
            "Hosted vs on-prem",
            [
                "Hosted: Job Rp 200rb · Loop Rp 300rb/bln (1 host pilot).",
                "On-prem (besar/kritis): model sama Pemerintah — Edge + Loop.",
                "Custom multi-host / SLA tertulis = quote WA.",
                "Kontrak pisah: build site vs keamanan wasit.",
            ],
        ),
        (
            "11e · Finance · Pemerintah",
            "On-prem Edge + Loop wajib",
            [
                "Lisensi Edge On-Prem: Rp 18jt / tahun (1 zona/DC ilustrasi).",
                "Loop On-Prem wajib: Rp 3,5jt / bulan.",
                "Tahun 1 tipikal 1 DC: jual ~Rp 60jt · COGS ~Rp 15jt · laba ~Rp 45jt · margin ~75%.",
                "Source & control plane SOC TIDAK termasuk.",
                "Pitching ≠ sudah live pengadaan SIPLah/E-Katalog.",
            ],
        ),
    ]

    for eyebrow, title, bullets in finance:
        def make(ey=eyebrow, ti=title, bu=bullets):
            def _s(slide, n, total):
                _title_block(slide, ey, ti)
                _bullets(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(4.8), bu, size=17)
                _footer(slide, n, total)

            return _s

        add(make())

    # --- 21 ROI ---
    def s_roi(slide, n, total):
        _title_block(slide, "12 · ROI", "Return untuk pelanggan & untuk Nexus")
        _card(slide, Inches(0.5), Inches(1.85), Inches(6.0), Inches(4.5))
        _textbox(slide, Inches(0.75), Inches(2.1), Inches(5.5), Inches(0.4), "ROI pelanggan (ilustrasi)", size=18, bold=True, color=TEAL)
        _bullets(
            slide,
            Inches(0.7),
            Inches(2.7),
            Inches(5.5),
            Inches(3.3),
            [
                "1 insiden deface/abuse sering > biaya Loop sebulan.",
                "Artefak mempercepat keputusan L0/L1 (kurangi meeting kosong).",
                "UMKM: site+pagar lebih murah dari jasa “bikin web + amankan” terpisah.",
                "Corporat/B2G: bukti residual untuk audit internal.",
            ],
            size=14,
        )
        _card(slide, Inches(6.8), Inches(1.85), Inches(6.0), Inches(4.5))
        _textbox(slide, Inches(7.05), Inches(2.1), Inches(5.5), Inches(0.4), "ROI Nexus (pilot)", size=18, bold=True, color=AMBER)
        _bullets(
            slide,
            Inches(7.0),
            Inches(2.7),
            Inches(5.5),
            Inches(3.3),
            [
                "Infra cash ≈ 0 (PC + tunnel) → margin % sehat.",
                "UMKM = volume; B2G = laba absolut.",
                "1 klien B2G tahun-1 ≈ puluhan juta laba kotor ilustrasi.",
                "Upsell alami: Starter → Tepi → Job → Loop → On-prem.",
            ],
            size=14,
        )
        _footer(slide, n, total)

    add(s_roi)

    # --- 22 Dampak aplikasi ---
    def s_dampak_app(slide, n, total):
        _title_block(slide, "13 · Dampak aplikasi", "Jika Nexus dipakai luas")
        impacts = [
            ("Teknis", "Lebih banyak kanal punya bukti pengendalian berulang, bukan PDF sekali."),
            ("Ekonomi", "UMKM dapat pagar terjangkau; institusi membayar wasit sesuai unit ekonomi."),
            ("Tata kelola", "Jejak approve L0/L1 + residual tertulis untuk pemilik risiko."),
            ("Industri lokal", "Alternatif jujur vs “SOC 24/7” yang tidak realistis di harga warung."),
            ("Negara / B2G", "Model lisensi runtime + Loop — IP tetap di vendor, operasi di DC klien."),
            ("Batasan jujur", "Bukan pengganti GRC bank penuh, pentest exploit, atau SLA data center."),
        ]
        positions = [(0.5, 1.85), (4.5, 1.85), (8.5, 1.85), (0.5, 4.35), (4.5, 4.35), (8.5, 4.35)]
        for (x, y), (t, d) in zip(positions, impacts):
            _card(slide, Inches(x), Inches(y), Inches(3.8), Inches(2.2))
            _textbox(slide, Inches(x + 0.2), Inches(y + 0.25), Inches(3.4), Inches(0.45), t, size=16, bold=True, color=TEAL)
            _textbox(slide, Inches(x + 0.2), Inches(y + 0.85), Inches(3.4), Inches(1.1), d, size=13, color=SOFT)
        _footer(slide, n, total)

    add(s_dampak_app)

    # --- 23 Roadmap ---
    def s_roadmap(slide, n, total):
        _title_block(slide, "14 · Roadmap & modal", "Pengembangan lanjut + balik modal (ilustrasi)")
        _bullets(
            slide,
            Inches(0.55),
            Inches(1.75),
            Inches(7.2),
            Inches(4.8),
            [
                "Q sekarang: demo Job Cowork + portal segmen + PC/tunnel pilot.",
                "+3–6 bln: packaging Edge berlisensi; 3–5 klien Loop hosted; 1 pilot on-prem staging.",
                "+6–12 bln: billing lebih rapi; multi-host Corporat; B2G kontrak pertama (jika ada).",
                "Biaya R&D+ops founder (ilustrasi): Rp 8–15jt/bulan blended (waktu + alat) — sesuaikan aktual Anda.",
                "Target balik modal kas: 1 klien B2G tahun-1 atau ~40–80 Loop UMKM/startup sehat.",
                "Skenario konservatif tanpa B2G: break-even ~6–12 bulan jika closing Loop stabil.",
                "Angka di slide = model internal PRICING_UNIT_ECONOMICS — bukan prospektus resmi.",
            ],
            size=14,
        )
        _card(slide, Inches(8.0), Inches(1.85), Inches(4.7), Inches(4.4))
        _textbox(slide, Inches(8.25), Inches(2.15), Inches(4.2), Inches(0.4), "Ringkas payback", size=16, bold=True, color=AMBER)
        _textbox(
            slide,
            Inches(8.25),
            Inches(2.7),
            Inches(4.2),
            Inches(3.2),
            "Modal kerja kecil\n(infra ≈ 0 di pilot)\n\nPayback cepat jika\nada 1× B2G\natau funnel Loop\nstartup/corporat\n\nUMKM = bahan bakar\nvolume & brand",
            size=15,
            color=SOFT,
        )
        _footer(slide, n, total)

    add(s_roadmap)

    # --- 24-28 Risk per segment ---
    risks = [
        (
            "15a · Risiko · UMKM",
            [
                "Support melebihi batch → margin habis. Mitigasi: template alert, naikkan harga, batasi volume.",
                "Ekspektasi “SOC 24/7” di Rp 20rb. Mitigasi: copy jujur di portal + onboarding WA.",
                "Churn tinggi. Mitigasi: upsell pagar bermakna; bundling site yang susah pindah.",
                "Abuse trial. Mitigasi: WA manual, bukan self-serve kartu kredit dulu.",
            ],
        ),
        (
            "15b · Risiko · Sekolah",
            [
                "Anggaran musiman / tahun ajaran. Mitigasi: paket tahunan & jadwal Job sebelum PPDB.",
                "PIC non-teknis. Mitigasi: artefak bahasa sederhana + WA.",
                "Hosting pihak ke-3 berubah. Mitigasi: dokumentasikan PROTECTED_HOST & DNS.",
                "Klaim berlebihan ke orang tua/komite. Mitigasi: batasi klaim di kontrak.",
            ],
        ),
        (
            "15c · Risiko · Startup",
            [
                "Pivot produk → host berubah. Mitigasi: Loop per host; quote ulang scope.",
                "Minta pentest exploit. Mitigasi: tegas wasit HTTP jinak ≠ Shannon.",
                "Bayar telat. Mitigasi: Job one-shot dulu; Loop setelah artefak pertama.",
                "Overload fitur. Mitigasi: ICP ketat — kanal web/API saja.",
            ],
        ),
        (
            "15d · Risiko · Corporat",
            [
                "Minta SOC penuh ke tangan mereka. Mitigasi: portal status ringan ≠ control plane.",
                "SLA data center di fase pilot PC. Mitigasi: kontrak jujur pilot; roadmap VPS/DC.",
                "Multi-host tanpa harga. Mitigasi: harga per host + custom quote.",
                "Vendor lock-in fear. Mitigasi: artefak portabel MD/JSON; on-prem path jelas.",
            ],
        ),
        (
            "15e · Risiko · Pemerintah",
            [
                "Pengadaan lama / HPS. Mitigasi: pitching dulu; jangan klaim E-Katalog selesai.",
                "Permintaan source code. Mitigasi: kebijakan tetap — lisensi runtime + escrow jika wajib.",
                "Air-gap / DC khusus. Mitigasi: custom quote; Loop tetap wajib.",
                "Gagal packaging lisensi. Mitigasi: milestone packaging sebelum janji go-live.",
            ],
        ),
    ]

    for eyebrow, bullets in risks:
        def make(ey=eyebrow, bu=bullets):
            def _s(slide, n, total):
                _title_block(slide, ey, "Mitigasi & risiko")
                _bullets(slide, Inches(0.6), Inches(1.9), Inches(12), Inches(4.8), bu, size=16)
                _footer(slide, n, total)

            return _s

        add(make())

    # --- Closing ---
    def s_close(slide, n, total):
        _textbox(slide, Inches(0.7), Inches(2.0), Inches(12), Inches(0.4), "PENUTUP", size=14, bold=True, color=TEAL)
        _textbox(slide, Inches(0.7), Inches(2.5), Inches(12), Inches(1), "The Ask", size=40, bold=True, color=WHITE)
        _bullets(
            slide,
            Inches(0.7),
            Inches(3.6),
            Inches(11),
            Inches(2.5),
            [
                "Dukungan untuk percepat packaging Edge berlisensi + 3–5 klien Loop berbayar.",
                "Intro ke ICP: fintech/integrator & pintu B2G staging (bukan klaim produksi selesai).",
                "Diskusi putaran modal kerja sesuai angka aktual founder (slide roadmap = kerangka).",
            ],
            size=17,
        )
        _textbox(slide, Inches(0.7), Inches(6.2), Inches(12), Inches(0.4), "Nexus Cyber — wasit yang bisa ditutup jujur.", size=14, color=MUTED)
        _footer(slide, n, total)

    add(s_close)

    # Disclaimer
    def s_disc(slide, n, total):
        _title_block(slide, "Lampiran", "Disclaimer pitching")
        _bullets(
            slide,
            Inches(0.6),
            Inches(1.9),
            Inches(12),
            Inches(4.8),
            [
                "Angka harga & margin dari docs/PRICING_UNIT_ECONOMICS.md — ilustrasi pilot, bukan laporan audit.",
                "Tidak mengklaim: eBPF/XDP nyata, SOC otonom 24/7, pentest exploit setara Shannon, billing Channel Starter produksi selesai.",
                "Lab target: portofolio di belakang WAF — bukti mesin, bukan produk yang dijual sebagai “portofolio SaaS”.",
                "SOC / Operator Console hanya internal Nexus — bukan deliverable pelanggan.",
                "Sumber: docs/PRODUCT_MODEL.md, COWORK_B2B.md, COWORK_B2G.md, DISTRIBUTION_PILOT.md, AGENTS.md.",
            ],
            size=15,
        )
        _footer(slide, n, total)

    add(s_disc)

    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        slide = new_slide(prs)
        fn(slide, i, total)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({total} slides)")


if __name__ == "__main__":
    build()
