"""
Script untuk Menggabungkan Desain Visual Pengguna dengan Konten Bisnis, Investor-Ready, & Market-Driven
- Format: 16:9 Widescreen (Light Mode Enterprise: Navy Blue & Emerald Green)
- Fitur Utama yang Ditambahkan:
  1. Slide 5: Tabel Komparasi Contreng (✔) & Silang (✖) + Competitive Advantages
  2. Slide 6: Market Size (TAM, SAM, SOM) & Market Positioning Matrix
  3. Slide 7: Cara Kerja Solusi Sederhana (3-Step Business Flow)
  4. Slide 8: Go-To-Market (GTM) Strategy (3 Pilar Distribusi)
  5. Slide 9: Dampak Sosial & Inklusi Finansial (Social Impact & ESG)
  6. Slide 10: Validasi Pasar, Traksi, & Minat Calon Adopter
  7. Slide 11: Model Bisnis & Proyeksi Finansial 3 Tahun
  8. Slide 13: Roadmap Eksekusi & Analisis ROI (ROI Investor 18–24 Bulan, Maks 2.5 Thn, ROI Klien < 3 Bulan)
  9. Slide 15-22: 8 Slide Appendix Q&A (Technical Deep-Dive, Dual-Brain AI, Kepatuhan POJK 30/2025, Detail COGS, Demo SOP)
"""

import os
import shutil
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# PALET WARNA LIGHT MODE ENTERPRISE
# ==========================================
BG_COLOR = RGBColor(248, 250, 252)       # Slate 50 (Soft Off-White)
CARD_BG = RGBColor(255, 255, 255)        # Pure White
CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200
NAVY_PRIMARY = RGBColor(15, 23, 42)      # Slate 900
NAVY_SECONDARY = RGBColor(30, 58, 138)   # Blue 900
EMERALD_ACCENT = RGBColor(5, 150, 105)   # Emerald 600
EMERALD_LIGHT = RGBColor(209, 250, 229)  # Emerald 100
CYAN_ACCENT = RGBColor(2, 132, 199)      # Sky 600
AMBER_ACCENT = RGBColor(217, 119, 6)     # Amber 600
ROSE_ACCENT = RGBColor(225, 29, 72)      # Rose 600
TEXT_PRIMARY = RGBColor(15, 23, 42)      # Slate 900
TEXT_SECONDARY = RGBColor(71, 85, 105)   # Slate 600
TEXT_MUTED = RGBColor(148, 163, 184)     # Slate 400
TABLE_HEADER_BG = RGBColor(30, 41, 59)   # Slate 800
TABLE_ALT_ROW = RGBColor(241, 245, 249)  # Slate 100

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def set_slide_background(slide, width, height):
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), width, height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = BG_COLOR
    bg_shape.line.fill.background()
    return bg_shape

def add_header(slide, title, category_badge="CORE", subtitle="", slide_num=""):
    # Badge
    badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4.5), Inches(0.35))
    tf_b = badge_box.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
    p_b = tf_b.paragraphs[0]
    is_core = "CORE" in category_badge
    p_b.text = f"● {category_badge}"
    p_b.font.size = Pt(10)
    p_b.font.bold = True
    p_b.font.name = FONT_HEADING
    p_b.font.color.rgb = EMERALD_ACCENT if is_core else CYAN_ACCENT

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(10.5), Inches(0.6))
    tf_t = title_box.text_frame
    tf_t.word_wrap = True
    tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(20)
    p_t.font.bold = True
    p_t.font.name = FONT_HEADING
    p_t.font.color.rgb = NAVY_PRIMARY

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(10.5), Inches(0.35))
        tf_s = sub_box.text_frame
        tf_s.word_wrap = True
        tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
        p_s = tf_s.paragraphs[0]
        p_s.text = subtitle
        p_s.font.size = Pt(11)
        p_s.font.name = FONT_BODY
        p_s.font.color.rgb = TEXT_SECONDARY

    # Slide Number
    if slide_num:
        num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.0), Inches(0.35))
        tf_n = num_box.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.alignment = PP_ALIGN.RIGHT
        p_n.text = f"Slide {slide_num}"
        p_n.font.size = Pt(10)
        p_n.font.bold = True
        p_n.font.color.rgb = TEXT_MUTED

def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = border_color
    card.line.width = Pt(1)
    return card

# ==========================================
# SLIDE 1: COVER & TIM PROFILING (PIDI MANDATORY)
# ==========================================
def populate_slide_1_cover(s1, w, h):
    set_slide_background(s1, w, h)

    # Accent top border
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), w, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = EMERALD_ACCENT
    top_bar.line.fill.background()

    # Competition Header
    ch_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_ch = ch_box.text_frame
    p_ch = tf_ch.paragraphs[0]
    p_ch.text = "PIDI DIGDAYA x HACKATHON 2026 · BANK INDONESIA, OJK, AFTECH, LPPI"
    p_ch.font.size = Pt(11)
    p_ch.font.bold = True
    p_ch.font.color.rgb = CYAN_ACCENT

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(1.3))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "NEXUS CYBER"
    p_t.font.size = Pt(36)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY_PRIMARY

    p_sub = tf_t.add_paragraph()
    p_sub.text = "Autonomous Edge Antibody Guardrail-as-a-Service (GaaS) untuk Ketahanan Siber Finansial"
    p_sub.font.size = Pt(13.5)
    p_sub.font.color.rgb = NAVY_SECONDARY
    p_sub.space_before = Pt(2)

    # Left Card: Problem Statement & Regulatory Mandate
    add_card(s1, Inches(0.8), Inches(2.25), Inches(5.7), Inches(3.6))
    tb_ps = s1.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.3), Inches(3.3))
    tf_ps = tb_ps.text_frame
    tf_ps.word_wrap = True

    p1 = tf_ps.paragraphs[0]
    p1.text = "PROBLEM STATEMENT RESMI (PIDI 2026)"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD_ACCENT

    p2 = tf_ps.add_paragraph()
    p2.text = "PS-1: Penguatan Ketahanan dan Inovasi Keuangan"
    p2.font.size = Pt(13)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_PRIMARY
    p2.space_before = Pt(4)

    p3 = tf_ps.add_paragraph()
    p3.text = "Sub-Topik: Manajemen Risiko Siber & Operasional Kanal Digital"
    p3.font.size = Pt(10)
    p3.font.color.rgb = TEXT_SECONDARY
    p3.space_before = Pt(2)

    p4 = tf_ps.add_paragraph()
    p4.text = "Kepatuhan Mandat Regulasi:"
    p4.font.size = Pt(10.5)
    p4.font.bold = True
    p4.font.color.rgb = NAVY_PRIMARY
    p4.space_before = Pt(8)

    p5 = tf_ps.add_paragraph()
    p5.text = "• POJK No. 30/2025 (Manajemen Risiko Siber Lembaga Keuangan / ITSK)\n• UU PDP No. 27/2022 (Perlindungan Data Pribadi Nasabah)\n• Standar ISO 27001 (Kontrol Audit Trail & Logging A.12.4)"
    p5.font.size = Pt(9.5)
    p5.font.color.rgb = TEXT_SECONDARY
    p5.space_before = Pt(2)

    # Right Card: Team Profiling & Affiliation
    add_card(s1, Inches(6.8), Inches(2.25), Inches(5.73), Inches(3.6))
    tb_tm = s1.shapes.add_textbox(Inches(7.0), Inches(2.4), Inches(5.33), Inches(3.3))
    tf_tm = tb_tm.text_frame
    tf_tm.word_wrap = True

    p6 = tf_tm.paragraphs[0]
    p6.text = "PROFILING TIM PESERTA & INSTITUSI"
    p6.font.size = Pt(11)
    p6.font.bold = True
    p6.font.color.rgb = CYAN_ACCENT

    p7 = tf_tm.add_paragraph()
    p7.text = "Nama Tim: TELULANG  |  ID: PIDIXHACKATHON_S0280"
    p7.font.size = Pt(12)
    p7.font.bold = True
    p7.font.color.rgb = NAVY_PRIMARY
    p7.space_before = Pt(4)

    p8 = tf_tm.add_paragraph()
    p8.text = "Asal Institusi: Telkom University"
    p8.font.size = Pt(10.5)
    p8.font.bold = True
    p8.font.color.rgb = EMERALD_ACCENT
    p8.space_before = Pt(2)

    p9 = tf_tm.add_paragraph()
    p9.text = "Susunan Anggota Tim Pengembang:"
    p9.font.size = Pt(10.5)
    p9.font.bold = True
    p9.font.color.rgb = NAVY_PRIMARY
    p9.space_before = Pt(6)

    tm_members = [
        ("1. Thoriq (Ketua Tim)", "Lead Cyber & Core Engine Architect"),
        ("2. AI Reasoning Engineer", "NEX-RED Cowork & Dual-Brain Model Lead"),
        ("3. Fullstack Specialist", "SOC Command Center & POJK Compliance Lead")
    ]
    for name, role in tm_members:
        p_m = tf_tm.add_paragraph()
        p_m.text = f"• {name} — {role}"
        p_m.font.size = Pt(9.5)
        p_m.font.color.rgb = TEXT_SECONDARY
        p_m.space_before = Pt(1)

    # Bottom Value Proposition Pill
    add_card(s1, Inches(0.8), Inches(6.0), Inches(11.73), Inches(1.0), bg_color=RGBColor(240, 253, 244), border_color=EMERALD_ACCENT)
    tb_vp = s1.shapes.add_textbox(Inches(1.0), Inches(6.05), Inches(11.33), Inches(0.9))
    tf_vp = tb_vp.text_frame
    tf_vp.word_wrap = True
    p_vp = tf_vp.paragraphs[0]
    p_vp.text = "VALUE PROPOSITION 1 KALIMAT:"
    p_vp.font.size = Pt(9.5)
    p_vp.font.bold = True
    p_vp.font.color.rgb = EMERALD_ACCENT

    p_vpt = tf_vp.add_paragraph()
    p_vpt.text = "\"Platform pengaman otonom (GaaS) yang mendeteksi celah kanal web/API perbankan, memasang antibodi di tepi dalam <1.2ms, dan memverifikasi penutupan celah secara nyata tanpa perlu merombak kode backend.\""
    p_vpt.font.size = Pt(10.5)
    p_vpt.font.italic = True
    p_vpt.font.color.rgb = NAVY_PRIMARY
    p_vpt.space_before = Pt(2)

# ==========================================
# SLIDE 5: COMPETITIVE ADVANTAGE & CHECKLIST MATRIX (✔ / ✖)
# ==========================================
def populate_slide_5(s5, w, h):
    set_slide_background(s5, w, h)
    add_header(s5, "Keunggulan Kompetitif & Matriks Perbandingan Solusi", "CORE 5/14", 
               "Perbandingan kemampuan objektif: Mengapa Nexus Cyber jauh lebih unggul bagi perbankan & fintech.", "05")

    # Table with Balanced 3-Tier Status (✔, ~, ✖)
    rows = 8
    cols = 5
    table_shape = s5.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(3.6))
    table = table_shape.table

    table.columns[0].width = Inches(3.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.97)
    table.columns[3].width = Inches(1.98)
    table.columns[4].width = Inches(1.98)

    headers = ["Parameter & Kapabilitas Kunci", "NEXUS CYBER", "Cloudflare / AWS", "Tenable Nessus", "Microsoft Copilot"]
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.text = head
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255, 255, 255)

    data_rows = [
        ("Mitigasi Celah Tepi Instan (<2 Menit)", "✔", "~", "✖", "~"),
        ("Wasit Uji Ulang Tembakan (Replay Proof)", "✔", "✖", "✖", "✖"),
        ("Proteksi DDoS Skala Global (Terabits)", "~", "✔", "✖", "✖"),
        ("Pemindaian Mendalam Ribuan Database CVE", "~", "✖", "✔", "~"),
        ("Analisis & Ringkasan Insiden AI (GenAI)", "✔", "~", "✖", "✔"),
        ("Keterjangkauan Biaya BPR & Fintech Daerah", "✔", "✖", "~", "~"),
        ("Kepatuhan Terintegrasi POJK No. 30/2025", "✔", "~", "~", "✖")
    ]

    for i, row in enumerate(data_rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            if j == 1:
                cell.fill.fore_color.rgb = EMERALD_LIGHT
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            if j == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(9.5)
                p.font.bold = True
                p.font.color.rgb = NAVY_PRIMARY
            else:
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(15)
                p.font.bold = True
                if val == "✔":
                    p.font.color.rgb = EMERALD_ACCENT
                elif val == "~":
                    p.font.color.rgb = AMBER_ACCENT
                else:
                    p.font.color.rgb = ROSE_ACCENT

    # Bottom 3 Direct Competitor Knockout Cards
    moats = [
        ("1. VS CLOUDFLARE / AWS WAF", 
         "• Kelemahan: Biaya ratusan juta/thn & butuh tim engineer setting rule.\n✔ Nexus: 90% Lebih Terjangkau bagi BPR/UMKM & aktif dalam 5 menit."),
        ("2. VS TENABLE / PENTEST", 
         "• Kelemahan: Hanya beri PDF tebal, perbaikan kode butuh >14 hari.\n✔ Nexus: Mitigasi Celah Tepi Instan (< 2 Menit) + Wasit Replay Proof."),
        ("3. VS MICROSOFT COPILOT", 
         "• Kelemahan: Hanya chatbot saran teks, tak bisa tahan trafik live.\n✔ Nexus: AI Otonom di Jalur Trafik (Inline Shield) secara real-time.")
    ]
    for i, (m_title, m_desc) in enumerate(moats):
        x = Inches(0.8 + i * 4.0)
        add_card(s5, x, Inches(5.6), Inches(3.7), Inches(1.5), bg_color=CARD_BG, border_color=EMERALD_ACCENT)
        tb = s5.shapes.add_textbox(x + Inches(0.15), Inches(5.68), Inches(3.4), Inches(1.35))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = m_title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = NAVY_PRIMARY
        
        for line in m_desc.split("\n"):
            p_d = tf.add_paragraph()
            p_d.text = line
            p_d.font.size = Pt(8.5)
            if line.startswith("✔"):
                p_d.font.bold = True
                p_d.font.color.rgb = EMERALD_ACCENT
            else:
                p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(1)

# ==========================================
# SLIDE 5B: 2x2 MARKET POSITIONING MATRIX (COMPETITOR & POSITIONING)
# ==========================================
def populate_slide_positioning(sp, w, h):
    set_slide_background(sp, w, h)
    
    # Top Left Badge: Competitor and Positioning
    badge_box = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.45))
    badge_box.fill.solid()
    badge_box.fill.fore_color.rgb = RGBColor(239, 246, 255) # Blue 50
    badge_box.line.color.rgb = RGBColor(147, 197, 253)     # Blue 300
    badge_box.line.width = Pt(1.2)
    tf_bb = badge_box.text_frame
    p_bb = tf_bb.paragraphs[0]
    p_bb.text = "Competitor and Positioning"
    p_bb.font.size = Pt(11)
    p_bb.font.bold = True
    p_bb.font.color.rgb = NAVY_SECONDARY
    p_bb.alignment = PP_ALIGN.CENTER

    # Top Right Header Tag
    tr_box = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(0.4), Inches(3.53), Inches(0.45))
    tr_box.fill.solid()
    tr_box.fill.fore_color.rgb = CARD_BG
    tr_box.line.color.rgb = CARD_BORDER
    tr_box.line.width = Pt(1)
    tf_tr = tr_box.text_frame
    p_tr = tf_tr.paragraphs[0]
    p_tr.text = "NEXUS CYBER  |  PIDI Hackathon 2026"
    p_tr.font.size = Pt(10.5)
    p_tr.font.bold = True
    p_tr.font.color.rgb = EMERALD_ACCENT
    p_tr.alignment = PP_ALIGN.CENTER

    # Center Coordinates for 2x2 Axes (Intersection at x=6.66 in, y=3.95 in)
    center_x = Inches(6.66)
    center_y = Inches(3.95)

    # Horizontal Axis Line (X-Axis)
    line_h = sp.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), center_y - Inches(0.015), Inches(11.33), Inches(0.03))
    line_h.fill.solid()
    line_h.fill.fore_color.rgb = RGBColor(37, 99, 235) # Blue 600
    line_h.line.fill.background()

    # Vertical Axis Line (Y-Axis)
    line_v = sp.shapes.add_shape(MSO_SHAPE.RECTANGLE, center_x - Inches(0.015), Inches(1.15), Inches(0.03), Inches(5.6))
    line_v.fill.solid()
    line_v.fill.fore_color.rgb = RGBColor(37, 99, 235) # Blue 600
    line_v.line.fill.background()

    # Axis Labels Pills
    # Top Y Label: Autonomous Closed-Loop
    lbl_top = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x - Inches(2.2), Inches(1.0), Inches(4.4), Inches(0.4))
    lbl_top.fill.solid()
    lbl_top.fill.fore_color.rgb = RGBColor(239, 246, 255)
    lbl_top.line.color.rgb = RGBColor(147, 197, 253)
    p = lbl_top.text_frame.paragraphs[0]
    p.text = "Autonomous Closed-Loop (Proteksi & Wasit Otonom)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Bottom Y Label: Passive Reporting
    lbl_bot = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, center_x - Inches(2.1), Inches(6.8), Inches(4.2), Inches(0.4))
    lbl_bot.fill.solid()
    lbl_bot.fill.fore_color.rgb = RGBColor(239, 246, 255)
    lbl_bot.line.color.rgb = RGBColor(147, 197, 253)
    p = lbl_bot.text_frame.paragraphs[0]
    p.text = "Passive Reporting (Laporan Statis & Advisory)"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Left X Label: High Cost & Complex Setup
    lbl_l = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), center_y - Inches(0.2), Inches(2.3), Inches(0.4))
    lbl_l.fill.solid()
    lbl_l.fill.fore_color.rgb = RGBColor(239, 246, 255)
    lbl_l.line.color.rgb = RGBColor(147, 197, 253)
    p = lbl_l.text_frame.paragraphs[0]
    p.text = "High Cost & Complex Setup"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # Right X Label: Cost-Efficient Plug & Play
    lbl_r = sp.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.23), center_y - Inches(0.2), Inches(2.3), Inches(0.4))
    lbl_r.fill.solid()
    lbl_r.fill.fore_color.rgb = RGBColor(239, 246, 255)
    lbl_r.line.color.rgb = RGBColor(147, 197, 253)
    p = lbl_r.text_frame.paragraphs[0]
    p.text = "Cost-Efficient Plug & Play"
    p.font.size = Pt(9.5)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    # ================= 4 QUADRANT CARDS =================

    # [1] TOP-RIGHT: WINNING POSITION (NEXUS CYBER)
    add_card(sp, Inches(6.9), Inches(1.55), Inches(5.43), Inches(2.2), bg_color=RGBColor(240, 253, 244), border_color=EMERALD_ACCENT)
    tb_nexus = sp.shapes.add_textbox(Inches(7.1), Inches(1.65), Inches(5.1), Inches(2.0))
    tf_nx = tb_nexus.text_frame
    tf_nx.word_wrap = True
    p = tf_nx.paragraphs[0]
    p.text = "★ NEXUS CYBER (Guardrail-as-a-Service)"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    nx_bullets = [
        "✔ Proteksi Tepi Instan (<1.2ms) + Wasit Uji Ulang (Replay Proof)",
        "✔ Sangat Terjangkau: UMKM (Rp 20rb - 35rb) & BPR/Fintech (Rp 300rb)",
        "✔ Zero-Code Onboarding: Selesai Terpasang 5 Menit Tanpa Ubah Backend"
    ]
    for b in nx_bullets:
        p_b = tf_nx.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(9.5)
        p_b.font.bold = True
        p_b.font.color.rgb = NAVY_PRIMARY
        p_b.space_before = Pt(3)

    # [2] TOP-LEFT: ENTERPRISE CLOUD WAF
    add_card(sp, Inches(1.0), Inches(1.55), Inches(5.4), Inches(2.2))
    tb_tl = sp.shapes.add_textbox(Inches(1.2), Inches(1.65), Inches(5.0), Inches(2.0))
    tf_tl = tb_tl.text_frame
    tf_tl.word_wrap = True
    p = tf_tl.paragraphs[0]
    p.text = "Cloudflare Enterprise / AWS WAF / Akamai"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(234, 88, 12) # Orange 600

    tl_bullets = [
        "• Proteksi DDoS Terabit Global & Jaringan Anycast CDN Luas",
        "✖ Biaya Sangat Mahal (Ratusan Juta/Thn) bagi BPR & Fintech Daerah",
        "✖ Rule Manual Kaku; Tidak Ada Wasit Uji Ulang Origin Otomatis"
    ]
    for b in tl_bullets:
        p_b = tf_tl.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_SECONDARY if "•" in b else ROSE_ACCENT
        p_b.space_before = Pt(3)

    # [3] BOTTOM-LEFT: VULNERABILITY SCANNERS & MANUAL PENTEST
    add_card(sp, Inches(1.0), Inches(4.35), Inches(5.4), Inches(2.3))
    tb_bl = sp.shapes.add_textbox(Inches(1.2), Inches(4.45), Inches(5.0), Inches(2.1))
    tf_bl = tb_bl.text_frame
    tf_bl.word_wrap = True
    p = tf_bl.paragraphs[0]
    p.text = "Tenable Nessus / Qualys / Konsultan Pentest Manual"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    bl_bullets = [
        "• Pemindaian Ribuan Database CVE Mendalam untuk Audit Tahunan",
        "✖ Hanya Menghasilkan Laporan PDF Statis Tanpa Proteksi Runtime Tepi",
        "✖ Biaya Sewa Lisensi/Jasa Mahal & Penanganan Celah Lambat (>14 Hari)"
    ]
    for b in bl_bullets:
        p_b = tf_bl.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_SECONDARY if "•" in b else ROSE_ACCENT
        p_b.space_before = Pt(3)

    # [4] BOTTOM-RIGHT: BASIC PASSIVE TOOLS & AI CHAT
    add_card(sp, Inches(6.9), Inches(4.35), Inches(5.43), Inches(2.3))
    tb_br = sp.shapes.add_textbox(Inches(7.1), Inches(4.45), Inches(5.1), Inches(2.1))
    tf_br = tb_br.text_frame
    tf_br.word_wrap = True
    p = tf_br.paragraphs[0]
    p.text = "Microsoft Security Copilot / Plugin WAF Dasar"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    br_bullets = [
        "• Asisten AI Chat Cerdas & Filter Web Hosting Lokal",
        "✖ Hanya Saran Teks/Kode Tanpa Kemampuan Blokir Trafik Tepi Instan",
        "✖ Sering False-Positive, Beban Token Tinggi & Memperberat Server"
    ]
    for b in br_bullets:
        p_b = tf_br.add_paragraph()
        p_b.text = b
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_SECONDARY if "•" in b else ROSE_ACCENT
        p_b.space_before = Pt(3)

# ==========================================
# SLIDE 6: MARKET SIZE (TAM/SAM/SOM) & MARKET POSITIONING
# ==========================================
def populate_slide_6(s6, w, h):
    set_slide_background(s6, w, h)
    add_header(s6, "Peluang Pasar (TAM, SAM, SOM) & Posisi Pasar", "CORE 6/14", 
               "Pasar keamanan digital perbankan dan fintech Indonesia yang berkembang pesat.", "06")

    # Left: Market Sizing (TAM, SAM, SOM)
    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_ms = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_ms = tb_ms.text_frame
    tf_ms.word_wrap = True

    p = tf_ms.paragraphs[0]
    p.text = "ESTIMASI UKURAN PASAR (MARKET SIZING)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    ms_items = [
        ("TAM (Total Addressable Market) — Rp 8,4 Triliun", "Total pasar belanja keamanan siber dan cloud protection di Indonesia pada 2026 (CAGR 18.5%)."),
        ("SAM (Serviceable Available Market) — Rp 1,9 Triliun", "Segmen perbankan digital, 1.400+ BPR/BPRS, 100+ Fintech P2P/Payment, dan BUMD yang diwajibkan kepatuhan POJK 30/2025."),
        ("SOM (Serviceable Obtainable Market) — Rp 48 Miliar", "Target penetrasi awal 5% pada ekosistem BPR & Fintech terlisensi dalam 3 tahun pertama operasi komersial.")
    ]
    for head, desc in ms_items:
        p_h = tf_ms.add_paragraph()
        p_h.text = head
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_ms.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # Right: Guardrail-as-a-Service (GaaS) Autonomous Sentinel
    add_card(s6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_mp = s6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_mp = tb_mp.text_frame
    tf_mp.word_wrap = True

    p = tf_mp.paragraphs[0]
    p.text = "KATEGORI BARU: GUARDRAIL-AS-A-SERVICE (GaaS)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p_sub = tf_mp.add_paragraph()
    p_sub.text = "Layanan Pengaman Siap Pakai — Agen Cerdas yang Mengawasi & Melindungi 24/7"
    p_sub.font.size = Pt(9)
    p_sub.font.italic = True
    p_sub.font.color.rgb = TEXT_SECONDARY
    p_sub.space_before = Pt(2)

    gaas_pillars = [
        ("1. Uji Titik Lemah Mandiri (Auto-Pentest)", "Secara rutin memeriksa dan menguji keamanan website sendiri sebelum ditemukan peretas."),
        ("2. Perisai Pelindung Selalu Aktif (Real-Time Guard)", "Menahan serangan berbahaya seketika dalam hitungan milidetik tanpa memperlambat akses nasabah."),
        ("3. Penyembuh Mandiri (Antibodi Digital)", "Menutup celah keamanan dan memulihkan website otomatis jika terjadi serangan atau defacement."),
        ("4. Perlindungan Celah Baru (Anti Zero-Day)", "Langsung beradaptasi kebal terhadap pola serangan baru tanpa perlu menunggu update manual.")
    ]
    for head, desc in gaas_pillars:
        p_h = tf_mp.add_paragraph()
        p_h.text = head
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(5)
        p_d = tf_mp.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY

# ==========================================
# SLIDE 7: CARA KERJA SOLUSI (BUSINESS-FRIENDLY 3-STEP FLOW)
# ==========================================
def populate_slide_7(s7, w, h):
    set_slide_background(s7, w, h)
    add_header(s7, "Bagaimana Solusi Bekerja: 3 Langkah Mudah Tanpa Ubah Kode", "CORE 7/14", 
               "Implementasi instan tanpa merombak infrastruktur backend perbankan yang sudah berjalan.", "07")

    steps_biz = [
        ("LANGKAH 1: KONEKSI DALAM 5 MENIT", "Zero-Downtime Reverse Proxy", CYAN_ACCENT, [
            "Cukup arahkan CNAME/DNS domain kanal web/API ke WAF Gateway Nexus (:8080).",
            "Tidak perlu mengubah satu baris pun kode sumber backend aplikasi yang sudah ada.",
            "Kompatibel langsung dengan cloud apa pun (AWS, Google Cloud, Vercel, VPS Lokal, On-Premises)."
        ]),
        ("LANGKAH 2: PERTAHANAN & WASIT OTOMATIS", "Dual-Brain AI + NEX-RED Cowork", EMERALD_ACCENT, [
            "Trafik serangan (SQLi, XSS, Deface Judi) diblokir instan di tepi dalam < 1.2 milidetik.",
            "Wasit cerdas menguji ketahanan server origin dan secara otomatis melahirkan antibodi baru.",
            "Tembakan ulang (replay verification) membuktikan celah benar-benar tertutup rapat."
        ]),
        ("LANGKAH 3: LAPORAN AUDIT 1-KLIK", "Kepatuhan Instan POJK 30/2025", NAVY_SECONDARY, [
            "Dashboard Command Center (:3001) menyajikan telemetri ancaman dan status risiko real-time.",
            "Ekspor laporan kepatuhan terstruktur (PDF/JSON) siap diserahkan kepada auditor OJK & BSSN.",
            "Notifikasi peringatan insiden langsung terkirim ke kanal Telegram tim operasional."
        ])
    ]

    for i, (stitle, ssub, color, items) in enumerate(steps_biz):
        x = Inches(0.8 + i * 4.0)
        add_card(s7, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s7.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        p_s = tf.add_paragraph()
        p_s.text = ssub
        p_s.font.size = Pt(12)
        p_s.font.bold = True
        p_s.font.color.rgb = NAVY_PRIMARY
        p_s.space_before = Pt(2)

        for item in items:
            p_i = tf.add_paragraph()
            p_i.text = f"• {item}"
            p_i.font.size = Pt(9.5)
            p_i.font.color.rgb = TEXT_SECONDARY
            p_i.space_before = Pt(8)

# ==========================================
# SLIDE 8: GO-TO-MARKET (GTM) STRATEGY
# ==========================================
def populate_slide_8(s8, w, h):
    set_slide_background(s8, w, h)
    add_header(s8, "Strategi Go-To-Market (GTM) & Skalasi Distribusi", "CORE 8/14", 
               "3 Jalur penetrasi pasar terarah untuk mengakuisisi institusi keuangan secara agresif.", "08")

    gtm_pillars = [
        ("PILAR 1: KEMITRAAN ASOSIASI BPR & FINTECH", EMERALD_ACCENT, [
            ("Target Channel", "Asosiasi Perbankan Daerah (Perbarindo) & Asosiasi Fintech Indonesia (Aftech)."),
            ("Model Program", "Program 'Klinik Kepatuhan POJK 30/2025' — Menawarkan audit celah gratis (L0 Job) untuk mengonversi menjadi langganan Loop Retainer berbayar."),
            ("Kelebihan", "Biaya akuisisi (CAC) sangat rendah karena penetrasi dilakukan secara kolektif.")
        ]),
        ("PILAR 2: KONSORSIUM AGENSI WEB FINANSIAL", CYAN_ACCENT, [
            ("Target Channel", "Software House & Agensi Pengembang Website/Portal Pemda & BPR."),
            ("Model Program", "Program Reseller / Co-Selling — Agensi memaketkan perlindungan Nexus Cyber langsung saat menjual jasa pembuatan web perbankan."),
            ("Kelebihan", "Ekspansi distribusi eksponensial tanpa perlu menambah banyak tim sales internal.")
        ]),
        ("PILAR 3: PRODUCT-LED GROWTH (STARTER TIER)", NAVY_SECONDARY, [
            ("Target Channel", "UMKM & Startup Digital yang mencari solusi web aman berbiaya murah."),
            ("Model Program", "Channel Starter (Rp 20rb) — Generator web instan + Header Shield di tepi sebagai pintu masuk (entry layer)."),
            ("Kelebihan", "Menciptakan basis pengguna akar rumput masif yang siap di-upsell ke tier Edge Shield & GaaS.")
        ])
    ]

    for i, (title, color, items) in enumerate(gtm_pillars):
        x = Inches(0.8 + i * 4.0)
        add_card(s8, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s8.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = color

        for h_text, d_text in items:
            p_h = tf.add_paragraph()
            p_h.text = h_text
            p_h.font.size = Pt(11)
            p_h.font.bold = True
            p_h.font.color.rgb = NAVY_PRIMARY
            p_h.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = d_text
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(2)

# ==========================================
# SLIDE 9: SOCIAL IMPACT & FINANCIAL INCLUSION (ESG)
# ==========================================
def populate_slide_9(s9, w, h):
    set_slide_background(s9, w, h)
    add_header(s9, "Dampak Sosial, Inklusi Keuangan, & Keberlanjutan (ESG)", "CORE 9/14", 
               "Demokratisasi keamanan siber kelas enterprise untuk melindungi nasabah akar rumput.", "09")

    impacts = [
        ("1. PERLINDUNGAN NASABAH AKAR RUMPUT", EMERALD_ACCENT, [
            ("Mencegah Kebocoran Rekening Rakyat", "Jutaan nasabah BPR dan UMKM di daerah rentan menjadi korban pencurian saldo akibat defacement & credential stuffing pada portal web bank daerah."),
            ("Dampak Nyata", "Nexus Cyber menjamin data simpanan dan pinjaman masyarakat kecil terlindungi 24/7 tanpa risiko disusupi malware/judi online.")
        ]),
        ("2. PENGUATAN KETAHANAN BPR & BUMD DAERAH", CYAN_ACCENT, [
            ("Menyelamatkan Bank Kecil dari Kebangkrutan", "Satu insiden siber dapat mematikan operasional BPR daerah karena biaya pemulihan yang sangat mahal."),
            ("Dampak Nyata", "Efisiensi biaya 90% memungkinkan bank terkecil sekalipun memiliki pertahanan setara bank multinasional, menjaga stabilitas ekonomi lokal.")
        ]),
        ("3. KEPERCAYAAN EKOSISTEM DIGITAL NASIONAL", NAVY_SECONDARY, [
            ("Mendukung Visi Bank Indonesia & OJK", "Kanal pembayaran digital yang aman mendorong akselerasi inklusi keuangan nasional dan elektronifikasi transaksi pemerintah daerah."),
            ("Dampak Nyata", "Meningkatkan skor indeks keamanan siber nasional dan memperkuat daya saing industri keuangan digital Indonesia di kancah ASEAN.")
        ])
    ]

    for i, (title, color, items) in enumerate(impacts):
        x = Inches(0.8 + i * 4.0)
        add_card(s9, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s9.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = color

        for h_text, d_text in items:
            p_h = tf.add_paragraph()
            p_h.text = h_text
            p_h.font.size = Pt(11)
            p_h.font.bold = True
            p_h.font.color.rgb = NAVY_PRIMARY
            p_h.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = d_text
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(2)

# ==========================================
# SLIDE 10A: EMPIRICAL TESTING, 10 EXPERT UAT & PERFORMANCE METRICS (CORE 10A/14)
# ==========================================
def populate_slide_10a(s10a, w, h):
    set_slide_background(s10a, w, h)
    add_header(s10a, "Hasil Uji Lapangan, UAT 10 Pakar Siber & Metrik Efektivitas Teknis", "CORE 10A/14", 
               "Teruji tangguh pada ekosistem akademik 2 universitas dan divalidasi oleh praktisi keamanan siber.", "10A")

    cards_data_s10a = [
        ("UJI COBA 2 UNIVERSITAS (TELKOM & UNMUL)", EMERALD_ACCENT, [
            ("Lingkungan Uji Nyata", "Simulasi portal akademik, API layanan mahasiswa, dan server riset publik di Telkom University & Universitas Mulawarman."),
            ("Stress Test 50.000+ Request", "Diuji dengan lonjakan payload serangan masif (OWASP Top 10, SQLi, Defacement Judi Online, dan Zero-Day)."),
            ("Hasil Ketahanan 99.6%", "Trafik serangan diblokir instan di tepi dengan latensi inspeksi 1.18 ms dan Zero Downtime.")
        ]),
        ("EVALUASI 10 PAKAR SIBER (UAT & SUS)", CYAN_ACCENT, [
            ("Panel 10 Praktisi Siber", "Blind testing & User Acceptance Testing (UAT) oleh 10 praktisi siber (SOC Analyst, Pentester, & Dosen Cybersecurity)."),
            ("Skor SUS: 89.2 / 100 (Grade A+)", "Mendapatkan predikat 'Exceptional Usability' untuk kemudahan navigasi dan integrasi operasional."),
            ("10/10 (100%) Pakar Sepakat", "Dashboard SOC (:3001) & tombol 1-klik laporan POJK 30/2025 mempercepat respon penanganan insiden hingga 85%.")
        ]),
        ("METRIK KINERJA & EFEKTIVITAS (UAS)", NAVY_SECONDARY, [
            ("99.6% Teratasi Otonom", "Mitigasi serangan berbahaya tertangani mandiri tanpa memerlukan campur tangan staf manusia."),
            ("Virtual Patching < 90 Detik", "99% lebih cepat dibanding siklus penambalan manual tradisional (> 14 hari)."),
            ("False Positive < 0.002%", "Akurasi penyaringan sangat presisi berkat sinergi Dual-Brain AI (Reflex & Reasoning)."),
            ("Wasit Replay Proof 100%", "Celah terverifikasi tuntas tertutup rapat (Status CLOSED_OK) melalui tembakan ulang otonom.")
        ])
    ]

    for i, (title, color, bullets) in enumerate(cards_data_s10a):
        x = Inches(0.8 + i * 4.0)
        add_card(s10a, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s10a.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color

        for b_title, b_desc in bullets:
            p_bt = tf.add_paragraph()
            p_bt.text = b_title
            p_bt.font.size = Pt(10)
            p_bt.font.bold = True
            p_bt.font.color.rgb = NAVY_PRIMARY
            p_bt.space_before = Pt(5)

            p_bd = tf.add_paragraph()
            p_bd.text = b_desc
            p_bd.font.size = Pt(8.8)
            p_bd.font.color.rgb = TEXT_SECONDARY
            p_bd.space_before = Pt(1)

# ==========================================
# SLIDE 10B: ADOPTER TRACTION & WILLINGNESS TO PAY (CORE 10B/14)
# ==========================================
def populate_slide_10b(s10b, w, h):
    set_slide_background(s10b, w, h)
    add_header(s10b, "Traksi Pengguna Nyata, Calon Mitra & Willingness to Pay", "CORE 10B/14", 
               "Minat adopsi konkret dari 3 instansi pendidikan & 2 UMKM dengan komitmen pembayaran 100%.", "10B")

    cards_data_s10b = [
        ("3 INSTANSI PENDIDIKAN & RISET", CYAN_ACCENT, [
            ("1. PUTI Telkom University", "Pusat Teknologi Informasi — Kebutuhan proteksi portal akademik & API layanan kampus dari serangan siber harian."),
            ("2. Universitas Mulawarman", "Perlindungan portal publik & server riset daerah dari kebocoran data dan defacement judi online."),
            ("3. UKM Cyber", "Validasi modul wasit, riset pentesting kolaboratif, dan pengujian live sandbox."),
            ("Willingness to Pay (WTP)", "Siap mengadopsi skema Retainer Edukasi (Rp 300rb - 500rb/bln) menghemat 90% biaya SOC.")
        ]),
        ("2 MITRA UMKM RIIL", EMERALD_ACCENT, [
            ("1. Mandiri Jaya (Belum Punya Web)", "Kasus Baru: Dibuatkan website profil bisnis instan yang otomatis dipagari keamanan siber sejak hari pertama."),
            ("2. Win Elektronik (Sudah Punya Web)", "Kasus Existing: Pemasangan Header Shield anti-defacement judi online pada toko online tanpa pindah hosting."),
            ("Daya Tarik: 2 Opsi Fleksibel", "UMKM baru langsung go-digital aman; UMKM berwebsite diproteksi tanpa pusing teknis."),
            ("Willingness to Pay (WTP)", "Bersedia membayar paket Starter Rp 20.000 – Rp 35.000/bulan.")
        ]),
        ("VALIDASI KOMERSIAL & KONVERSI", NAVY_SECONDARY, [
            ("100% Willingness to Pay", "Seluruh calon mitra yang diwawancarai menyetujui struktur harga yang ditawarkan."),
            ("5 Komitmen Pilot (LOI)", "Siap menandatangani Letter of Intent (LOI) untuk implementasi Sandbox Q1 2026."),
            ("Model Dual-Branch Terbukti", "Solusi mencakup 100% kebutuhan UMKM baik yang belum maupun sudah memiliki website.")
        ])
    ]

    for i, (title, color, bullets) in enumerate(cards_data_s10b):
        x = Inches(0.8 + i * 4.0)
        add_card(s10b, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s10b.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for b_title, b_desc in bullets:
            p_bt = tf.add_paragraph()
            p_bt.text = b_title
            p_bt.font.size = Pt(10)
            p_bt.font.bold = True
            p_bt.font.color.rgb = NAVY_PRIMARY
            p_bt.space_before = Pt(5)

            p_bd = tf.add_paragraph()
            p_bd.text = b_desc
            p_bd.font.size = Pt(9)
            p_bd.font.color.rgb = TEXT_SECONDARY
            p_bd.space_before = Pt(1)

# ==========================================
# SLIDE 11: FINANCIAL MODEL & 3-YEAR PROJECTIONS
# ==========================================
def populate_slide_11(s11, w, h):
    set_slide_background(s11, w, h)
    add_header(s11, "Model Bisnis 3 Segmen & Proyeksi Finansial 3 Tahun", "CORE 11/14", 
               "Monetisasi terukur berbasis volume dan langganan berulang dengan margin kotor tinggi (~78%).", "11")

    # Left: 3 Business Model Segments
    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_l = s11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "3 SEGMEN MODEL BISNIS & UNIT ECONOMICS"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    streams = [
        ("1. UMKM / Usaha Mikro (Volume-Based)", "Rp 20.000 – 35.000 / bln", "• Belum Punya Web: Auto-generate Website Baru + Shield Terpasang.\n• Sudah Punya Web: Header Shield Instan tanpa Pindah Hosting.\nCOGS Rp 5.000/bln · Gross Margin 85% · Target: 300 -> 5.000 UMKM."),
        ("2. BPR, Fintech & Kampus (Loop Retainer)", "Rp 300.000 – 500.000 / bln", "Wasit otomatis + virtual patching instan + kepatuhan POJK. COGS Rp 70.000/bln · Margin 78% · Target: 35 -> 400+ Institusi."),
        ("3. B2G & Enterprise (On-Premises License)", "Rp 25 – 50 Juta / thn", "Dedicated Edge Engine on-prem + SLA support + audit forensik. COGS Rp 3 Jt/thn · Margin 90% · Target: 3 -> 30 Kontrak.")
    ]
    for head, pr, d in streams:
        p_h = tf_l.add_paragraph()
        p_h.text = f"{head} — {pr}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(6)
        p_d = tf_l.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY

    # Right: 3-Year Projections
    add_card(s11, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_r = s11.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "PROYEKSI PENDAPATAN & PROFITABILITAS 3 TAHUN"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    proj = [
        ("Tahun 1 (2026): Validasi & Pilot", "Pendapatan: Rp 450 Juta (Gross Profit Rp 345 Jt)", "Target: 35 BPR/Fintech + 300 UMKM + 3 Pilot B2G. Gross Margin 76.6%."),
        ("Tahun 2 (2027): Ekspansi B2B Nasional", "Pendapatan: Rp 1,85 Miliar (EBITDA Rp 720 Jt)", "Target: 150 BPR/Fintech + 1.500 UMKM + 12 Enterprise. EBITDA Positif."),
        ("Tahun 3 (2028): Penetrasi B2G & Skalasi", "Pendapatan: Rp 5,20 Miliar (EBITDA Rp 2,8 Miliar)", "Target: 400+ BPR/Fintech + 5.000 UMKM + 30 Enterprise."),
        ("Unit Economics Sehat", "Gross Margin 78% · LTV/CAC 7.2x · Payback 4.5 Bulan", "Biaya operasional server sangat efisien berkat Go Concurrency Engine.")
    ]
    for head, pr, d in proj:
        p_h = tf_r.add_paragraph()
        p_h.text = head
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(5)
        p_sub = tf_r.add_paragraph()
        p_sub.text = pr
        p_sub.font.size = Pt(9.5)
        p_sub.font.bold = True
        p_sub.font.color.rgb = EMERALD_ACCENT
        p_d = tf_r.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_SECONDARY

# ==========================================
# SLIDE 13: ROADMAP, ROI & FUNDING ASK (CORE 13/14)
# ==========================================
def populate_slide_13(s13, w, h):
    set_slide_background(s13, w, h)
    add_header(s13, "Roadmap, Analisis ROI/BEP, dan Kebutuhan Pendanaan", "CORE 13/14", 
               "Periode pengembalian investor 18–24 bulan (Maks 2.5 tahun) dan alokasi pendanaan Rp 500 Juta.", "13")

    # Left: ROI Analysis & BEP
    add_card(s13, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_roi = s13.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_roi = tb_roi.text_frame
    tf_roi.word_wrap = True

    p = tf_roi.paragraphs[0]
    p.text = "ANALISIS ROI INVESTOR & BREAK-EVEN POINT (BEP)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    roi_items = [
        ("ROI Klien Lembaga Keuangan (< 3 Bulan)", "Penghematan rata-rata kerugian insiden Rp 3,8 Miliar & denda regulasi dibanding biaya langganan Loop Rp 3,6 Jt/thn $\rightarrow$ ROI Klien > 1.000%."),
        ("Titik Impas (Break-Even Point / BEP)", "Tercapai pada Bulan ke-14 (Awal Tahun ke-2) dengan 35 klien institusi aktif + 300 UMKM (menutup OPEX Rp 30 Jt/bln)."),
        ("Periode Pengembalian Investor (18 – 24 Bulan)", "Modal investasi tahap awal kembali penuh dalam 1.5 – 2.0 Tahun (di bawah batas maksimal 2.5 tahun) berkat margin kotor 78%."),
        ("Potensi Return Tahun ke-3 (3.0x – 3.5x MOIC)", "Berbasis EBITDA Rp 2,8 Miliar dengan multiplier valuasi SaaS keamanan siber 5x – 8x.")
    ]
    for head, d in roi_items:
        p_h = tf_roi.add_paragraph()
        p_h.text = f"✔ {head}"
        p_h.font.size = Pt(10)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(5)
        p_d = tf_roi.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # Right: Funding Ask (Rp 500 Juta) & Use of Funds
    add_card(s13, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_rd = s13.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_rd = tb_rd.text_frame
    tf_rd.word_wrap = True

    p = tf_rd.paragraphs[0]
    p.text = "FUNDING ASK: RP 500.000.000 & ALOKASI DANA"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    p_target = tf_rd.add_paragraph()
    p_target.text = "Target Pendanaan: Rp 500 Juta (Runway 18 Bulan Menuju Profitabel)"
    p_target.font.size = Pt(9.5)
    p_target.font.bold = True
    p_target.font.color.rgb = EMERALD_ACCENT
    p_target.space_before = Pt(3)

    allocations = [
        ("1. 40% (Rp 200 Juta) — Product R&D & AI Infrastructure", "Server edge proxy cluster, optimasi runtime model AI, hardware testbed, dan audit kepatuhan ISO 27001 / BSSN."),
        ("2. 35% (Rp 175 Juta) — Sales, Partnership & GTM", "Klinik kepatuhan bersama Perbarindo & Aftech, roadshow instansi daerah, dan program kemitraan reseller."),
        ("3. 15% (Rp 75 Juta) — Legalitas, Perizinan & OJK Sandbox", "Pendaftaran HAKI, izin PSE Kominfo, dan pendampingan kepatuhan OJK Regulatory Sandbox."),
        ("4. 10% (Rp 50 Juta) — Cadangan Kas / Dana Darurat (Reserve)", "Buffer kas operasional untuk menjamin kestabilan likuiditas selama 18 bulan.")
    ]
    for q_title, q_desc in allocations:
        p_h = tf_rd.add_paragraph()
        p_h.text = q_title
        p_h.font.size = Pt(9.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(4)
        p_d = tf_rd.add_paragraph()
        p_d.text = q_desc
        p_d.font.size = Pt(8.5)
        p_d.font.color.rgb = TEXT_SECONDARY

# ==========================================
# PART B: 8 APPENDIX SLIDES
# ==========================================
def add_appendix_slides(prs):
    blank_slide_layout = prs.slide_layouts[6]
    w, h = prs.slide_width, prs.slide_height

    # SLIDE A1: DEFENSE DELTA
    sa1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa1, w, h)
    add_header(sa1, "Appendix A1: Logika Wasit Defense Delta & Matriks Status Penutupan", "ADDITIONAL · Q&A", 
               "Detail taksonomi status penutupan Job Cowork untuk pembuktian ilmiah wasit jujur.", "A1")

    labels = [
        ("origin_open", "Celah Terbuka di Backend", AMBER_ACCENT, "Origin server menerima payload serangan; WAF belum memblokir. Risiko nyata teridentifikasi."),
        ("waf_blocked", "Ditahan di Tepi WAF", CYAN_ACCENT, "WAF berhasil memblokir serangan (403), meskipun backend origin mungkin masih memiliki kelemahan."),
        ("both_held", "Keduanya Menahan", EMERALD_ACCENT, "Baik WAF maupun origin backend sama-sama menolak payload serangan secara aman."),
        ("replay_held", "Tembakan Ulang Tertahan", EMERALD_ACCENT, "Setelah antibodi dipasang, tembakan ulang (replay attack) tetap berhasil ditolak 403."),
        ("replay_missed", "Uji Ulang Gagal", AMBER_ACCENT, "Antibodi gagal menahan variasi serangan baru. Job otomatis ditutup sebagai CLOSED_GAP (Residual Eksplisit)."),
        ("antibody_learned", "Antibodi Terverifikasi", NAVY_SECONDARY, "Antibodi zero-day sukses disimpan di memori imun host (Postgres) dan di-share ke RAM cache.")
    ]

    for i, (l_name, l_title, color, l_desc) in enumerate(labels):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(1.8 + row * 2.5)
        add_card(sa1, x, y, Inches(3.7), Inches(2.3))
        tb = sa1.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.4), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Label: {l_name}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        p_t = tf.add_paragraph()
        p_t.text = l_title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_PRIMARY
        p_t.space_before = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = l_desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(4)

    # SLIDE A2: AI 8-POINT MANDATORY STANDARD
    sa2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa2, w, h)
    add_header(sa2, "Appendix A2: Standar Teknis Wajib AI & Machine Learning (8 Poin PIDI)", "ADDITIONAL · Q&A", 
               "Pembuktian transparansi sistem AI: Bukan sekadar gimmick 'Powered by AI'.", "A2")

    ai_points = [
        ("1. Input Data", "Payload HTTP lengkap: Header, URL Query, Body, Client IP, Method."),
        ("2. Sumber Data", "Dataset serangan web riil (OWASP Top 10, CVE feeds) & traffic logs lab."),
        ("3. Preprocessing Data", "Normalisasi kanonik (stripping comments, URL decode, UTF-8 normalization)."),
        ("4. Output Model", "Klasifikasi biner (Attack/Legit), Threat Type, dan Confidence Score (0.0-1.0)."),
        ("5. Pemanfaatan Keputusan", "Memicu respon 403 instan atau merumuskan draf antibodi virtual patch baru."),
        ("6. Metrik Performa", "Latensi Reflex <1.2ms, Zero False Positive pada 10.000 legitimate HTTP requests."),
        ("7. Pengujian & Limitasi", "Diuji via live replay attack; memerlukan fallback jika model lokal belum siap."),
        ("8. Human Oversight", "Gerbang persetujuan L0/L1 wajib diverifikasi operator sebelum deploy produksi.")
    ]

    for i, (point_title, point_desc) in enumerate(ai_points):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 1.25)
        add_card(sa2, x, y, Inches(5.7), Inches(1.15))
        tb = sa2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), Inches(5.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = point_title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

        p_d = tf.add_paragraph()
        p_d.text = point_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(2)

    # SLIDE A3: DUAL-BRAIN AI ARCHITECTURE
    sa3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa3, w, h)
    add_header(sa3, "Appendix A3: Arsitektur Dual-Brain AI (Otak Kiri vs Otak Kanan)", "ADDITIONAL · Q&A", 
               "Pemisahan tugas komputasi sinkron latensi rendah dengan penalaran asinkron mendalam.", "A3")

    add_card(sa3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_db1 = sa3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_db1 = tb_db1.text_frame
    tf_db1.word_wrap = True

    p = tf_db1.paragraphs[0]
    p.text = "OTAK KIRI: REFLEX FILTER (nex-ai-reflex)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    items_db1 = [
        ("Eksekusi Sinkron", "Berjalan langsung pada alur request HTTP utama di internal/ai/reflex_filter.go."),
        ("Latensi Ultra-Rendah", "Kecepatan eksekusi < 1.2 milidetik menggunakan regex terkompilasi."),
        ("Deteksi Pola Klasik", "Menangani SQL Injection, Cross-Site Scripting (XSS), Path Traversal, RCE, dan defacement."),
        ("RAM-First Cache", "Pola antibodi yang sudah dipelajari langsung dimasukkan ke memori RAM untuk blokir instan 403.")
    ]
    for head, d in items_db1:
        p_h = tf_db1.add_paragraph()
        p_h.text = f"✔ {head}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_db1.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    add_card(sa3, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_db2 = sa3.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_db2 = tb_db2.text_frame
    tf_db2.word_wrap = True

    p = tf_db2.paragraphs[0]
    p.text = "OTAK KANAN: REASONING ENGINE (nex-ai-protect)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    items_db2 = [
        ("Eksekusi Asinkron", "Dipanggil via Goroutine di latar belakang (timeout 30s) tanpa memblokir trafik pengguna."),
        ("Penalaran Mendalam", "Menganalisis intensi peretas, ancaman persisten tingkat lanjut (APT), dan anomali zero-day."),
        ("Formulasi Antibodi", "Menghasilkan aturan virtual patch baru dan menyimpannya ke tabel PostgreSQL antibody_audits."),
        ("Kemandirian Model Lokal", "Berjalan pada Ollama lokal host; dilarang keras fallback ke model publik eksternal.")
    ]
    for head, d in items_db2:
        p_h = tf_db2.add_paragraph()
        p_h.text = f"⚡ {head}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_db2.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # SLIDE A4: COMPLIANCE MAPPING
    sa4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa4, w, h)
    add_header(sa4, "Appendix A4: Pemetaan Kepatuhan Regulasi (POJK 30/2025 & ISO 27001)", "ADDITIONAL · Q&A", 
               "Kesesuaian langsung arsitektur Nexus Cyber dengan pasal-pasal regulasi perbankan.", "A4")

    add_card(sa4, Inches(0.8), Inches(1.8), Inches(11.73), Inches(5.1))
    tb_comp = sa4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.7))
    tf_comp = tb_comp.text_frame
    tf_comp.word_wrap = True

    comp_items = [
        ("POJK No. 30/2025 Pasal Ketahanan Siber", "Mewajibkan penyelenggara ITSK memiliki sistem deteksi dan mitigasi insiden siber secara aktif.", "Nexus Cyber menyediakan proteksi tepi always-on dan virtual patch otomatis dalam hitungan menit."),
        ("POJK No. 30/2025 Pasal Uji Ketahanan Berkala", "Mewajibkan simulasi serangan siber dan uji penetrasi berkala.", "Wasit NEX-RED menjalankan Job Cowork berkala dengan pengujian tembakan ulang (replay verification) otomatis."),
        ("ISO 27001 Kontrol A.12.4 (Logging & Monitoring)", "Mewajibkan pencatatan log keamanan yang aman dari manipulasi dan dapat diaudit.", "Tabel PostgreSQL threat_logs menyimpan rekaman telemetri dengan UUID v4, SHA-256 prev_hash, dan target domain index."),
        ("UU PDP No. 27/2022 Pasal Perlindungan Kerahasiaan", "Mewajibkan enkripsi data nasabah dari risiko intersepsi dan kebocoran.", "PQC Shield berbasis NIST ML-KEM-768 melindungi payload komunikasi data sensitif dari ancaman komputasi kuantum.")
    ]

    for i, (reg, mand, sol) in enumerate(comp_items):
        p_r = tf_comp.paragraphs[0] if i == 0 else tf_comp.add_paragraph()
        p_r.text = reg
        p_r.font.size = Pt(11)
        p_r.font.bold = True
        p_r.font.color.rgb = NAVY_PRIMARY
        if i > 0: p_r.space_before = Pt(8)

        p_m = tf_comp.add_paragraph()
        p_m.text = f"• Mandat Regulasi: {mand}"
        p_m.font.size = Pt(9.5)
        p_m.font.color.rgb = TEXT_SECONDARY

        p_s = tf_comp.add_paragraph()
        p_s.text = f"✔ Implementasi Nexus: {sol}"
        p_s.font.size = Pt(9.5)
        p_s.font.bold = True
        p_s.font.color.rgb = EMERALD_ACCENT

    # SLIDE A5: DEMO SOP & 3-TIER BACKUP PLAN
    sa5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa5, w, h)
    add_header(sa5, "Appendix A5: SOP Demonstrasi Live & 3-Tier Demo Backup Plan", "ADDITIONAL · Q&A", 
               "Protokol mitigasi kegagalan teknis saat presentasi panggung offline maupun online.", "A5")

    tiers = [
        ("TIER 1: LIVE DEMO (PILIHAN UTAMA)", EMERALD_ACCENT, [
            ("Lingkungan Demo", "WAF Gateway lokal :8080 terhubung ke target live portofolio Vercel."),
            ("SOP Alur (90 Detik)", "1. Akses halaman target normal → 2. Tembakkan payload serangan (SQLi/Deface) → 3. WAF memblokir 403 → 4. Buka Command Center SOC (:3001) melihat telemetri log."),
            ("Pemeriksaan Pra-Demo", "Pastikan Docker Desktop aktif dan model Ollama nex-ai-* siap.")
        ]),
        ("TIER 2: VIDEO OFFLINE (CADANGAN 1)", CYAN_ACCENT, [
            ("Format File", "Rekaman MP4 1080p durasi 75 detik tersimpan di harddisk lokal laptop."),
            ("Skenario Penggunaan", "Diputar langsung jika koneksi internet venue/hotspot mengalami gangguan."),
            ("Konten Rekaman", "Screen-recording proses deteksi serangan, notifikasi Telegram, dan audit trail.")
        ]),
        ("TIER 3: SCREENSHOT FLOW (CADANGAN 2)", AMBER_ACCENT, [
            ("Format", "Rangkaian screenshot alur langkah demi langkah pada slide appendix."),
            ("Skenario Penggunaan", "Digunakan jika laptop mengalami kendala pemutaran media video."),
            ("Konten", "Gambar input payload, respon 403 terminal, dan visualisasi SOC dashboard.")
        ])
    ]

    for i, (ttitle, color, items) in enumerate(tiers):
        x = Inches(0.8 + i * 4.0)
        add_card(sa5, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = sa5.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ttitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for head, d in items:
            p_h = tf.add_paragraph()
            p_h.text = head
            p_h.font.size = Pt(10.5)
            p_h.font.bold = True
            p_h.font.color.rgb = NAVY_PRIMARY
            p_h.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = d
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(2)

    # SLIDE A6: DETAILED FINANCIAL MODEL & COGS
    sa6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa6, w, h)
    add_header(sa6, "Appendix A6: Struktur Biaya Pokok (COGS) & Proyeksi LTV / CAC", "ADDITIONAL · Q&A", 
               "Analisis profitabilitas unit economics untuk meyakinkan investor dan juri bisnis.", "A6")

    add_card(sa6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_cogs = sa6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_cogs = tb_cogs.text_frame
    tf_cogs.word_wrap = True

    p = tf_cogs.paragraphs[0]
    p.text = "STRUKTUR BIAYA POKOK LAYANAN (COGS PER HOST)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    cogs_items = [
        ("Komputasi Gateway Go (VPS Cloud)", "Rp 45.000 / bln", "High-efficiency Go runtime memungkinkan 1 VPS Rp 200rb melayani 4-5 host kanal."),
        ("Inferensi AI Lokal (Ollama)", "Rp 15.000 / bln", "Zero API token cost berkat bobot model mandiri nex-ai-protect."),
        ("Storage Audit Log (PostgreSQL)", "Rp 10.000 / bln", "Penyimpanan log terindeks dengan efisiensi kompresi data tinggi."),
        ("Total COGS per Host", "Rp 70.000 / bln", "Harga Jual Loop: Rp 300.000 / bln $\rightarrow$ Gross Margin 76.6%.")
    ]
    for head, pr, d in cogs_items:
        p_h = tf_cogs.add_paragraph()
        p_h.text = f"{head}: {pr}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_cogs.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    add_card(sa6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_ltv = sa6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_ltv = tb_ltv.text_frame
    tf_ltv.word_wrap = True

    p = tf_ltv.paragraphs[0]
    p.text = "METRIK SKALABILITAS BISNIS (LTV / CAC)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    ltv_items = [
        ("Customer Acquisition Cost (CAC)", "Rp 1.500.000", "Diperoleh melalui kemitraan asosiasi BPR/Fintech dan program pilot terarah."),
        ("Customer Lifetime Value (LTV)", "Rp 10.800.000", "Berdasarkan masa retensi rata-rata 36 bulan dengan skema retainership bulanan."),
        ("Rasio LTV / CAC", "7.2x (Sangat Sehat)", "Standar industri SaaS B2B ideal adalah > 3.0x."),
        ("Payback Period", "4.5 Bulan", "Modal akuisisi kembali dalam waktu kurang dari satu semester.")
    ]
    for head, pr, d in ltv_items:
        p_h = tf_ltv.add_paragraph()
        p_h.text = f"{head}: {pr}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_ltv.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # SLIDE A7: ELIGIBILITY CHECKLIST
    sa7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa7, w, h)
    add_header(sa7, "Appendix A7: Checklist Evaluasi Mandiri 12 Kriteria PIDI 2026", "ADDITIONAL · Q&A", 
               "Verifikasi kelayakan menyeluruh sebelum submission final dan pitching juri.", "A7")

    add_card(sa7, Inches(0.8), Inches(1.8), Inches(11.73), Inches(5.1))
    tb_chk = sa7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.7))
    tf_chk = tb_chk.text_frame
    tf_chk.word_wrap = True

    checks = [
        ("1. Problem", "Bukti data resmi ancaman siber dan sanksi regulasi POJK 30/2025."),
        ("2. Alignment", "Konsisten memilih Problem Statement 1 (Manajemen Risiko Finansial)."),
        ("3. Solution", "Penjelasan sederhana tanpa jargon: Ukur → Kendalikan → Uji."),
        ("4. Prototype", "Functional Prototype Go WAF (:8080) & NEX-RED (:3004) berjalan nyata."),
        ("5. Technical", "Arsitektur 3-Layer jelas dari UX hingga database ISO 27001."),
        ("6. Testing", "Data benchmarking latensi (<1.2ms) dan akurasi wasit (99.4%)."),
        ("7. Impact", "Metrik evaluasi baseline vs target kuantitatif (MTTR <5 menit)."),
        ("8. Market", "Target pasar 1.400+ BPR & 100+ Fintech tervalidasi regulasi."),
        ("9. Differentiation", "Matriks komparasi jelas membedakan diri dari WAF & scanner lama."),
        ("10. Team", "Pembagian peran dan kepemilikan kode sumber nyata per anggota."),
        ("11. Roadmap", "Rencana pasca-hackathon milestone-based Q1-Q4 konkret."),
        ("12. Transparency", "Membedakan tegas fitur aktif vs batasan eBPF stub & roadmap.")
    ]

    for i, (chk_title, chk_desc) in enumerate(checks):
        p_c = tf_chk.paragraphs[0] if i == 0 else tf_chk.add_paragraph()
        p_c.text = f"✔ [LULUS] {chk_title}: {chk_desc}"
        p_c.font.size = Pt(9.5)
        p_c.font.color.rgb = NAVY_PRIMARY
        if i > 0: p_c.space_before = Pt(2)

    # SLIDE A8: CLOSING
    sa8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa8, w, h)

    add_card(sa8, Inches(2.0), Inches(1.5), Inches(9.33), Inches(4.5), bg_color=CARD_BG, border_color=EMERALD_ACCENT)
    tb_cl = sa8.shapes.add_textbox(Inches(2.3), Inches(1.8), Inches(8.73), Inches(3.9))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True

    p = tf_cl.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "NEXUS CYBER"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    p_sub = tf_cl.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "Solusi Manajemen Risiko Siber Kanal Digital Terpercaya untuk Ekosistem Keuangan Indonesia"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = EMERALD_ACCENT
    p_sub.space_before = Pt(4)

    p_div = tf_cl.add_paragraph()
    p_div.alignment = PP_ALIGN.CENTER
    p_div.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
    p_div.font.size = Pt(10)
    p_div.font.color.rgb = CARD_BORDER
    p_div.space_before = Pt(10)

    p_contact = tf_cl.add_paragraph()
    p_contact.alignment = PP_ALIGN.CENTER
    p_contact.text = "SIAP BERKOLABORASI & MEMULAI PROGRAM PILOT BERSAMA INSTITUSI ANDA"
    p_contact.font.size = Pt(12)
    p_contact.font.bold = True
    p_contact.font.color.rgb = NAVY_PRIMARY
    p_contact.space_before = Pt(10)

    p_info = tf_cl.add_paragraph()
    p_info.alignment = PP_ALIGN.CENTER
    p_info.text = "Repository: github.com/Thbetyfu/NEXUS-CYBER-FASE3 · Web Portal: nexus-gaas-web (:3003)\nKontak Tim: thoriq@nexus-cyber.test / WhatsApp: +62 895-6033-58692"
    p_info.font.size = Pt(11)
    p_info.font.color.rgb = TEXT_SECONDARY
    p_info.space_before = Pt(8)

def main():
    source_pptx = r"d:\NEXUS\Lomba\power point\terbaru\PITCHDECK_PIDIXHACKATHON_S0280_TELULANG.pptx"
    backup_pptx = r"d:\NEXUS\Lomba\power point\terbaru\PITCHDECK_PIDIXHACKATHON_S0280_TELULANG_ORIGINAL_BACKUP.pptx"
    output_merged_pptx = r"d:\NEXUS\Lomba\power point\terbaru\PITCHDECK_PIDIXHACKATHON_S0280_TELULANG_COMBINED.pptx"

    # 1. Pastikan backup file asli ada
    if not os.path.exists(backup_pptx):
        shutil.copy2(source_pptx, backup_pptx)
        print(f"[BACKUP] Backup asli tersimpan di: {backup_pptx}")

    # 2. Muat presentasi visual pengguna dari backup asli agar bersih
    prs = Presentation(backup_pptx)
    w = prs.slide_width
    h = prs.slide_height

    # Populate Slide 1: PIDI Mandatory Cover & Team Profiling
    print("[1/11] Mengisi Slide 1: Cover Resmi PIDI, Profil Tim & Problem Statement (PS-1)...")
    populate_slide_1_cover(prs.slides[0], w, h)

    # Populate Slide 5: Checklist Matrix (Contreng ✔ / Silang ✖) & Competitor Advantage
    print("[2/11] Mengisi Slide 5: Matriks Komparasi Contreng & Silang...")
    populate_slide_5(prs.slides[4], w, h)

    # Insert Slide 5B: 2x2 Market Positioning Matrix (Competitor & Positioning)
    print("[3/11] Menambahkan Slide 5B: 2x2 Market Positioning Matrix...")
    s5b = prs.slides.add_slide(prs.slide_layouts[6])
    populate_slide_positioning(s5b, w, h)
    prs.slides._sldIdLst.insert(5, prs.slides._sldIdLst[-1])

    # Populate Slide 6: Market Sizing (TAM, SAM, SOM) & GaaS Sentinel (now at index 6)
    print("[3/10] Mengisi Slide 6: Market Size (TAM, SAM, SOM) & Kategori GaaS...")
    populate_slide_6(prs.slides[6], w, h)

    # Populate Slide 7: Cara Kerja Solusi (Business-Friendly 3-Step Flow) (now at index 7)
    print("[4/10] Mengisi Slide 7: Alur Bisnis 3 Langkah Sederhana...")
    populate_slide_7(prs.slides[7], w, h)

    # Populate Slide 8: Go-To-Market (GTM) Strategy (now at index 8)
    print("[5/10] Mengisi Slide 8: Strategi Go-To-Market (GTM)...")
    populate_slide_8(prs.slides[8], w, h)

    # Populate Slide 9: Social Impact & Financial Inclusion (ESG) (now at index 9)
    print("[6/10] Mengisi Slide 9: Dampak Sosial & Inklusi Keuangan (ESG)...")
    populate_slide_9(prs.slides[9], w, h)

    # Populate Slide 10A: Market & Regulatory Validation (now at index 10)
    print("[7/10] Mengisi Slide 10A: Validasi Masalah, Regulasi & Bukti Teknis...")
    populate_slide_10a(prs.slides[10], w, h)

    # Insert Slide 10B: Real Adopter Traction (3 Instansi + 2 UMKM) & WTP (insert at index 11)
    print("[8/10] Menambahkan Slide 10B: Traksi 3 Instansi + 2 UMKM & WTP...")
    s10b = prs.slides.add_slide(prs.slide_layouts[6])
    populate_slide_10b(s10b, w, h)
    prs.slides._sldIdLst.insert(11, prs.slides._sldIdLst[-1])

    # Populate Slide 11: 3-Segment Business Model & 3-Year Projections (now at index 12)
    print("[9/10] Mengisi Slide 11: Model Bisnis 3 Segmen & Proyeksi Finansial...")
    populate_slide_11(prs.slides[12], w, h)

    # Populate Slide 13: Roadmap, ROI Analysis & Funding Ask (Rp 500 Juta) (now at index 14)
    print("[10/10] Mengisi Slide 13: Roadmap, ROI/BEP & Funding Ask Rp 500 Juta...")
    populate_slide_13(prs.slides[14], w, h)

    # Append 8 Appendix Q&A Slides
    print("[APPENDIX] Menambahkan 8 Slide Lampiran Q&A...")
    add_appendix_slides(prs)

    targets = [
        r"d:\NEXUS\Lomba\power point\terbaru\PITCHDECK_PIDIXHACKATHON_S0280_TELULANG_V2.pptx",
        output_merged_pptx,
        source_pptx
    ]

    for target in targets:
        try:
            prs.save(target)
            print(f"[SUKSES] File tersimpan sempurna di: {target}")
        except PermissionError:
            print(f"[PERINGATAN] File {target} sedang dibuka di PowerPoint (terkunci).")

    print(f"Total Slide: {len(prs.slides)} (16 Slide Inti + 8 Slide Lampiran Q&A = {len(prs.slides)} Slides)")

if __name__ == "__main__":
    main()
