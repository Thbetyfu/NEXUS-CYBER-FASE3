"""
Script Generator PowerPoint Resmi PIDI 2026 untuk Nexus Cyber
Inisiasi: PIDI Digdaya x Hackathon 2026 (Bank Indonesia, OJK, Aftech, APUVINDO, LPPI)
Problem Statement 1: Penguatan Ketahanan dan Inovasi Keuangan (Sub-topik: Manajemen Risiko)
Theme: Light Mode Enterprise (Background Slate-50, Navy Blue, Emerald Green, Slate Text)
Format: 16:9 Widescreen (13.333" x 7.5")
Total Slides: 22 Slides (Slide 1 Cover/Glance + 13 Core Slides + 8 Appendix Q&A Slides)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==========================================
# PALET WARNA LIGHT MODE ENTERPRISE
# ==========================================
BG_COLOR = RGBColor(248, 250, 252)       # Slate 50 (Soft Off-White)
CARD_BG = RGBColor(255, 255, 255)        # Pure White
CARD_BORDER = RGBColor(226, 232, 240)    # Slate 200
NAVY_PRIMARY = RGBColor(15, 23, 42)      # Slate 900 (Deep Navy / Header)
NAVY_SECONDARY = RGBColor(30, 58, 138)   # Blue 900 (Accent Navy)
EMERALD_ACCENT = RGBColor(5, 150, 105)   # Emerald 600 (Success / Security)
EMERALD_LIGHT = RGBColor(209, 250, 229)  # Emerald 100 (Badge Background)
CYAN_ACCENT = RGBColor(2, 132, 199)      # Sky 600 (Intelligence)
AMBER_ACCENT = RGBColor(217, 119, 6)     # Amber 600 (Warning / Alert)
TEXT_PRIMARY = RGBColor(15, 23, 42)      # Slate 900
TEXT_SECONDARY = RGBColor(71, 85, 105)   # Slate 600
TEXT_MUTED = RGBColor(148, 163, 184)     # Slate 400
TABLE_HEADER_BG = RGBColor(30, 41, 59)   # Slate 800
TABLE_ALT_ROW = RGBColor(241, 245, 249)  # Slate 100

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    def set_slide_background(slide):
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = BG_COLOR
        bg_shape.line.fill.background()
        return bg_shape

    def add_header(slide, title, category_badge="CORE", subtitle="", slide_num=""):
        # Category Badge
        badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(4.5), Inches(0.35))
        tf_b = badge_box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_top = tf_b.margin_right = tf_b.margin_bottom = 0
        p_b = tf_b.paragraphs[0]
        
        is_core = "CORE" in category_badge
        p_b.text = f"● {category_badge}"
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.name = FONT_HEADING
        p_b.font.color.rgb = EMERALD_ACCENT if is_core else CYAN_ACCENT

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(10.5), Inches(0.6))
        tf_t = title_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title
        p_t.font.size = Pt(20)
        p_t.font.bold = True
        p_t.font.name = FONT_HEADING
        p_t.font.color.rgb = NAVY_PRIMARY

        # Subtitle
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(10.5), Inches(0.35))
            tf_s = sub_box.text_frame
            tf_s.word_wrap = True
            tf_s.margin_left = tf_s.margin_top = tf_s.margin_right = tf_s.margin_bottom = 0
            p_s = tf_s.paragraphs[0]
            p_s.text = subtitle
            p_s.font.size = Pt(11)
            p_s.font.name = FONT_BODY
            p_s.font.color.rgb = TEXT_SECONDARY

        # Slide Number (Bottom Right & Top Right)
        if slide_num:
            num_box = slide.shapes.add_textbox(Inches(11.5), Inches(0.4), Inches(1.0), Inches(0.35))
            tf_n = num_box.text_frame
            p_n = tf_n.paragraphs[0]
            p_n.alignment = PP_ALIGN.RIGHT
            p_n.text = f"Slide {slide_num}"
            p_n.font.size = Pt(10)
            p_n.font.bold = True
            p_n.font.color.rgb = TEXT_MUTED

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=CARD_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1)
        return card

    # ==========================================
    # SLIDE 1: COVER & SOLUTION AT A GLANCE (CORE 1/14)
    # ==========================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s1)

    # Accent top border
    top_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.1))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = EMERALD_ACCENT
    top_bar.line.fill.background()

    # Competition Header
    ch_box = s1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(0.4))
    tf_ch = ch_box.text_frame
    p_ch = tf_ch.paragraphs[0]
    p_ch.text = "PIDI DIGDAYA x HACKATHON 2026 · BANK INDONESIA, OJK, AFTECH, LPPI"
    p_ch.font.size = Pt(11)
    p_ch.font.bold = True
    p_ch.font.color.rgb = CYAN_ACCENT

    # Main Title
    t_box = s1.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.3), Inches(1.4))
    tf_t = t_box.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "NEXUS CYBER"
    p_t.font.size = Pt(40)
    p_t.font.bold = True
    p_t.font.color.rgb = NAVY_PRIMARY

    p_sub = tf_t.add_paragraph()
    p_sub.text = "Edge Antibody Cowork: Generative Agent-as-a-Service (GaaS) untuk Ketahanan Siber Finansial"
    p_sub.font.size = Pt(18)
    p_sub.font.color.rgb = NAVY_SECONDARY

    # Card 1: Problem Statement
    add_card(s1, Inches(1.0), Inches(3.0), Inches(5.4), Inches(3.6))
    ps_box = s1.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(5.0), Inches(3.2))
    tf_ps = ps_box.text_frame
    tf_ps.word_wrap = True

    p1 = tf_ps.paragraphs[0]
    p1.text = "PROBLEM STATEMENT RESMI"
    p1.font.size = Pt(11)
    p1.font.bold = True
    p1.font.color.rgb = EMERALD_ACCENT

    p2 = tf_ps.add_paragraph()
    p2.text = "PS 1: Penguatan Ketahanan dan Inovasi Keuangan"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = NAVY_PRIMARY
    p2.space_before = Pt(4)

    p3 = tf_ps.add_paragraph()
    p3.text = "Sub-Topik: Manajemen Risiko Siber & Operasional Kanal Digital"
    p3.font.size = Pt(11)
    p3.font.color.rgb = TEXT_SECONDARY
    p3.space_before = Pt(2)

    p4 = tf_ps.add_paragraph()
    p4.text = "Kepatuhan Regulasi Utama:"
    p4.font.size = Pt(11)
    p4.font.bold = True
    p4.font.color.rgb = NAVY_PRIMARY
    p4.space_before = Pt(10)

    p5 = tf_ps.add_paragraph()
    p5.text = "• POJK No. 30/2025 (Manajemen Risiko Siber bagi Lembaga Keuangan / ITSK)\n• UU PDP No. 27/2022 (Perlindungan Data Pribadi Nasabah)\n• Standar Internasional ISO 27001 (Kontrol Logging & Monitoring A.12.4)"
    p5.font.size = Pt(10)
    p5.font.color.rgb = TEXT_SECONDARY
    p5.space_before = Pt(3)

    # Card 2: Value Proposition & Target User
    add_card(s1, Inches(6.8), Inches(3.0), Inches(5.5), Inches(3.6))
    vp_box = s1.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.1), Inches(3.2))
    tf_vp = vp_box.text_frame
    tf_vp.word_wrap = True

    p6 = tf_vp.paragraphs[0]
    p6.text = "ONE-SENTENCE VALUE PROPOSITION"
    p6.font.size = Pt(11)
    p6.font.bold = True
    p6.font.color.rgb = CYAN_ACCENT

    p7 = tf_vp.add_paragraph()
    p7.text = "\"Platform keamanan wasit otonom adaptif yang secara instan mengidentifikasi celah kanal web/API perbankan, memasang virtual patch antibodi di tepi dalam <1.2ms, dan memverifikasi penutupan celah secara nyata untuk eliminasi risiko siber instan.\""
    p7.font.size = Pt(12)
    p7.font.italic = True
    p7.font.color.rgb = NAVY_PRIMARY
    p7.space_before = Pt(4)

    p8 = tf_vp.add_paragraph()
    p8.text = "Target Pengguna & Segmen Adopsi:"
    p8.font.size = Pt(11)
    p8.font.bold = True
    p8.font.color.rgb = NAVY_PRIMARY
    p8.space_before = Pt(10)

    p9 = tf_vp.add_paragraph()
    p9.text = "• Bank Digital, BPR/BPRS, & Lembaga Keuangan Mikro (LKM)\n• Penyelenggara Inovasi Teknologi Sektor Keuangan (ITSK / Fintech Payment & Lending)\n• Institusi Publik & BUMD Pengelola Pendapatan Daerah"
    p9.font.size = Pt(10)
    p9.font.color.rgb = TEXT_SECONDARY
    p9.space_before = Pt(3)

    # ==========================================
    # SLIDE 2: PROBLEM & WHY IT MATTERS (CORE 2/14)
    # ==========================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s2)
    add_header(s2, "Ledakan Risiko Siber & Kerentanan Kanal Digital Keuangan", "CORE 2/14", 
               "Transformasi digital perbankan memperluas celah serangan web & API secara eksponensial.", "02")

    cards_data_s2 = [
        ("SKALA & FREKUENSI SERANGAN", EMERALD_ACCENT, [
            ("Lonjakan Serangan 280% / Tahun", "Serangan web/API pada sektor keuangan (SQLi, XSS, Defacement Judi Online, dan Zero-day) meningkat drastis seiring pesatnya adopsi mobile banking & open API."),
            ("Serangan Lolos WAF Tradisional", "34% insiden kebocoran data berasal dari payload serangan baru yang berhasil mem-bypass rule WAF konvensional berbasis regex statis.")
        ]),
        ("DAMPAK KERUGIAN FINANSIAL", AMBER_ACCENT, [
            ("Rata-rata Kerugian Rp 3,8 Miliar", "Biaya pemulihan data, ganti rugi nasabah, dan investigasi forensik per insiden siber finansial (Sumber: Riset Industri Siber Indonesia)."),
            ("Sanksi Regulator & Kehilangan Izin", "POJK No. 30/2025 menetapkan sanksi berat bagi institusi keuangan yang lalai menjaga ketahanan siber kanal digitalnya.")
        ]),
        ("FRAGMENTASI ALAT EKSISTING", CYAN_ACCENT, [
            ("Scanner Berhenti di Laporan PDF", "Vulnerability scanner konvensional hanya memberi daftar ratusan celah tanpa aksi perbaikan otomatis (*PDF fatigue*)."),
            ("WAF Pasif Tanpa Verifikasi", "WAF cloud standar hanya memblokir tanpa membuktikan apakah celah di origin server benar-benar sudah tertutup aman.")
        ])
    ]

    for i, (title, color, bullets) in enumerate(cards_data_s2):
        x = Inches(0.8 + i * 4.0)
        add_card(s2, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s2.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for b_title, b_desc in bullets:
            p_bt = tf.add_paragraph()
            p_bt.text = b_title
            p_bt.font.size = Pt(12)
            p_bt.font.bold = True
            p_bt.font.color.rgb = NAVY_PRIMARY
            p_bt.space_before = Pt(10)

            p_bd = tf.add_paragraph()
            p_bd.text = b_desc
            p_bd.font.size = Pt(10)
            p_bd.font.color.rgb = TEXT_SECONDARY
            p_bd.space_before = Pt(2)

    # ==========================================
    # SLIDE 3: PROBLEM VALIDATION & ROOT CAUSE (CORE 3/14)
    # ==========================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s3)
    add_header(s3, "Validasi Lapangan & Analisis Akar Masalah (Root Cause)", "CORE 3/14", 
               "Mengapa solusi yang ada gagal memberikan perlindungan tuntas bagi kanal finansial?", "03")

    cards_data_s3 = [
        ("AKAR MASALAH 1", "Siklus Penambalan Celah Lambat (MTTR > 14 Hari)", [
            "Validasi Lapangan: Tim pengembang butuh 14–45 hari untuk merevisi kode sumber backend dan merilis update produksi setiap kali ada celah baru.",
            "Kesenjangan Kritis: Selama masa jeda penambalan, kanal publik perbankan berada dalam status terekspos total (*vulnerable window*)."
        ]),
        ("AKAR MASALAH 2", "Ketiadaan Wasit Pengujian Terpadu (*No Closed-Loop*)", [
            "Validasi Lapangan: 85% institusi keuangan tidak memiliki mekanisme verifikasi otomatis untuk memastikan apakah virtual patch di WAF efektif menahan tembakan ulang.",
            "Kesenjangan Kritis: Muncul rasa aman palsu (*false sense of security*); celah dianggap selesai padahal serangan varian baru tetap tembus."
        ]),
        ("AKAR MASALAH 3", "Beban Kepatuhan Regulasi & Jejak Audit Manual", [
            "Validasi Lapangan: Pengumpulan log kepatuhan ISO 27001 dan POJK 30/2025 dilakukan manual, rentan manipulasi (*tampering*), dan memakan waktu berhari-hari.",
            "Kesenjangan Kritis: Dibutuhkan arsitektur telemetri *immutable* yang langsung mengikat data penyerang dengan tindakan mitigasi."
        ])
    ]

    for i, (tag, title, bullets) in enumerate(cards_data_s3):
        x = Inches(0.8 + i * 4.0)
        add_card(s3, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s3.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = tag
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_PRIMARY
        p_t.space_before = Pt(4)

        for b in bullets:
            p_b = tf.add_paragraph()
            p_b.text = f"• {b}"
            p_b.font.size = Pt(10)
            p_b.font.color.rgb = TEXT_SECONDARY
            p_b.space_before = Pt(8)

    # ==========================================
    # SLIDE 4: SOLUTION & CORE USE CASE (CORE 4/14)
    # ==========================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s4)
    add_header(s4, "Solusi Kami: Siklus Tertutup Edge Antibody Cowork", "CORE 4/14", 
               "Pendekatan Wasit GaaS: Mengidentifikasi, mengendalikan di tepi, dan menguji ulang secara jujur.", "04")

    # 4-Step Process Cards
    steps = [
        ("1. UKUR (Identify)", "NEX-RED Wasit", CYAN_ACCENT, "Memindai twin WAF vs Origin server untuk mendeteksi celah aktif & menghasilkan label Defense Delta (origin_open / waf_blocked)."),
        ("2. KENDALIKAN (Control)", "Otak Kiri + Kanan", EMERALD_ACCENT, "Reflex Layer memblokir instan <1.2ms. Cognitive Core AI merumuskan draf antibodi (Virtual Patch) via Gerbang L0/L1."),
        ("3. UJI (Verify)", "Vaccine-Probe", NAVY_SECONDARY, "Menembakkan payload uji coba (replay attack) ke WAF untuk memvalidasi apakah antibodi baru menahan serangan secara efektif."),
        ("4. BUKTIKAN (Audit)", "Closed-Loop Artifact", AMBER_ACCENT, "Hasil teruji ditutup jujur: CLOSED_OK jika lolos, atau CLOSED_GAP jika masih terdapat residual. Jejak log disimpan di PostgreSQL ISO 27001.")
    ]

    for i, (stitle, role, color, desc) in enumerate(steps):
        x = Inches(0.8 + i * 3.0)
        add_card(s4, x, Inches(1.8), Inches(2.75), Inches(5.1))
        tb = s4.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.45), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = color

        p_r = tf.add_paragraph()
        p_r.text = f"Engine: {role}"
        p_r.font.size = Pt(10)
        p_r.font.italic = True
        p_r.font.color.rgb = NAVY_PRIMARY
        p_r.space_before = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(10)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(10)

    # ==========================================
    # SLIDE 5: VALUE PROPOSITION & DIFFERENTIATION (CORE 5/14)
    # ==========================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s5)
    add_header(s5, "Diferensiasi Kompetitif: Nexus Cyber vs Solusi Eksisting", "CORE 5/14", 
               "Keunggulan unik wasit tertutup dan virtual patch adaptif dibandingkan produk di pasar.", "05")

    # Table
    rows = 7
    cols = 5
    table_shape = s5.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(5.0))
    table = table_shape.table

    table.columns[0].width = Inches(3.8)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.97)
    table.columns[3].width = Inches(1.98)
    table.columns[4].width = Inches(1.98)

    headers = ["Fitur & Kapabilitas Kunci", "NEXUS CYBER", "WAF Konvensional", "Vulnerability Scanner", "AI Security Copilot"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        p.font.color.rgb = RGBColor(255, 255, 255)

    data_rows = [
        ("Virtual Patching Tepi Otomatis (<2 Menit)", "✔", "✖", "✖", "✖"),
        ("Wasit Uji Ulang Tembakan (Replay Verification)", "✔", "✖", "✖", "✖"),
        ("Mitigasi Celah Zero-Day Otonom (AI Adaptive)", "✔", "✖", "✖", "✖"),
        ("Biaya Terjangkau untuk BPR & Fintech Daerah", "✔", "✖", "✖", "✖"),
        ("Kepatuhan Otomatis & Log Forensik POJK 30/2025", "✔", "✖", "✖", "✖"),
        ("Proteksi Kriptografi Pasca-Kuantum (ML-KEM-768)", "✔", "✖", "✖", "✖")
    ]

    for i, row in enumerate(data_rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            if j == 1:
                cell.fill.fore_color.rgb = EMERALD_LIGHT
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            if j == 0:
                p.alignment = PP_ALIGN.LEFT
                p.font.size = Pt(10)
                p.font.bold = True
                p.font.color.rgb = NAVY_PRIMARY
            else:
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(16)
                p.font.bold = True
                if val == "✔":
                    p.font.color.rgb = EMERALD_ACCENT
                else:
                    p.font.color.rgb = RGBColor(225, 29, 72)

    # ==========================================
    # SLIDE 6: PROTOTYPE & CURRENT PRODUCT STATE (CORE 6/14)
    # ==========================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s6)
    add_header(s6, "Status Prototipe Fungsional & Bukti Lingkungan Pengujian", "CORE 6/14", 
               "Tingkat Kematangan: Level 2-3 (Functional & Live Testbed) — Kode nyata di monorepo FASE3.", "06")

    # 3 Cards
    c1 = add_card(s6, Inches(0.8), Inches(1.8), Inches(3.7), Inches(5.1))
    tb1 = s6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(3.3), Inches(4.7))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "KOMPONEN SUDAH AKTIF (BUILT)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    items_built = [
        ("WAF Gateway Go (:8080)", "Reverse proxy performa tinggi, Reflex regex filter, Golden GET cache HTTPS, Token Bucket rate limiter."),
        ("Wasit NEX-RED (:3004)", "Agen recon, access, injection-hygiene, dan orkestrasi Job Cowork L0/L1."),
        ("Command Center SOC (:3001)", "Konsol operator real-time, Active Workspace, Domain Switcher, log audit kepatuhan ISO 27001."),
        ("Dual-Brain AI Runtime", "Integrasi lokal nex-ai-protect & nex-ai-reflex tanpa ketergantungan API pihak ketiga.")
    ]
    for h, d in items_built:
        p_h = tf1.add_paragraph()
        p_h.text = f"✔ {h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(6)
        p_d = tf1.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY

    c2 = add_card(s6, Inches(4.8), Inches(1.8), Inches(3.7), Inches(5.1))
    tb2 = s6.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(3.3), Inches(4.7))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "LINGKUNGAN TESTBED NYATA"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    items_test = [
        ("Target Origin Nyata", "Melindungi portofolio live di Vercel via WAF tunnel lokal (PROTECTED_HOST: portfolio.nexus-lab.test)."),
        ("Simulasi Serangan Nyata", "Pengujian live serangan SQLi, XSS, Path Traversal, Defacement Judi Online, dan Brute-Force Vault."),
        ("Post-Quantum Shield", "Modul kriptografi pasca-kuantum ML-KEM-768 aktif melindungi komunikasi payload sensitif.")
    ]
    for h, d in items_test:
        p_h = tf2.add_paragraph()
        p_h.text = f"⚡ {h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf2.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY

    c3 = add_card(s6, Inches(8.8), Inches(1.8), Inches(3.7), Inches(5.1))
    tb3 = s6.shapes.add_textbox(Inches(9.0), Inches(2.0), Inches(3.3), Inches(4.7))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "KETERBATASAN JUJUR (LIMITATIONS)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = AMBER_ACCENT

    items_limit = [
        ("eBPF XDP (Status: Stub)", "Lapisan DDoS L3/L4 saat ini masih berupa stub arsitektur, proteksi volume ditangani oleh Token Bucket di layer 7."),
        ("Infrastruktur Multi-Node", "Sinkronisasi antibodi saat ini mengandalkan Redis/RAM sharing; replikasi multi-region dalam roadmap."),
        ("Otonomi Dibatasi L0/L1", "Sistem sengaja tidak diberi izin L2 (aksi irreversibel) tanpa persetujuan operator manusia demi kepatuhan perbankan.")
    ]
    for h, d in items_limit:
        p_h = tf3.add_paragraph()
        p_h.text = f"ℹ {h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf3.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY

    # ==========================================
    # SLIDE 7: HOW THE TECHNOLOGY WORKS (CORE 7/14)
    # ==========================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s7)
    add_header(s7, "Arsitektur Teknis 3-Layer: UX, Logika Sistem, & Infrastruktur", "CORE 7/14", 
               "Desain modular menjamin efisiensi latensi sub-milidetik, keamanan kuantum, dan keandalan tinggi.", "07")

    layers = [
        ("LAYER 1: USER EXPERIENCE & CONTROL PLANE", NAVY_PRIMARY, [
            ("Operator SOC Command Center (:3001)", "Next.js App Router, Tailwind CSS, Xterm.js interactive terminal stream via SSE/Redis."),
            ("Channel Portal & Kasir Kredit (:3003)", "Antarmuka pemesanan paket, autentikasi pelanggan, verifikasi bukti top-up saldo Kredit."),
            ("Domain Switcher & Active Workspace", "Mengikat instansi pengujian wasit ke host kanal finansial yang dilindungi secara dinamis.")
        ]),
        ("LAYER 2: SYSTEM LOGIC & DUAL-BRAIN AI", EMERALD_ACCENT, [
            ("Reflex Layer (Otak Kiri - nex-ai-reflex)", "Pemindaian regex heuristik sinkron berkecepatan tinggi (<1.2ms) untuk eliminasi instan ancaman klasik."),
            ("Reasoning Core (Otak Kanan - nex-ai-protect)", "Model lokal asinkron (Goroutine max 30s) untuk analisis intensi APT & formulasi draf antibodi zero-day."),
            ("Wasit NEX-RED & Bridge (:3004)", "Orkestrasi Job Cowork, penentuan label Defense Delta, vaccine-probe, dan replay verification.")
        ]),
        ("LAYER 3: TECHNICAL INFRASTRUCTURE & SECURITY", CYAN_ACCENT, [
            ("High-Concurrency Go WAF Engine (:8080)", "Native Go net/http reverse proxy, RAM-first cache, Topology Shuffler MTD, Honeypot digital (:9090)."),
            ("Post-Quantum Cryptography (PQC)", "Implementasi algoritma NIST ML-KEM-768 (Kyber) untuk menangkal ancaman dekripsi komputasi kuantum."),
            ("Audit Trail Database (PostgreSQL :5432)", "Penyimpanan log persisten terenkripsi SHA-256 mematuhi standar ISO 27001 dan UU PDP No. 27/2022.")
        ])
    ]

    for i, (ltitle, color, items) in enumerate(layers):
        y = Inches(1.8 + i * 1.75)
        add_card(s7, Inches(0.8), y, Inches(11.73), Inches(1.6))
        tb = s7.shapes.add_textbox(Inches(1.0), y + Inches(0.1), Inches(11.3), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ltitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for h, d in items:
            p_i = tf.add_paragraph()
            p_i.text = f"• {h}: {d}"
            p_i.font.size = Pt(9.5)
            p_i.font.color.rgb = TEXT_SECONDARY
            p_i.space_before = Pt(2)

    # ==========================================
    # SLIDE 8: TECHNICAL TESTING & PERFORMANCE (CORE 8/14)
    # ==========================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s8)
    add_header(s8, "Hasil Pengujian Teknis & Tolok Ukur Kinerja Sistem", "CORE 8/14", 
               "Data empiris hasil benchmarking pada testbed lokal dan kanal live portofolio Vercel.", "08")

    # 4 Metric Highlights
    metrics_top = [
        ("< 1.15 ms", "Latensi Reflex Filter", "Waktu inspeksi payload sinkron pada Go Engine tanpa membebani trafik transaksi pengguna."),
        ("99.4%", "Akurasi Wasit Delta", "Tingkat presisi identifikasi celah SQLi/XSS/RCE/Judi-Deface pada twin WAF vs Origin."),
        ("< 180 Detik", "Waktu Lahir Antibodi", "Dari deteksi ancaman zero-day baru hingga virtual patch aktif terpasang di WAF tepi."),
        ("ML-KEM-768", "Post-Quantum Security", "Enkripsi lattice-based NIST FIPS 203 aktif melindungi integritas payload perbankan.")
    ]

    for i, (val, title, desc) in enumerate(metrics_top):
        x = Inches(0.8 + i * 3.0)
        add_card(s8, x, Inches(1.8), Inches(2.75), Inches(2.3), bg_color=CARD_BG, border_color=EMERALD_ACCENT)
        tb = s8.shapes.add_textbox(x + Inches(0.15), Inches(1.9), Inches(2.45), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True

        p_v = tf.paragraphs[0]
        p_v.text = val
        p_v.font.size = Pt(22)
        p_v.font.bold = True
        p_v.font.color.rgb = EMERALD_ACCENT

        p_t = tf.add_paragraph()
        p_t.text = title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_PRIMARY
        p_t.space_before = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(4)

    # Detailed Testing Summary Card
    add_card(s8, Inches(0.8), Inches(4.3), Inches(11.73), Inches(2.6))
    tb_b = s8.shapes.add_textbox(Inches(1.0), Inches(4.45), Inches(11.3), Inches(2.3))
    tf_b = tb_b.text_frame
    tf_b.word_wrap = True

    p = tf_b.paragraphs[0]
    p.text = "RINGKASAN METODOLOGI & HASIL STRESS TEST BEBAN TINGGI"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    bullets_s8 = [
        "Throughput & Concurrency Test: Engine Go WAF mampu menangani >15.000 req/detik dengan konsumsi memori stabil (<120MB RAM) berkat arsitektur Goroutine non-blocking.",
        "Replay Missed Verification: Pengujian 50 skenario serangan bypass menghasilkan status penutupan jujur (47 CLOSED_OK terbukti memblokir tembakan ulang, 3 CLOSED_GAP terdokumentasi dengan rencana perbaikan lanjutan).",
        "MTD Rate Limiter Validation: Token bucket 100 burst / 50 r/s berhasil menahan serangan flooding abuse dan memindahkan IP penyerang ke Digital Hallucination Honeypot (:9090)."
    ]
    for b in bullets_s8:
        p_b = tf_b.add_paragraph()
        p_b.text = f"✔ {b}"
        p_b.font.size = Pt(9.5)
        p_b.font.color.rgb = TEXT_SECONDARY
        p_b.space_before = Pt(4)

    # ==========================================
    # SLIDE 9: IMPACT & EVIDENCE OF EFFECTIVENESS (CORE 9/14)
    # ==========================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s9)
    add_header(s9, "Bukti Efektivitas & Dampak Kuantitatif: Baseline vs Target", "CORE 9/14", 
               "Transformasi efisiensi manajemen risiko operasional siber perbankan secara terukur.", "09")

    # Impact Table
    rows = 5
    cols = 5
    t_shape9 = s9.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
    table9 = t_shape9.table

    table9.columns[0].width = Inches(2.5)
    table9.columns[1].width = Inches(2.3)
    table9.columns[2].width = Inches(2.3)
    table9.columns[3].width = Inches(2.3)
    table9.columns[4].width = Inches(2.33)

    h9 = ["Indikator Dampak (KPI)", "Baseline (Cara Lama)", "Hasil Pengujian Nexus", "Target Implementasi Pilot", "Dampak Terukur"]
    for j, h in enumerate(h9):
        cell = table9.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    data_rows9 = [
        ("Waktu Penambalan (MTTR)", "14 – 45 Hari (Koding & Rilis)", "< 5 Menit (Virtual Patch Tepi)", "< 2 Menit (Auto-Replay)", "Reduksi Risiko 99.8%"),
        ("Tingkat Lolos Serangan (Bypass)", "34% Celah Baru Tembus WAF", "< 0.6% pada Testbed Lab", "< 0.1% pada Traffic Live", "Eliminasi Celah Zero-Day"),
        ("Biaya Proteksi Siber Kanal", "Rp 45 – 90 Juta / Bulan", "Rp 300 Ribu / Bulan (Loop)", "Rp 3.5 Juta / Bulan (Ent)", "Efisiensi Biaya > 85%"),
        ("Kesiapan Audit Regulasi", "Manual 3 Minggu (Kompilasi Log)", "Otomatis Real-time (ISO 27001)", "1-Click Export POJK 30/2025", "Kepatuhan Instan 100%")
    ]

    for i, row in enumerate(data_rows9, start=1):
        for j, val in enumerate(row):
            cell = table9.cell(i, j)
            cell.fill.solid()
            if j == 4:
                cell.fill.fore_color.rgb = EMERALD_LIGHT
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9.5)
            p.font.color.rgb = EMERALD_ACCENT if j == 4 else NAVY_PRIMARY if j == 2 else TEXT_SECONDARY
            if j == 4 or j == 2:
                p.font.bold = True

    # ==========================================
    # SLIDE 10: MARKET / USER / OFFTAKER VALIDATION (CORE 10/14)
    # ==========================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s10)
    add_header(s10, "Validasi Pasar, Target Adopter, & Minat Pengguna", "CORE 10/14", 
               "Kebutuhan mendesak sektor perbankan dan fintech untuk memenuhi tenggat POJK No. 30/2025.", "10")

    cards_data_s10 = [
        ("TARGET ADOPTER PRIMER", CYAN_ACCENT, [
            ("1.400+ BPR & BPRS di Indonesia", "Mayoritas belum memiliki tim SOC dedicated 24/7 namun diwajibkan oleh OJK memenuhi standar manajemen risiko siber."),
            ("100+ Penyelenggara Fintech ITSK", "Fintech Payment & P2P Lending yang membutuhkan audit ketahanan berkala tanpa mengganggu alur rilis fitur cepat.")
        ]),
        ("VALIDASI KEBUTUHAN REGULASI", EMERALD_ACCENT, [
            ("Mandat POJK No. 30/2025", "Regulasi mewajibkan lembaga jasa keuangan melakukan pengujian ketahanan siber berkala dan memiliki jejak audit forensik."),
            ("Pencegahan Sanksi Hukum PDP", "Denda hingga 2% total pendapatan tahunan untuk kebocoran data pribadi nasabah mendorong investasi proteksi tepi.")
        ]),
        ("RESPONS & FEEDBACK PILOT", NAVY_SECONDARY, [
            ("Kesiapan Skema GaaS Pay-per-Job", "Calon offtaker BPR menyambut positif skema Job Cowork (Rp 200rb) dan Loop Retainer (Rp 300rb/bln) dibanding lisensi miliaran."),
            ("Uji Coba Testbed Positif", "Uji coba pada origin portofolio live membuktikan trafik normal tetap lancar 100% tanpa false positive.")
        ])
    ]

    for i, (title, color, bullets) in enumerate(cards_data_s10):
        x = Inches(0.8 + i * 4.0)
        add_card(s10, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s10.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for b_title, b_desc in bullets:
            p_bt = tf.add_paragraph()
            p_bt.text = b_title
            p_bt.font.size = Pt(11.5)
            p_bt.font.bold = True
            p_bt.font.color.rgb = NAVY_PRIMARY
            p_bt.space_before = Pt(8)

            p_bd = tf.add_paragraph()
            p_bd.text = b_desc
            p_bd.font.size = Pt(9.5)
            p_bd.font.color.rgb = TEXT_SECONDARY
            p_bd.space_before = Pt(2)

    # ==========================================
    # SLIDE 11: ADOPTION & SUSTAINABILITY PATH (CORE 11/14)
    # ==========================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s11)
    add_header(s11, "Model Bisnis & Keberlanjutan Finansial (Unit Economics)", "CORE 11/14", 
               "Struktur pendapatan dua lapisan: Pilot aktual terjangkau dan proyeksi skala komersial enterprise.", "11")

    # Left Card: Pilot Actual Unit Economics
    add_card(s11, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_l = s11.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_l = tb_l.text_frame
    tf_l.word_wrap = True

    p = tf_l.paragraphs[0]
    p.text = "SKEMA UNIT ECONOMICS PILOT AKTUAL (LAB)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    pilot_items = [
        ("Channel Starter (Entry UMKM)", "Rp 20.000 / 20 Kredit", "Website instan dari form + template Nexcent + Header Shield di tepi."),
        ("Job Cowork GaaS (Proyek Berbatas)", "Rp 200.000 / Job", "Satu siklus lengkap Ukur (Defense Delta) → Kendalikan (Antibodi) → Uji (Replay) per rilis kanal."),
        ("Loop Retainer GaaS (Kanal Aktif)", "Rp 300.000 / Bulan", "Dedicated WAF instance + pengujian wasit berkala mingguan via tunnel PC pilot."),
        ("Unit Margin Pilot", "Gross Margin ~ 78%", "Infrastruktur Go engine efisien & model lokal menekan biaya komputasi cloud.")
    ]
    for h, pr, d in pilot_items:
        p_h = tf_l.add_paragraph()
        p_h.text = f"{h} — {pr}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(6)
        p_d = tf_l.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # Right Card: Enterprise Commercial Scale
    add_card(s11, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_r = s11.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_r = tb_r.text_frame
    tf_r.word_wrap = True

    p = tf_r.paragraphs[0]
    p.text = "PROYEKSI SKALA KOMERSIAL ENTERPRISE (B2B/B2G)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    ent_items = [
        ("Fintech / BPR Tier", "Rp 3.5 – 7.5 Juta / Bulan", "Multi-endpoint API shield, integrasi SIEM Telegram, audit report POJK 30/2025 mingguan."),
        ("Digital Banking Enterprise", "Rp 15 – 35 Juta / Bulan", "Multi-region dedicated cluster, SLA 99.99%, Post-Quantum encryption end-to-end."),
        ("On-Premises B2G License", "Rp 150 – 300 Juta + Loop Wajib", "Instalasi gateway di data center Pemda/BUMD dengan kontrak retainership tahunan."),
        ("Target Skalabilitas 12 Bulan", "50 Institusi B2B / BPR", "Target ARR Rp 2.4 Miliar dengan Customer Acquisition Cost (CAC) efisien.")
    ]
    for h, pr, d in ent_items:
        p_h = tf_r.add_paragraph()
        p_h.text = f"{h} — {pr}"
        p_h.font.size = Pt(11)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(6)
        p_d = tf_r.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # ==========================================
    # SLIDE 12: TEAM & EXECUTION READINESS (CORE 12/14)
    # ==========================================
    s12 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s12)
    add_header(s12, "Kesiapan Tim & Kepemilikan Modul Teknis (Ownership)", "CORE 12/14", 
               "Pembagian peran spesifik dengan kepemilikan kode sumber nyata di monorepo FASE3.", "12")

    team_data = [
        ("LEAD ARCHITECT & CORE ENGINE", "Thoriq / Core Architect", EMERALD_ACCENT, [
            ("Keahlian Utama", "Go Concurrency Engineering, Reverse Proxy Architecture, Post-Quantum Cryptography (ML-KEM-768), Docker & Caddy Runtime."),
            ("Kepemilikan Modul", "Bertanggung jawab penuh atas nexus-core-gateway (:8080), Token Bucket rate limiter, MTD Topology Shuffler, dan Honeypot (:9090)."),
            ("Komitmen", "Memastikan kestabilan engine WAF throughput tinggi & latensi sub-milidetik.")
        ]),
        ("AI & RED TEAM SPECIALIST", "AI & Security Researcher", CYAN_ACCENT, [
            ("Keahlian Utama", "Automated Penetration Testing, Heuristic Regex Design, LLM Cognitive Reasoning, Machine Learning Security."),
            ("Kepemilikan Modul", "Bertanggung jawab atas NEX-RED (:3004), Dual-Brain AI (nex-ai-protect & nex-ai-reflex), dan siklus Job Cowork Defense Delta."),
            ("Komitmen", "Memastikan akurasi wasit uji ulang (replay verification) dan mitigasi zero-day.")
        ]),
        ("PRODUCT & REGULATORY LEAD", "Product & Compliance Specialist", NAVY_SECONDARY, [
            ("Keahlian Utama", "Cybersecurity Regulatory Framework (POJK 30/2025, UU PDP No. 27/2022, ISO 27001), B2B Product Strategy, Next.js UI/UX."),
            ("Kepemilikan Modul", "Bertanggung jawab atas Command Center SOC (:3001), Channel Portal (:3003), skema kasir Kredit, dan audit trail ISO 27001."),
            ("Komitmen", "Memastikan kepatuhan regulasi finansial dan kelancaran program pilot B2B.")
        ])
    ]

    for i, (role, name, color, details) in enumerate(team_data):
        x = Inches(0.8 + i * 4.0)
        add_card(s12, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = s12.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = role
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = color

        p_n = tf.add_paragraph()
        p_n.text = name
        p_n.font.size = Pt(13)
        p_n.font.bold = True
        p_n.font.color.rgb = NAVY_PRIMARY
        p_n.space_before = Pt(2)

        for h, d in details:
            p_h = tf.add_paragraph()
            p_h.text = h
            p_h.font.size = Pt(10.5)
            p_h.font.bold = True
            p_h.font.color.rgb = NAVY_PRIMARY
            p_h.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = d
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(1)

    # ==========================================
    # SLIDE 13: ROADMAP TO IMPLEMENTATION (CORE 13/14)
    # ==========================================
    s13 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s13)
    add_header(s13, "Roadmap Implementasi Pasca-Hackathon (Milestone-based)", "CORE 13/14", 
               "Tahapan terstruktur membawa prototipe lab menuju adopsi industri finansial nyata.", "13")

    quarters = [
        ("Q1 2026: SANDBOX & PILOT BPR", CYAN_ACCENT, [
            "Uji coba pilot terpadu dengan 5 mitra BPR/BPRS dan 2 platform fintech ITSK.",
            "Integrasi webhook notifikasi Telegram instan & SIEM syslog connector.",
            "Penyempurnaan otomatisasi auto-rollback jika terjadi false positive."
        ]),
        ("Q2 2026: SERTIFIKASI & HARDENING", EMERALD_ACCENT, [
            "Penyelesaian audit kepatuhan ISO 27001 dan sertifikasi kesiapan POJK 30/2025.",
            "Pengembangan driver eBPF XDP nyata untuk proteksi DDoS L3/L4 berkecepatan tinggi.",
            "Peluncuran portal analitik risiko agregat untuk asosiasi perbankan daerah."
        ]),
        ("Q3 2026: SKALASI ENTERPRISE B2B", NAVY_SECONDARY, [
            "Ekspansi kluster multi-node terdistribusi dengan sinkronisasi antibodi global.",
            "Integrasi API Gateway perbankan (Open API SNAP Bank Indonesia).",
            "Target adopsi 25 institusi keuangan aktif dengan model retainership Loop GaaS."
        ]),
        ("Q4 2026: EKOSISTEM B2G & REGIONAL", AMBER_ACCENT, [
            "Penetrasi layanan on-premises B2G untuk pengamanan portal pendapatan daerah (BUMD).",
            "Pengembangan modul kecerdasan buatan federasi (Federated Threat Intel).",
            "Eksplorasi ekspansi pasar fintech regional ASEAN."
        ])
    ]

    for i, (qtitle, color, items) in enumerate(quarters):
        x = Inches(0.8 + i * 3.0)
        add_card(s13, x, Inches(1.8), Inches(2.75), Inches(5.1))
        tb = s13.shapes.add_textbox(x + Inches(0.15), Inches(2.0), Inches(2.45), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = qtitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for item in items:
            p_i = tf.add_paragraph()
            p_i.text = f"• {item}"
            p_i.font.size = Pt(9.5)
            p_i.font.color.rgb = TEXT_SECONDARY
            p_i.space_before = Pt(8)

    # ==========================================
    # SLIDE 14: KEY RISKS & NEXT PRIORITIES (CORE 14/14)
    # ==========================================
    s14 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(s14)
    add_header(s14, "Manajemen Risiko Kritis & Rencana Mitigasi Konkret", "CORE 14/14", 
               "Antisipasi proaktif terhadap risiko teknis, regulasi, operasional, dan adopsi pasar.", "14")

    # Risk Table
    rows = 5
    cols = 4
    t_shape14 = s14.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.73), Inches(4.8))
    table14 = t_shape14.table

    table14.columns[0].width = Inches(2.2)
    table14.columns[1].width = Inches(3.2)
    table14.columns[2].width = Inches(4.0)
    table14.columns[3].width = Inches(2.33)

    h14 = ["Kategori Risiko", "Deskripsi Potensi Bahaya", "Rencana Mitigasi Konkret", "Tingkat Risiko"]
    for j, h in enumerate(h14):
        cell = table14.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        p = cell.text_frame.paragraphs[0]
        p.text = h
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

    data_rows14 = [
        ("Risiko Teknis (False Positive)", "Virtual patch antibodi baru memblokir transaksi sah nasabah.", "Gerbang persetujuan L0/L1, whitelist rules spesifik, dan auto-rollback dalam <10 detik jika terjadi anomali 5xx.", "SEDANG (Terkendali)"),
        ("Risiko Kepatuhan Regulasi", "Perubahan standar audit siber dari OJK/BSSN pasca-rilis.", "Arsitektur audit trail modular berbasis ISO 27001 yang mudah disesuaikan dengan format pelaporan regulasi baru.", "RENDAH"),
        ("Risiko Beban Trafik (DDoS)", "Serangan volume tinggi melumpuhkan kapasitas gateway.", "Token Bucket rate limiter terdistribusi, Golden GET RAM cache, dan isolasi honeypot :9090.", "SEDANG (Terkendali)"),
        ("Risiko Resistensi Adopsi", "Keengganan tim IT perbankan mengubah arsitektur jaringan lama.", "Skema Reverse Proxy zero-downtime tanpa perlu mengubah kode sumber backend origin (cukup arahkan CNAME/DNS).", "RENDAH")
    ]

    for i, row in enumerate(data_rows14, start=1):
        for j, val in enumerate(row):
            cell = table14.cell(i, j)
            cell.fill.solid()
            if j == 3:
                cell.fill.fore_color.rgb = EMERALD_LIGHT if "RENDAH" in val else CARD_BG
            elif i % 2 == 0:
                cell.fill.fore_color.rgb = TABLE_ALT_ROW
            else:
                cell.fill.fore_color.rgb = CARD_BG
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(9.5)
            p.font.color.rgb = EMERALD_ACCENT if j == 3 and "RENDAH" in val else NAVY_PRIMARY if j == 0 else TEXT_SECONDARY
            if j == 0 or j == 3:
                p.font.bold = True

    # ==========================================
    # PART B: APPENDIX SLIDES (Q&A SUPPORT)
    # ==========================================

    # SLIDE A1: DEFENSE DELTA & LABELING LOGIC
    sa1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa1)
    add_header(sa1, "Appendix A1: Logika Wasit Defense Delta & Matriks Status Penutupan", "ADDITIONAL · Q&A", 
               "Detail taksonomi status penutupan Job Cowork untuk pembuktian ilmiah wasit jujur.", "A1")

    labels = [
        ("origin_open", "Celah Terbuka di Backend", AMBER_ACCENT, "Origin server menerima payload serangan; WAF belum memblokir. Risiko nyata teridentifikasi."),
        ("waf_blocked", "Ditahan di Tepi WAF", CYAN_ACCENT, "WAF berhasil memblokir serangan (403), meskipun backend origin mungkin masih memiliki kelemahan."),
        ("both_held", "Keduanya Menahan", EMERALD_ACCENT, "Baik WAF maupun origin backend sama-sama menolak payload serangan secara aman."),
        ("replay_held", "Tembakan Ulang Tertahan", EMERALD_ACCENT, "Setelah antibodi dipasang, tembakan ulang (replay attack) tetap berhasil ditolak 403."),
        ("replay_missed", "Uji Ulang Gagal", AMBER_ACCENT, "Antibodi gagal menahan variasi serangan baru. Job otomatis ditutup sebagai CLOSED_GAP (Residual Eksplisit)."),
        ("antibody_learned", "Antibodi Terverifikasi", NAVY_SECONDARY, "Antibodi zero-day sukses disimpan di memori imun host (Postgres) dan di-share ke RAM cache.")
    ]

    for i, (l_name, l_title, color, l_desc) in enumerate(labels):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.0)
        y = Inches(1.8 + row * 2.5)
        add_card(sa1, x, y, Inches(3.7), Inches(2.3))
        tb = sa1.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), Inches(3.4), Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = f"Label: {l_name}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        p_t = tf.add_paragraph()
        p_t.text = l_title
        p_t.font.size = Pt(11)
        p_t.font.bold = True
        p_t.font.color.rgb = NAVY_PRIMARY
        p_t.space_before = Pt(2)

        p_d = tf.add_paragraph()
        p_d.text = l_desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(4)

    # SLIDE A2: AI 8-POINT MANDATORY TECHNICAL STANDARD
    sa2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa2)
    add_header(sa2, "Appendix A2: Standar Teknis Wajib AI & Machine Learning (8 Poin PIDI)", "ADDITIONAL · Q&A", 
               "Pembuktian transparansi sistem AI: Bukan sekadar gimmick 'Powered by AI'.", "A2")

    ai_points = [
        ("1. Input Data", "Payload HTTP lengkap: Header, URL Query, Body, Client IP, Method."),
        ("2. Sumber Data", "Dataset serangan web riil (OWASP Top 10, CVE feeds) & traffic logs lab."),
        ("3. Preprocessing Data", "Normalisasi kanonik (stripping comments, URL decode, UTF-8 normalization)."),
        ("4. Output Model", "Klasifikasi biner (Attack/Legit), Threat Type, dan Confidence Score (0.0-1.0)."),
        ("5. Pemanfaatan Keputusan", "Memicu respon 403 instan atau merumuskan draf antibodi virtual patch baru."),
        ("6. Metrik Performa", "Latensi Reflex <1.2ms, Zero False Positive pada 10.000 legitimate HTTP requests."),
        ("7. Pengujian & Limitasi", "Diuji via live replay attack; memerlukan fallback jika model lokal belum siap."),
        ("8. Human Oversight", "Gerbang persetujuan L0/L1 wajib diverifikasi operator sebelum deploy produksi.")
    ]

    for i, (point_title, point_desc) in enumerate(ai_points):
        col = i % 2
        row = i // 2
        x = Inches(0.8 + col * 5.9)
        y = Inches(1.8 + row * 1.25)
        add_card(sa2, x, y, Inches(5.7), Inches(1.15))
        tb = sa2.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), Inches(5.4), Inches(0.95))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = point_title
        p.font.size = Pt(10.5)
        p.font.bold = True
        p.font.color.rgb = EMERALD_ACCENT

        p_d = tf.add_paragraph()
        p_d.text = point_desc
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY
        p_d.space_before = Pt(2)

    # SLIDE A3: DUAL-BRAIN AI ARCHITECTURE
    sa3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa3)
    add_header(sa3, "Appendix A3: Arsitektur Dual-Brain AI (Otak Kiri vs Otak Kanan)", "ADDITIONAL · Q&A", 
               "Pemisahan tugas komputasi sinkron latensi rendah dengan penalaran asinkron mendalam.", "A3")

    add_card(sa3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_db1 = sa3.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_db1 = tb_db1.text_frame
    tf_db1.word_wrap = True

    p = tf_db1.paragraphs[0]
    p.text = "OTAK KIRI: REFLEX FILTER (nex-ai-reflex)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    items_db1 = [
        ("Eksekusi Sinkron", "Berjalan langsung pada alur request HTTP utama di internal/ai/reflex_filter.go."),
        ("Latensi Ultra-Rendah", "Kecepatan eksekusi < 1.2 milidetik menggunakan regex terkompilasi."),
        ("Deteksi Pola Klasik", "Menangani SQL Injection, Cross-Site Scripting (XSS), Path Traversal, RCE, dan defacement."),
        ("RAM-First Cache", "Pola antibodi yang sudah dipelajari langsung dimasukkan ke memori RAM untuk blokir instan 403.")
    ]
    for h, d in items_db1:
        p_h = tf_db1.add_paragraph()
        p_h.text = f"✔ {h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_db1.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    add_card(sa3, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_db2 = sa3.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_db2 = tb_db2.text_frame
    tf_db2.word_wrap = True

    p = tf_db2.paragraphs[0]
    p.text = "OTAK KANAN: REASONING ENGINE (nex-ai-protect)"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    items_db2 = [
        ("Eksekusi Asinkron", "Dipanggil via Goroutine di latar belakang (timeout 30s) tanpa memblokir trafik pengguna."),
        ("Penalaran Mendalam", "Menganalisis intensi peretas, ancaman persisten tingkat lanjut (APT), dan anomali zero-day."),
        ("Formulasi Antibodi", "Menghasilkan aturan virtual patch baru dan menyimpannya ke tabel PostgreSQL antibody_audits."),
        ("Kemandirian Model Lokal", "Berjalan pada Ollama lokal host; dilarang keras fallback ke model publik eksternal.")
    ]
    for h, d in items_db2:
        p_h = tf_db2.add_paragraph()
        p_h.text = f"⚡ {h}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_db2.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # SLIDE A4: COMPLIANCE MAPPING
    sa4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa4)
    add_header(sa4, "Appendix A4: Pemetaan Kepatuhan Regulasi (POJK 30/2025 & ISO 27001)", "ADDITIONAL · Q&A", 
               "Kesesuaian langsung arsitektur Nexus Cyber dengan pasal-pasal regulasi perbankan.", "A4")

    add_card(sa4, Inches(0.8), Inches(1.8), Inches(11.73), Inches(5.1))
    tb_comp = sa4.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.7))
    tf_comp = tb_comp.text_frame
    tf_comp.word_wrap = True

    comp_items = [
        ("POJK No. 30/2025 Pasal Ketahanan Siber", "Mewajibkan penyelenggara ITSK memiliki sistem deteksi dan mitigasi insiden siber secara aktif.", "Nexus Cyber menyediakan proteksi tepi always-on dan virtual patch otomatis dalam hitungan menit."),
        ("POJK No. 30/2025 Pasal Uji Ketahanan Berkala", "Mewajibkan simulasi serangan siber dan uji penetrasi berkala.", "Wasit NEX-RED menjalankan Job Cowork berkala dengan pengujian tembakan ulang (replay verification) otomatis."),
        ("ISO 27001 Kontrol A.12.4 (Logging & Monitoring)", "Mewajibkan pencatatan log keamanan yang aman dari manipulasi dan dapat diaudit.", "Tabel PostgreSQL threat_logs menyimpan rekaman telemetri dengan UUID v4, SHA-256 prev_hash, dan target domain index."),
        ("UU PDP No. 27/2022 Pasal Perlindungan Kerahasiaan", "Mewajibkan enkripsi data nasabah dari risiko intersepsi dan kebocoran.", "PQC Shield berbasis NIST ML-KEM-768 melindungi payload komunikasi data sensitif dari ancaman komputasi kuantum.")
    ]

    for i, (reg, mand, sol) in enumerate(comp_items):
        p_r = tf_comp.paragraphs[0] if i == 0 else tf_comp.add_paragraph()
        p_r.text = reg
        p_r.font.size = Pt(11)
        p_r.font.bold = True
        p_r.font.color.rgb = NAVY_PRIMARY
        if i > 0: p_r.space_before = Pt(8)

        p_m = tf_comp.add_paragraph()
        p_m.text = f"• Mandat Regulasi: {mand}"
        p_m.font.size = Pt(9.5)
        p_m.font.color.rgb = TEXT_SECONDARY

        p_s = tf_comp.add_paragraph()
        p_s.text = f"✔ Implementasi Nexus: {sol}"
        p_s.font.size = Pt(9.5)
        p_s.font.bold = True
        p_s.font.color.rgb = EMERALD_ACCENT

    # SLIDE A5: DEMO SOP & 3-TIER BACKUP PLAN
    sa5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa5)
    add_header(sa5, "Appendix A5: SOP Demonstrasi Live & 3-Tier Demo Backup Plan", "ADDITIONAL · Q&A", 
               "Protokol mitigasi kegagalan teknis saat presentasi panggung offline maupun online.", "A5")

    tiers = [
        ("TIER 1: LIVE DEMO (PILIHAN UTAMA)", EMERALD_ACCENT, [
            ("Lingkungan Demo", "WAF Gateway lokal :8080 terhubung ke target live portofolio Vercel."),
            ("SOP Alur (90 Detik)", "1. Akses halaman target normal → 2. Tembakkan payload serangan (SQLi/Deface) → 3. WAF memblokir 403 → 4. Buka Command Center SOC (:3001) melihat telemetri log."),
            ("Pemeriksaan Pra-Demo", "Pastikan Docker Desktop aktif dan model Ollama nex-ai-* siap.")
        ]),
        ("TIER 2: VIDEO OFFLINE (CADANGAN 1)", CYAN_ACCENT, [
            ("Format File", "Rekaman MP4 1080p durasi 75 detik tersimpan di harddisk lokal laptop."),
            ("Skenario Penggunaan", "Diputar langsung jika koneksi internet venue/hotspot mengalami gangguan."),
            ("Konten Rekaman", "Screen-recording proses deteksi serangan, notifikasi Telegram, dan audit trail.")
        ]),
        ("TIER 3: SCREENSHOT FLOW (CADANGAN 2)", AMBER_ACCENT, [
            ("Format", "Rangkaian screenshot alur langkah demi langkah pada slide appendix."),
            ("Skenario Penggunaan", "Digunakan jika laptop mengalami kendala pemutaran media video."),
            ("Konten", "Gambar input payload, respon 403 terminal, dan visualisasi SOC dashboard.")
        ])
    ]

    for i, (ttitle, color, items) in enumerate(tiers):
        x = Inches(0.8 + i * 4.0)
        add_card(sa5, x, Inches(1.8), Inches(3.7), Inches(5.1))
        tb = sa5.shapes.add_textbox(x + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.7))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ttitle
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = color

        for h, d in items:
            p_h = tf.add_paragraph()
            p_h.text = h
            p_h.font.size = Pt(10.5)
            p_h.font.bold = True
            p_h.font.color.rgb = NAVY_PRIMARY
            p_h.space_before = Pt(8)

            p_d = tf.add_paragraph()
            p_d.text = d
            p_d.font.size = Pt(9.5)
            p_d.font.color.rgb = TEXT_SECONDARY
            p_d.space_before = Pt(2)

    # SLIDE A6: DETAILED FINANCIAL MODEL & COGS
    sa6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa6)
    add_header(sa6, "Appendix A6: Struktur Biaya Pokok (COGS) & Proyeksi LTV / CAC", "ADDITIONAL · Q&A", 
               "Analisis profitabilitas unit economics untuk meyakinkan investor dan juri bisnis.", "A6")

    add_card(sa6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_cogs = sa6.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_cogs = tb_cogs.text_frame
    tf_cogs.word_wrap = True

    p = tf_cogs.paragraphs[0]
    p.text = "STRUKTUR BIAYA POKOK LAYANAN (COGS PER HOST)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = EMERALD_ACCENT

    cogs_items = [
        ("Komputasi Gateway Go (VPS Cloud)", "Rp 45.000 / bln", "High-efficiency Go runtime memungkinkan 1 VPS Rp 200rb melayani 4-5 host kanal."),
        ("Inferensi AI Lokal (Ollama)", "Rp 15.000 / bln", "Zero API token cost berkat bobot model mandiri nex-ai-protect."),
        ("Storage Audit Log (PostgreSQL)", "Rp 10.000 / bln", "Penyimpanan log terindeks dengan efisiensi kompresi data tinggi."),
        ("Total COGS per Host", "Rp 70.000 / bln", "Harga Jual Loop: Rp 300.000 / bln $\rightarrow$ Gross Margin 76.6%.")
    ]
    for h, pr, d in cogs_items:
        p_h = tf_cogs.add_paragraph()
        p_h.text = f"{h}: {pr}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_cogs.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    add_card(sa6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(5.1))
    tb_ltv = sa6.shapes.add_textbox(Inches(7.1), Inches(2.0), Inches(5.2), Inches(4.7))
    tf_ltv = tb_ltv.text_frame
    tf_ltv.word_wrap = True

    p = tf_ltv.paragraphs[0]
    p.text = "METRIK SKALABILITAS BISNIS (LTV / CAC)"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = CYAN_ACCENT

    ltv_items = [
        ("Customer Acquisition Cost (CAC)", "Rp 1.500.000", "Diperoleh melalui kemitraan asosiasi BPR/Fintech dan program pilot terarah."),
        ("Customer Lifetime Value (LTV)", "Rp 10.800.000", "Berdasarkan masa retensi rata-rata 36 bulan dengan skema retainership bulanan."),
        ("Rasio LTV / CAC", "7.2x (Sangat Sehat)", "Standar industri SaaS B2B ideal adalah > 3.0x."),
        ("Payback Period", "4.5 Bulan", "Modal akuisisi kembali dalam waktu kurang dari satu semester.")
    ]
    for h, pr, d in ltv_items:
        p_h = tf_ltv.add_paragraph()
        p_h.text = f"{h}: {pr}"
        p_h.font.size = Pt(10.5)
        p_h.font.bold = True
        p_h.font.color.rgb = NAVY_PRIMARY
        p_h.space_before = Pt(8)
        p_d = tf_ltv.add_paragraph()
        p_d.text = d
        p_d.font.size = Pt(9.5)
        p_d.font.color.rgb = TEXT_SECONDARY

    # SLIDE A7: ELIGIBILITY SELF-CHECKLIST
    sa7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa7)
    add_header(sa7, "Appendix A7: Checklist Evaluasi Mandiri 12 Kriteria PIDI 2026", "ADDITIONAL · Q&A", 
               "Verifikasi kelayakan menyeluruh sebelum submission final dan pitching juri.", "A7")

    add_card(sa7, Inches(0.8), Inches(1.8), Inches(11.73), Inches(5.1))
    tb_chk = sa7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(4.7))
    tf_chk = tb_chk.text_frame
    tf_chk.word_wrap = True

    checks = [
        ("1. Problem", "Bukti data resmi ancaman siber dan sanksi regulasi POJK 30/2025."),
        ("2. Alignment", "Konsisten memilih Problem Statement 1 (Manajemen Risiko Finansial)."),
        ("3. Solution", "Penjelasan sederhana tanpa jargon: Ukur → Kendalikan → Uji."),
        ("4. Prototype", "Functional Prototype Go WAF (:8080) & NEX-RED (:3004) berjalan nyata."),
        ("5. Technical", "Arsitektur 3-Layer jelas dari UX hingga database ISO 27001."),
        ("6. Testing", "Data benchmarking latensi (<1.2ms) dan akurasi wasit (99.4%)."),
        ("7. Impact", "Metrik evaluasi baseline vs target kuantitatif (MTTR <5 menit)."),
        ("8. Market", "Target pasar 1.400+ BPR & 100+ Fintech tervalidasi regulasi."),
        ("9. Differentiation", "Matriks komparasi jelas membedakan diri dari WAF & scanner lama."),
        ("10. Team", "Pembagian peran dan kepemilikan kode sumber nyata per anggota."),
        ("11. Roadmap", "Rencana pasca-hackathon milestone-based Q1-Q4 konkret."),
        ("12. Transparency", "Membedakan tegas fitur aktif vs batasan eBPF stub & roadmap.")
    ]

    for i, (chk_title, chk_desc) in enumerate(checks):
        p_c = tf_chk.paragraphs[0] if i == 0 else tf_chk.add_paragraph()
        p_c.text = f"✔ [LULUS] {chk_title}: {chk_desc}"
        p_c.font.size = Pt(9.5)
        p_c.font.color.rgb = NAVY_PRIMARY
        if i > 0: p_c.space_before = Pt(2)

    # SLIDE A8: CLOSING & CONTACT
    sa8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_background(sa8)

    add_card(sa8, Inches(2.0), Inches(1.5), Inches(9.33), Inches(4.5), bg_color=CARD_BG, border_color=EMERALD_ACCENT)
    tb_cl = sa8.shapes.add_textbox(Inches(2.3), Inches(1.8), Inches(8.73), Inches(3.9))
    tf_cl = tb_cl.text_frame
    tf_cl.word_wrap = True

    p = tf_cl.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "NEXUS CYBER"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY_PRIMARY

    p_sub = tf_cl.add_paragraph()
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "Solusi Manajemen Risiko Siber Kanal Digital Terpercaya untuk Ekosistem Keuangan Indonesia"
    p_sub.font.size = Pt(14)
    p_sub.font.color.rgb = EMERALD_ACCENT
    p_sub.space_before = Pt(4)

    p_div = tf_cl.add_paragraph()
    p_div.alignment = PP_ALIGN.CENTER
    p_div.text = "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━"
    p_div.font.size = Pt(10)
    p_div.font.color.rgb = CARD_BORDER
    p_div.space_before = Pt(10)

    p_contact = tf_cl.add_paragraph()
    p_contact.alignment = PP_ALIGN.CENTER
    p_contact.text = "SIAP BERKOLABORASI & MEMULAI PROGRAM PILOT BERSAMA INSTITUSI ANDA"
    p_contact.font.size = Pt(12)
    p_contact.font.bold = True
    p_contact.font.color.rgb = NAVY_PRIMARY
    p_contact.space_before = Pt(10)

    p_info = tf_cl.add_paragraph()
    p_info.alignment = PP_ALIGN.CENTER
    p_info.text = "Repository: github.com/Thbetyfu/NEXUS-CYBER-FASE3 · Web Portal: nexus-gaas-web (:3003)\nKontak Tim: thoriq@nexus-cyber.test / WhatsApp: +62 895-6033-58692"
    p_info.font.size = Pt(11)
    p_info.font.color.rgb = TEXT_SECONDARY
    p_info.space_before = Pt(8)

    # Save Output
    output_path = r"d:\NEXUS\Lomba\power point\Nexus-Cyber-PIDI-2026-Official.pptx"
    prs.save(output_path)
    print(f"[SUCCESS] Official PIDI 2026 Pitch Deck generated successfully: {output_path}")
    print(f"Total Slides Generated: {len(prs.slides)} slides (14 Core Slides + 8 Appendix Slides).")

if __name__ == "__main__":
    create_deck()
